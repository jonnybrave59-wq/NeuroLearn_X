from __future__ import annotations

import asyncio
import csv
import hashlib
import io
import json
import logging
import os
import random
import re
import secrets
import time
import uuid
import zipfile
from collections import defaultdict, deque
from contextlib import asynccontextmanager
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
from urllib.parse import urlsplit

from fastapi import (
    Depends,
    FastAPI,
    File,
    HTTPException,
    Query,
    Request,
    Response,
    UploadFile,
    status,
)
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from starlette.middleware.httpsredirect import HTTPSRedirectMiddleware
from sqlalchemy import and_, delete, func, or_, select, text, update
from sqlalchemy.exc import IntegrityError
from sqlalchemy.orm import Session

from .algorithms import (
    adaptive_pathway_score,
    prerequisite_ancestors,
    topological_order,
    would_create_cycle,
)
from .database import Base, SessionLocal, engine, get_db
from .ml import predict_activity_load, predict_student_cognitive_load, train_ensemble
from .models import (
    Activity,
    ActivityConcept,
    AnswerChoice,
    Assessment,
    AssessmentAssignment,
    AssessmentQuestion,
    AssessmentAttempt,
    AuditLog,
    Concept,
    ConsentRecord,
    CognitiveLoadPrediction,
    ExpertEvaluation,
    InteractionLog,
    ItemResponse,
    LearningGap,
    LearningSummary,
    LoginHistory,
    MasteryRecord,
    MentalEffortRating,
    Misconception,
    MisconceptionHistory,
    ModelVersion,
    PathwayRecommendation,
    PathwayStep,
    PathwayVersion,
    PrerequisiteEdge,
    Question,
    StudentProfile,
    UploadedDocument,
    User,
    TeacherIntervention,
    TutoringSession,
    TutoringResponse,
)
from .onboarding import (
    ONBOARDING_DIAGNOSTIC_ITEMS,
    ensure_onboarding_diagnostic,
)
from .schemas import (
    ActivityInput,
    ActivityBulkInput,
    AssessmentInput,
    AssessmentStatusInput,
    AttemptInput,
    ConceptInput,
    DocumentGenerationInput,
    EdgeInput,
    EvaluationInput,
    ForgotPasswordInput,
    LoginInput,
    MentalEffortInput,
    MisconceptionInput,
    PasswordInput,
    PathwayAssignmentInput,
    PathwayPreviewInput,
    QuestionBankInput,
    QuestionBatchInput,
    QuestionBulkEditInput,
    QuestionInput,
    ResetInput,
    SettingsInput,
    StudentActionInput,
    StudentRegistrationInput,
    TargetInput,
    TutoringResponseInput,
    TutoringSessionInput,
    InterventionInput,
)
from .security import (
    COOKIE_NAME,
    create_access_token,
    current_user,
    hash_password,
    require_role,
    verify_password,
)
from .services import (
    audit,
    build_pathway_preview,
    generate_pathways,
    get_setting,
    latest_mastery_map,
    learning_content_for,
    mental_effort_category,
    record_pathway_evidence,
    recalculate_mastery,
    save_settings,
    serialize_pathway,
    settings_payload,
)
from .tutoring import (
    start_tutoring_session,
    structured_solution,
    submit_tutoring_response,
    update_summary_effort,
)


APP_ENV = os.getenv("APP_ENV", "development").strip().lower()
RENDER_EXTERNAL_HOSTNAME = os.getenv("RENDER_EXTERNAL_HOSTNAME", "").strip()
PUBLIC_APP_URL = os.getenv("PUBLIC_APP_URL", "").strip()
if not PUBLIC_APP_URL and RENDER_EXTERNAL_HOSTNAME:
    PUBLIC_APP_URL = f"https://{RENDER_EXTERNAL_HOSTNAME}"
COOKIE_SECURE = os.getenv("COOKIE_SECURE", "0") == "1"
COOKIE_SAMESITE = os.getenv("COOKIE_SAMESITE", "lax").strip().lower()
if COOKIE_SAMESITE not in {"lax", "strict", "none"}:
    raise RuntimeError("COOKIE_SAMESITE must be lax, strict, or none")

logging.basicConfig(
    level=getattr(logging, os.getenv("LOG_LEVEL", "INFO").upper(), logging.INFO),
    format="%(asctime)s %(levelname)s %(name)s %(message)s",
)
logger = logging.getLogger("neurolearnx")


def normalized_origin(value: str) -> str | None:
    """Return an exact scheme/host/port origin or reject unsafe input."""
    value = value.strip().rstrip("/")
    if not value or value == "*":
        return None
    parts = urlsplit(value)
    if (
        parts.scheme not in {"http", "https", "capacitor"}
        or not parts.netloc
        or "*" in parts.netloc
    ):
        return None
    if parts.username or parts.password:
        return None
    return f"{parts.scheme}://{parts.netloc}"


def configured_origins() -> tuple[str, ...]:
    values = [
        item
        for item in os.getenv("ALLOWED_ORIGINS", "").split(",")
        if item.strip()
    ]
    values.extend([PUBLIC_APP_URL, os.getenv("FRONTEND_ORIGIN", "")])
    values.extend(
        item
        for item in os.getenv("CAPACITOR_ORIGINS", "").split(",")
        if item.strip()
    )
    if APP_ENV != "production" and os.getenv("REPLIT_DEPLOYMENT") != "1":
        values.append("http://127.0.0.1:5173")
    origins = [origin for value in values if (origin := normalized_origin(value))]
    return tuple(dict.fromkeys(origins))


ALLOWED_ORIGINS = configured_origins()
PRODUCTION = APP_ENV == "production" or os.getenv("REPLIT_DEPLOYMENT") == "1"
if PRODUCTION:
    public_origin = normalized_origin(PUBLIC_APP_URL)
    if not public_origin or not public_origin.startswith("https://"):
        raise RuntimeError("Production requires a clean HTTPS PUBLIC_APP_URL")
    if not COOKIE_SECURE:
        raise RuntimeError("Production requires COOKIE_SECURE=1")
    if public_origin not in ALLOWED_ORIGINS:
        raise RuntimeError("Production ALLOWED_ORIGINS must include the public app")
    cross_origin_clients = [
        origin for origin in ALLOWED_ORIGINS if origin != public_origin
    ]
    if cross_origin_clients and COOKIE_SAMESITE != "none":
        raise RuntimeError(
            "Cross-origin production clients require COOKIE_SAMESITE=none"
        )


@asynccontextmanager
async def lifespan(_app: FastAPI):
    if os.getenv("CREATE_TABLES_ON_STARTUP", "1") == "1":
        Base.metadata.create_all(engine)
    with SessionLocal() as db:
        diagnostic = ensure_onboarding_diagnostic(db)
        if diagnostic is None:
            logger.warning(
                "The onboarding diagnostic could not be prepared because 30 eligible "
                "multiple-choice diagnostic items are not available."
            )
    try:
        yield
    finally:
        engine.dispose()
        logger.info("Database connection pool closed during application shutdown")


app = FastAPI(
    title="NeuroLearn-X API",
    version="1.3.1",
    description="Explainable adaptive learning research prototype.",
    lifespan=lifespan,
)
if PRODUCTION:
    # Uvicorn trusts the deployment proxy configured by scripts/replit-run.sh,
    # so X-Forwarded-Proto is resolved before this middleware checks the scheme.
    app.add_middleware(HTTPSRedirectMiddleware)
app.add_middleware(
    CORSMiddleware,
    allow_origins=list(ALLOWED_ORIGINS),
    allow_credentials=True,
    allow_methods=["GET", "POST", "PUT", "DELETE", "OPTIONS"],
    allow_headers=[
        "Accept",
        "Content-Type",
        "Origin",
        "X-Requested-With",
        "X-Request-ID",
    ],
)

RATE_BUCKETS: dict[str, deque[float]] = defaultdict(deque)


def enforce_rate_limit(
    request: Request,
    scope: str,
    limit: int,
    window: int = 60,
    identity: str | None = None,
):
    client = request.client.host if request.client else "unknown"
    key = f"{scope}:{identity or client}"
    now = time.monotonic()
    bucket = RATE_BUCKETS[key]
    while bucket and bucket[0] <= now - window:
        bucket.popleft()
    if len(bucket) >= limit:
        raise HTTPException(
            status_code=429,
            detail="Too many requests. Please wait briefly and try again.",
        )
    bucket.append(now)


@app.middleware("http")
async def security_headers(request: Request, call_next):
    started = time.monotonic()
    supplied_request_id = (request.headers.get("x-request-id") or "").strip()
    request_id = (
        supplied_request_id
        if re.fullmatch(r"[A-Za-z0-9._:-]{8,80}", supplied_request_id)
        else uuid.uuid4().hex
    )
    request.state.request_id = request_id
    origin = (request.headers.get("origin") or "").rstrip("/")
    request_origin = normalized_origin(str(request.base_url))
    unsafe_method = request.method.upper() not in {"GET", "HEAD", "OPTIONS"}
    missing_authenticated_origin = (
        PRODUCTION
        and unsafe_method
        and COOKIE_NAME in request.cookies
        and not origin
    )
    if missing_authenticated_origin:
        response = JSONResponse(
            status_code=403,
            content={"detail": "Authenticated changes require an approved origin"},
        )
    elif (
        unsafe_method
        and origin
        and origin != request_origin
        and origin not in ALLOWED_ORIGINS
    ):
        response = JSONResponse(
            status_code=403,
            content={"detail": "Request origin is not approved"},
        )
    else:
        response = await call_next(request)
    if (
        (
            response.status_code == 401
            or (
                response.status_code == 403
                and request.url.path
                in {"/api/student/dashboard", "/api/teacher/dashboard"}
            )
        )
        and request.url.path.startswith("/api/")
        and request.url.path
        not in {
            "/api/auth/login",
            "/api/auth/register/student",
            "/api/auth/forgot-password",
        }
    ):
        response.delete_cookie(
            COOKIE_NAME,
            path="/",
            secure=COOKIE_SECURE,
            samesite=COOKIE_SAMESITE,
        )
    response.headers["X-Content-Type-Options"] = "nosniff"
    response.headers["X-Request-ID"] = request_id
    response.headers["Referrer-Policy"] = "no-referrer"
    response.headers["Cross-Origin-Opener-Policy"] = "same-origin"
    response.headers["Permissions-Policy"] = (
        "camera=(), microphone=(), geolocation=(), payment=(), usb=()"
    )
    response.headers["Content-Security-Policy"] = (
        "default-src 'self'; "
        "script-src 'self'; "
        "style-src 'self' 'unsafe-inline'; "
        "img-src 'self' data: blob:; "
        "font-src 'self' data:; "
        "connect-src 'self' https:; "
        "object-src 'none'; base-uri 'self'; frame-ancestors 'none'"
    )
    if COOKIE_SECURE:
        response.headers["Strict-Transport-Security"] = (
            "max-age=31536000; includeSubDomains"
        )
    if request.url.path.startswith(("/api/auth/", "/api/student/", "/api/teacher/")):
        response.headers["Cache-Control"] = "no-store, max-age=0"
        response.headers["Pragma"] = "no-cache"
    elapsed_ms = (time.monotonic() - started) * 1000
    logger.info(
        "request_id=%s %s %s -> %s %.1fms",
        request_id,
        request.method,
        request.url.path,
        response.status_code,
        elapsed_ms,
    )
    return response


@app.exception_handler(Exception)
async def unhandled_exception(request: Request, error: Exception):
    logger.error(
        "request_id=%s unhandled server error on %s %s",
        getattr(request.state, "request_id", "unknown"),
        request.method,
        request.url.path,
        exc_info=(type(error), error, error.__traceback__),
    )
    return JSONResponse(
        status_code=500,
        content={"detail": "Internal server error"},
    )


def user_payload(user: User) -> dict[str, Any]:
    return {
        "id": user.id,
        "participant_code": user.participant_code,
        "username": user.username,
        "email": user.email,
        "display_name": user.display_name,
        "first_name": user.first_name,
        "last_name": user.last_name,
        "role": user.role,
        "must_change_password": user.must_change_password,
        "is_demo": user.is_demo,
        "account_status": user.account_status,
        "created_at": user.created_at,
        "last_sign_in_at": user.last_sign_in_at,
    }


def concept_payload(concept: Concept) -> dict[str, Any]:
    return {
        "id": concept.id,
        "code": concept.code,
        "name": concept.name,
        "subject": concept.subject,
        "description": concept.description,
        "difficulty": concept.difficulty,
        "active": concept.active,
    }


def activity_payload(
    db: Session,
    activity: Activity,
    include_questions: bool = False,
    include_correct_answers: bool = False,
):
    concept_ids = list(
        db.scalars(
            select(ActivityConcept.concept_id).where(
                ActivityConcept.activity_id == activity.id
            )
        )
    )
    payload = {
        "id": activity.id,
        "title": activity.title,
        "description": activity.description,
        "activity_type": activity.activity_type,
        "difficulty": activity.difficulty,
        "estimated_minutes": activity.estimated_minutes,
        "instructions": activity.instructions,
        "resource_url": activity.resource_url,
        "active": activity.active,
        "is_diagnostic": activity.is_diagnostic,
        "is_onboarding_diagnostic": activity.is_onboarding_diagnostic,
        "concept_ids": concept_ids,
        "dependencies": {
            "assigned_students": db.scalar(
                select(func.count(func.distinct(AssessmentAssignment.student_id)))
                .join(Assessment, Assessment.id == AssessmentAssignment.assessment_id)
                .where(Assessment.activity_id == activity.id)
            ) or 0,
            "attempts": db.scalar(
                select(func.count(AssessmentAttempt.id)).where(
                    AssessmentAttempt.activity_id == activity.id
                )
            ) or 0,
            "results": db.scalar(
                select(func.count(AssessmentAttempt.id)).where(
                    AssessmentAttempt.activity_id == activity.id,
                    AssessmentAttempt.submitted_at.is_not(None),
                )
            ) or 0,
            "pathway_steps": db.scalar(
                select(func.count(PathwayStep.id)).where(
                    PathwayStep.activity_id == activity.id
                )
            ) or 0,
        },
    }
    if include_questions:
        questions = list(
            db.scalars(
                select(Question)
                .where(Question.activity_id == activity.id, Question.active.is_(True))
                .order_by(Question.position)
            )
        )
        payload["questions"] = []
        for question in questions:
            question_data = {
                "id": question.id,
                "concept_id": question.concept_id,
                "prompt": question.prompt,
                "question_type": question.question_type,
                "hint": question.hint,
                "points": question.points,
                "position": question.position,
                "choices": [
                    {
                        "id": choice.id,
                        "text": choice.text,
                        "position": choice.position,
                        **(
                            {"is_correct": choice.is_correct}
                            if include_correct_answers
                            else {}
                        ),
                    }
                    for choice in db.scalars(
                        select(AnswerChoice)
                        .where(AnswerChoice.question_id == question.id)
                        .order_by(AnswerChoice.position)
                    )
                ],
            }
            if include_correct_answers:
                question_data["feedback"] = question.feedback
            payload["questions"].append(question_data)
    return payload


def clean_content(value: str, maximum: int = 5000) -> str:
    value = value.replace("\x00", " ")
    value = "".join(
        character
        for character in value
        if character in "\n\t" or ord(character) >= 32
    )
    return re.sub(r"[ \t]+", " ", value).strip()[:maximum]


def uploaded_document_payload(document: UploadedDocument) -> dict[str, Any]:
    return {
        "id": document.id,
        "original_filename": document.original_filename,
        "file_type": document.file_type,
        "file_size": document.file_size,
        "processing_status": document.processing_status,
        "text_length": len(document.extracted_text or ""),
        "text_preview": clean_content(document.extracted_text or "", 450),
        "analysis": document.analysis or {},
        "created_at": document.created_at,
    }


def question_bank_payload(db: Session, question: Question) -> dict[str, Any]:
    concept = db.get(Concept, question.concept_id)
    source = (
        db.get(UploadedDocument, question.source_document_id)
        if question.source_document_id
        else None
    )
    choices = list(
        db.scalars(
            select(AnswerChoice)
            .where(AnswerChoice.question_id == question.id)
            .order_by(AnswerChoice.position)
        )
    )
    correct_answer = question.correct_answer
    if not correct_answer:
        correct_answer = next(
            (choice.text for choice in choices if choice.is_correct), ""
        )
    return {
        "id": question.id,
        "activity_id": question.activity_id,
        "concept_id": question.concept_id,
        "concept": concept.name if concept else "Unknown concept",
        "prompt": question.prompt,
        "question_type": question.question_type,
        "correct_answer": correct_answer,
        "explanation": question.explanation or question.feedback,
        "hint": question.hint,
        "difficulty": question.difficulty_label,
        "cognitive_level": question.cognitive_level,
        "subject": question.subject,
        "topic": question.topic,
        "learning_competency": question.learning_competency,
        "source_type": question.source_type,
        "source_document_id": question.source_document_id,
        "source_document": source.original_filename if source else None,
        "source_locator": question.source_locator,
        "solution_steps": question.solution_steps,
        "solution_structure": question.solution_structure or {},
        "estimated_cognitive_demand": question.estimated_cognitive_demand,
        "prerequisite_concept_id": question.prerequisite_concept_id,
        "validation_status": question.validation_status,
        "validation_flags": question.validation_flags or [],
        "generation_metadata": question.generation_metadata or {},
        "distractor_rationales": question.distractor_rationales or {},
        "is_calculation": question.is_calculation,
        "status": question.status,
        "points": question.points,
        "choices": [
            {
                "id": choice.id,
                "text": choice.text,
                "is_correct": choice.is_correct,
                "position": choice.position,
                "misconception_id": choice.misconception_id,
                "misconception_confidence": choice.misconception_confidence,
                "mapping_status": choice.mapping_status,
                "misconception": (
                    {
                        "code": misconception.code,
                        "name": misconception.name,
                        "explanation": misconception.explanation,
                        "remediation_instruction": misconception.remediation_instruction,
                    }
                    if choice.misconception_id
                    and (misconception := db.get(Misconception, choice.misconception_id))
                    else None
                ),
            }
            for choice in choices
        ],
        "created_at": question.created_at,
        "updated_at": question.updated_at,
    }


def validate_office_archive(extension: str, data: bytes):
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            expanded = sum(item.file_size for item in archive.infolist())
            if expanded > 50 * 1024 * 1024:
                raise HTTPException(
                    status_code=413,
                    detail="The expanded document is too large to process safely",
                )
            names = set(archive.namelist())
            required_member = {
                ".docx": "word/document.xml",
                ".pptx": "ppt/presentation.xml",
            }[extension]
            if "[Content_Types].xml" not in names or required_member not in names:
                raise HTTPException(
                    status_code=422,
                    detail="The office document signature does not match its file type",
                )
    except zipfile.BadZipFile:
        raise HTTPException(status_code=422, detail="The office document is unreadable")


def extract_document_text(extension: str, data: bytes) -> str:
    try:
        if extension == ".txt":
            text_value = data.decode("utf-8-sig")
        elif extension == ".pdf":
            from pypdf import PdfReader

            reader = PdfReader(io.BytesIO(data))
            if len(reader.pages) > 250:
                raise HTTPException(status_code=413, detail="PDF page limit exceeded")
            pages = []
            for index, page in enumerate(reader.pages, start=1):
                try:
                    page_text = page.extract_text(extraction_mode="layout") or ""
                except TypeError:
                    page_text = page.extract_text() or ""
                marker = "" if len(page_text.strip()) >= 20 else "\n[Unreadable page: OCR unavailable or no text layer]"
                pages.append(f"[Page {index}]\n{page_text}{marker}")
            text_value = "\n".join(pages)
        elif extension == ".docx":
            validate_office_archive(extension, data)
            from docx import Document

            document = Document(io.BytesIO(data))
            paragraphs = []
            for paragraph in document.paragraphs:
                if not paragraph.text.strip():
                    continue
                style = (paragraph.style.name if paragraph.style else "").casefold()
                marker = "[Heading] " if "heading" in style or style == "title" else ""
                paragraphs.append(f"{marker}{paragraph.text}")
            text_value = "\n".join(paragraphs)
        elif extension == ".pptx":
            validate_office_archive(extension, data)
            from pptx import Presentation

            presentation = Presentation(io.BytesIO(data))
            text_value = "\n".join(
                f"[Slide {number}]\n"
                + "\n".join(
                    shape.text
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text
                )
                for number, slide in enumerate(presentation.slides, start=1)
            )
        else:
            raise HTTPException(status_code=415, detail="Unsupported file type")
    except HTTPException:
        raise
    except Exception:
        raise HTTPException(
            status_code=422,
            detail="The uploaded document could not be read safely",
        )
    text_value = clean_content(text_value, int(os.getenv("MAX_EXTRACTED_TEXT_CHARS", "250000")))
    if len(text_value) < 40:
        raise HTTPException(
            status_code=422,
            detail="The document does not contain enough readable text",
        )
    return text_value


def analyze_material(text_value: str) -> dict[str, Any]:
    lines = [
        clean_content(line, 1000)
        for line in text_value.splitlines()
        if clean_content(line, 1000)
    ]
    content_lines = [line for line in lines if not re.fullmatch(r"\[(?:Page|Slide) \d+\]", line)]
    title = next(
        (
            re.sub(r"^\[Heading\]\s*", "", line)
            for line in content_lines
            if len(re.sub(r"^\[Heading\]\s*", "", line)) <= 140
        ),
        "Uploaded learning material",
    )
    headings = [
        re.sub(r"^\[Heading\]\s*", "", line)
        for line in content_lines
        if line.startswith("[Heading]")
        or (
            len(line) <= 90
            and (
                line.isupper()
                or line.endswith(":")
                or (
                    len(line.split()) <= 9
                    and not re.search(r"[.!?=]$", line)
                    and sum(word[:1].isupper() for word in line.split())
                    >= max(1, len(line.split()) // 2)
                )
            )
        )
    ][:12]
    sentences = source_sentences(text_value)
    definitions = [
        sentence
        for sentence in sentences
        if re.search(r"\b(is defined as|means|refers to| is an? | is the )\b", sentence, re.I)
    ][:10]
    formulas = [
        line for line in content_lines
        if re.search(r"[=≈±√ΔΣ]|(?:\b[A-Za-z]\s*=\s*)", line)
    ][:12]
    examples = [
        sentence
        for sentence in sentences
        if re.search(r"\b(example|given|solution|calculate|suppose)\b", sentence, re.I)
        or (
            any(re.search(r"\b(worked example|example|solution)\b", heading, re.I) for heading in headings)
            and re.search(r"\d.*(?:=|\bN\b|kg|m/s|J\b|Pa\b|W\b)", sentence)
        )
    ][:10]
    competencies = [
        sentence
        for sentence in sentences
        if re.search(
            r"\b(explain|calculate|identify|analy[sz]e|apply|compare|determine|solve|describe|learners? should|students? should)\b",
            sentence,
            re.I,
        )
    ][:10]
    relationships = [
        sentence
        for sentence in sentences
        if re.search(
            r"\b(causes?|depends? on|proportional|related (?:to|by)|therefore|because|conserved|increases?|decreases?)\b",
            sentence,
            re.I,
        )
    ][:10]
    misconceptions = [
        sentence
        for sentence in sentences
        if re.search(r"\b(common error|misconception|do not confuse|incorrect|mistake|should not|avoid)\b", sentence, re.I)
    ][:10]
    words = re.findall(r"[A-Za-z][A-Za-z'-]{3,}", " ".join(content_lines).casefold())
    stop = {
        "this", "that", "with", "from", "have", "will", "when", "where", "which",
        "their", "there", "these", "those", "into", "using", "used", "than", "then",
        "also", "such", "each", "between", "because", "about", "your", "does",
    }
    counts: dict[str, int] = defaultdict(int)
    for word in words:
        if word not in stop:
            counts[word] += 1
    key_concepts = [
        word.title()
        for word, count in sorted(counts.items(), key=lambda item: (-item[1], item[0]))
        if count >= 2
    ][:12]
    normalized = " ".join(content_lines).casefold()
    mathematics_terms = len(re.findall(
        r"\b(equation|algebra|function|geometry|probability|calculate|formula|variable|slope|trigonometry|mathematics)\b|[=+\-*/]",
        normalized,
    ))
    physics_terms = len(re.findall(
        r"\b(force|motion|velocity|energy|momentum|newton|physics|acceleration|vector)\b",
        normalized,
    ))
    subject = (
        "Mathematics"
        if mathematics_terms >= max(2, physics_terms)
        else "Physics"
        if physics_terms >= 2
        else "General education"
    )
    filipino_markers = len(re.findall(r"\b(ang|mga|ng|sa|ay|para|mula|bilang)\b", normalized))
    english_markers = len(re.findall(r"\b(the|and|of|to|is|for|with|from)\b", normalized))
    language = "Filipino" if filipino_markers > english_markers else "English"
    page_numbers = sorted({int(value) for value in re.findall(r"\[Page (\d+)\]", text_value)})
    unreadable_pages = sorted({int(value) for value in re.findall(
        r"\[Page (\d+)\]\s*\[Unreadable page:", text_value
    )})
    variable_names = sorted(set(re.findall(r"\b([A-Za-z])\s*=", " ".join(formulas))))
    prerequisites = [
        sentence for sentence in sentences
        if re.search(r"\b(prerequisite|prior knowledge|before learning|requires?|recall)\b", sentence, re.I)
    ][:8]
    difficulty_score = min(5, 1 + len(formulas) // 3 + len(examples) // 4)
    learner_level = (
        "Advanced secondary or early tertiary"
        if difficulty_score >= 4
        else "Secondary"
        if difficulty_score >= 2
        else "Introductory"
    )
    return {
        "title": title,
        "module_title": title,
        "detected_subject": subject,
        "detected_language": language,
        "headings": list(dict.fromkeys(headings)),
        "main_topic": headings[0] if headings else title,
        "main_topics": list(dict.fromkeys(headings[:5] or [title])),
        "subtopics": list(dict.fromkeys(headings[1:12])),
        "key_concepts": key_concepts,
        "important_concepts": key_concepts,
        "definitions": definitions,
        "facts": sentences[:10],
        "formulas": formulas,
        "worked_examples": examples,
        "competencies": competencies,
        "learning_objectives": competencies,
        "relationships": relationships,
        "misconceptions": misconceptions,
        "formula_variables": variable_names,
        "prerequisites": prerequisites,
        "estimated_learner_level": learner_level,
        "estimated_difficulty": ["Introductory", "Moderate", "Advanced", "Advanced", "Advanced"][difficulty_score - 1],
        "pages_used_as_evidence": page_numbers,
        "unreadable_pages": unreadable_pages,
        "ocr_status": (
            "Some pages have no readable text layer; install an OCR provider to process them."
            if unreadable_pages
            else "Not required; all pages contained readable text."
        ),
        "method": "Local deterministic text analysis; no external AI service was used.",
        "limitations": (
            "The analysis identifies explicit text patterns. Scanned images, implicit "
            "meaning, and complex mathematical notation may require teacher review."
        ),
    }


def source_passages(text_value: str) -> list[dict[str, str]]:
    locator = "Uploaded material"
    container_locator = "Uploaded material"
    passages: list[dict[str, str]] = []
    buffered_lines: list[str] = []

    def flush_buffer() -> None:
        if not buffered_lines:
            return
        for sentence in source_sentences(" ".join(buffered_lines)):
            passages.append({"text": sentence, "source_locator": locator})
        buffered_lines.clear()

    for line in text_value.splitlines():
        value = clean_content(line, 900)
        if not value:
            continue
        marker = re.fullmatch(r"\[(Page|Slide) (\d+)\]", value)
        if marker:
            flush_buffer()
            container_locator = f"{marker.group(1)} {marker.group(2)}"
            locator = container_locator
            continue
        if value.startswith("[Heading]"):
            flush_buffer()
            heading = value.removeprefix("[Heading]").strip()
            locator = (
                f'{container_locator} · Heading "{heading}"'
                if container_locator != "Uploaded material"
                else f'Heading "{heading}"'
            )
            continue
        likely_plain_heading = (
            len(value) <= 90
            and len(value.split()) <= 9
            and not re.search(r"[.!?=]$", value)
            and sum(word[:1].isupper() for word in value.split())
            >= max(1, len(value.split()) // 2)
        )
        if likely_plain_heading:
            flush_buffer()
            locator = (
                f'{container_locator} · Heading "{value}"'
                if container_locator != "Uploaded material"
                else f'Heading "{value}"'
            )
            continue
        buffered_lines.append(value)
    flush_buffer()
    return passages


def source_sentences(text_value: str) -> list[str]:
    sentences = [
        clean_content(sentence, 600)
        for sentence in re.split(r"(?<=[.!?])\s+|\n+", text_value)
    ]
    return list(dict.fromkeys(
        sentence for sentence in sentences if 35 <= len(sentence) <= 600
    ))


CONCEPT_GROUNDING_PATTERNS = {
    "GM-AE": r"\b(algebra|expression|variable|coefficient)\b",
    "GM-LE": r"\b(linear equation|solve|unknown|equality)\b",
    "GM-FG": r"\b(function|graph|slope|coordinate)\b",
    "GM-SN": r"\b(scientific notation|power of ten|exponent)\b",
    "GM-UC": r"\b(convert|conversion|unit|dimensional)\b",
    "GM-TR": r"\b(trigon|sine|cosine|tangent|angle)\b",
    "GP-SV": r"\b(vector|scalar|component|resultant)\b",
    "GP-MK": r"\b(motion|kinematic|velocity|displacement|acceleration|speed)\b",
    "GP-NL": r"\b(force|newton|inertia|mass and acceleration)\b",
    "GP-WE": r"\b(work|energy|power|joule)\b",
    "GP-MI": r"\b(momentum|impulse|collision)\b",
}


def material_supports_concept(concept: Concept, text_value: str) -> bool:
    pattern = CONCEPT_GROUNDING_PATTERNS.get(concept.code)
    if pattern:
        return bool(re.search(pattern, text_value, re.I))
    terms = re.findall(r"[A-Za-z][A-Za-z'-]{3,}", concept.name)
    return any(re.search(rf"\b{re.escape(term)}\b", text_value, re.I) for term in terms)


def generated_question_data(
    sentences: list[str],
    index: int,
    question_type: str,
    concept: Concept,
    topic: str,
) -> tuple[str, str, list[dict[str, Any]]]:
    statement = sentences[index % len(sentences)]
    if question_type == "Multiple choice":
        prompt = (
            f"Based on the learning material, which statement best supports "
            f"{topic} in item {index + 1}?"
        )
        alternatives = [
            candidate
            for candidate in sentences
            if candidate.casefold() != statement.casefold()
        ]
        alternatives.extend(
            [
                f"{concept.name} is unrelated to the physical situation described.",
                "The relationship is independent of every variable named in the material.",
                "The measured quantities cannot be compared using the stated principle.",
            ]
        )
        unique_alternatives = []
        for candidate in alternatives:
            if candidate.casefold() not in {
                statement.casefold(),
                *(value.casefold() for value in unique_alternatives),
            }:
                unique_alternatives.append(candidate)
            if len(unique_alternatives) == 3:
                break
        choices = [{"text": statement, "is_correct": True}] + [
            {"text": candidate, "is_correct": False}
            for candidate in unique_alternatives
        ]
        return prompt, statement, choices
    if question_type == "True or false":
        return (
            f"True or false: {statement}",
            "True",
            [
                {"text": "True", "is_correct": True},
                {"text": "False", "is_correct": False},
            ],
        )
    if question_type == "Identification":
        return (
            f"Identify the physics concept aligned with this statement: {statement}",
            concept.name,
            [],
        )
    return (
        f"In 1–3 sentences, explain how this statement relates to {concept.name}: {statement}",
        statement,
        [],
    )


def generated_material_question(
    passages: list[dict[str, str]],
    index: int,
    question_type: str,
    concept: Concept,
    topic: str,
    include_calculations: bool,
) -> dict[str, Any]:
    passage = passages[index % len(passages)]
    statement = passage["text"]
    locator = passage["source_locator"]
    calculation = include_calculations and bool(
        re.search(
            r"[=≈]|\b(calculate|formula|velocity|force|energy|momentum|work|power)\b",
            statement,
            re.I,
        )
    )
    flags = ["Teacher review required: generated locally from extracted source text."]
    solution_steps = (
        f"Use the relationship at {locator}: identify quantities, convert units, "
        "substitute, solve, and check the result."
        if calculation
        else f"Find the supporting statement at {locator} and compare it with the response."
    )
    choices: list[dict[str, Any]] = []
    rationales: dict[str, str] = {}
    calculation_templates = {
        "GM-AE": (r"algebra|expression|variable", "Evaluate 2x + 5 when x = 4.", "13", ["18", "11", "8"], ["Substitute x = 4.", "Multiply 2(4) = 8.", "Add 5 to obtain 13."]),
        "GM-LE": (r"equation|unknown|solve", "Solve 3x + 6 = 21.", "x = 5", ["x = 9", "x = 7", "x = 3"], ["Subtract 6 from both sides: 3x = 15.", "Divide both sides by 3: x = 5."]),
        "GM-FG": (r"graph|slope|function", "Find the slope between (1, 3) and (5, 11).", "2", ["4", "8", "0.5"], ["Compute the change in y: 11 - 3 = 8.", "Compute the change in x: 5 - 1 = 4.", "Divide: slope = 8/4 = 2."]),
        "GM-SN": (r"scientific notation|power of ten|exponent", "Calculate (2 × 10³)(3 × 10²).", "6 × 10⁵", ["6 × 10⁶", "5 × 10⁵", "6 × 10¹"], ["Multiply coefficients: 2 × 3 = 6.", "Add exponents: 3 + 2 = 5.", "Write 6 × 10⁵."]),
        "GM-UC": (r"convert|conversion|km/h|unit", "Convert 72 km/h to m/s.", "20 m/s", ["72 m/s", "25.9 m/s", "2 m/s"], ["Multiply by 1000 m/1 km.", "Multiply by 1 h/3600 s.", "Cancel units and calculate 72/3.6 = 20 m/s."]),
        "GM-TR": (r"trigon|cos|sine|tangent", "Find the horizontal component of a 10 N vector at 60°.", "5 N", ["8.66 N", "10 N", "20 N"], ["Use Fx = F cos θ.", "Substitute Fx = 10 cos 60°.", "Calculate Fx = 5 N."]),
        "GP-SV": (r"vector|component|resultant", "Find the magnitude of perpendicular components 3 N and 4 N.", "5 N", ["7 N", "1 N", "12 N"], ["Use R = √(Rx² + Ry²).", "Substitute R = √(3² + 4²).", "Calculate R = √25 = 5 N."]),
        "GP-MK": (r"velocity|kinematic|displacement|acceleration", "An object travels 60 m in 5 s at constant velocity. Find its velocity.", "12 m/s", ["300 m/s", "55 m/s", "0.083 m/s"], ["Use v = d/t.", "Substitute v = 60 m / 5 s.", "Calculate v = 12 m/s."]),
        "GP-NL": (r"force|newton", "A 24 N net force acts on a 6 kg object. Find its acceleration.", "4 m/s²", ["144 m/s²", "18 m/s²", "0.25 m/s²"], ["Use ΣF = ma.", "Rearrange to a = ΣF/m.", "Substitute a = 24 N / 6 kg = 4 m/s²."]),
        "GP-WE": (r"work|energy|power|force", "A 20 N force moves an object 3 m in the force direction. Find the work.", "60 J", ["23 J", "6.67 J", "17 J"], ["Use W = Fd cos θ.", "The force is parallel, so cos 0° = 1.", "Calculate W = 20(3) = 60 J."]),
        "GP-MI": (r"momentum|impulse", "Find the momentum of a 2 kg cart moving at 5 m/s.", "10 kg·m/s", ["2.5 kg·m/s", "7 kg·m/s", "25 kg·m/s"], ["Use p = mv.", "Substitute p = 2 kg(5 m/s).", "Calculate p = 10 kg·m/s."]),
    }
    template = calculation_templates.get(concept.code)
    if include_calculations and question_type == "Multiple choice" and template and index == 0:
        support_pattern, calculation_prompt, answer, distractor_values, steps = template
        if re.search(support_pattern, statement, re.I):
            choices = [{"text": answer, "is_correct": True}] + [
                {"text": value, "is_correct": False} for value in distractor_values
            ]
            rationales = {
                answer: "Correct: the cited relationship is applied with compatible units.",
                distractor_values[0]: "Common error: multiplying or combining values without isolating the requested quantity.",
                distractor_values[1]: "Common error: using the wrong arithmetic operation or relationship.",
                distractor_values[2]: "Common error: reversing a ratio or omitting a required exponent.",
            }
            return {
                "prompt": f"Using the relationship supported at {locator}, {calculation_prompt}",
                "correct_answer": answer,
                "choices": choices,
                "source_locator": locator,
                "solution_steps": " ".join(steps),
                "validation_flags": flags,
                "distractor_rationales": rationales,
                "is_calculation": True,
            }
        flags.append("No solvable calculation aligned with this source passage; generated a conceptual item instead.")
        calculation = False
    elif calculation:
        calculation = False
        solution_steps = f"Find the supporting statement at {locator} and compare it with the response."
    if question_type == "Multiple choice":
        focus_terms = [
            word
            for word in re.findall(r"[A-Za-z][A-Za-z'-]{3,}", statement)
            if word.casefold() not in {"that", "this", "with", "from", "using", "material"}
        ][:3]
        focus = ", ".join(focus_terms) or topic
        prompt = f"Which statement matches the uploaded material's discussion of {focus}?"
        alternatives = [
            f"{concept.name} uses the named quantities, but their relationship is reversed.",
            f"The material treats {concept.name} as independent of the quantities it names.",
            f"The relationship can be applied without checking units or the stated conditions.",
        ]
        distractors = []
        for candidate in alternatives:
            normalized = candidate.casefold()
            if normalized != statement.casefold() and all(
                normalized != item.casefold() for item in distractors
            ):
                distractors.append(candidate)
            if len(distractors) == 3:
                break
        choices = [{"text": statement, "is_correct": True}] + [
            {"text": candidate, "is_correct": False} for candidate in distractors
        ]
        correct_answer = statement
        rationales = {
            statement: "Correct: directly supported by the cited material.",
            **{
                item: "Distractor: not supported by the cited passage or changes its relationship."
                for item in distractors
            },
        }
    elif question_type == "True or false":
        prompt = f"True or false, according to {locator}: {statement}"
        correct_answer = "True"
        choices = [
            {"text": "True", "is_correct": True},
            {"text": "False", "is_correct": False},
        ]
        rationales = {
            "True": "Correct: this statement appears in the cited source.",
            "False": "Distractor: it contradicts the cited source.",
        }
        calculation = False
    elif question_type == "Identification":
        prompt = f"Identify the concept connected to this material statement: {statement}"
        correct_answer = concept.name
        calculation = False
    else:
        prompt = f"In 1–3 sentences, explain how this statement relates to {concept.name}: {statement}"
        correct_answer = statement
    return {
        "prompt": prompt,
        "correct_answer": correct_answer,
        "choices": choices,
        "source_locator": locator,
        "solution_steps": solution_steps,
        "validation_flags": flags,
        "distractor_rationales": rationales,
        "is_calculation": calculation,
    }


def validate_generated_question(
    data: dict[str, Any],
    concept: Concept,
    competency: str,
    prior_prompts: set[str],
) -> tuple[str, list[str]]:
    flags = list(data["validation_flags"])
    prompt_key = re.sub(r"\W+", " ", data["prompt"].casefold()).strip()
    if prompt_key in prior_prompts:
        flags.append("Potential repeated question idea.")
    if concept.name.casefold() not in data["prompt"].casefold() and not competency.strip():
        flags.append("Concept or competency alignment needs confirmation.")
    if not data["correct_answer"].strip():
        flags.append("Correct answer is missing.")
    choices = data["choices"]
    if choices:
        if sum(1 for choice in choices if choice["is_correct"]) != 1:
            flags.append("The item must have exactly one correct choice.")
        normalized = [choice["text"].strip().casefold() for choice in choices]
        if len(normalized) != len(set(normalized)):
            flags.append("Duplicate answer choices detected.")
    if data["is_calculation"] and not re.search(r"[=≈]|\d", data["correct_answer"]):
        flags.append("Calculation solvability needs teacher confirmation.")
    status_value = (
        "Ready for review"
        if len(flags) == 1 and flags[0].startswith("Teacher review required")
        else "Needs review"
    )
    return status_value, flags


def ai_configuration() -> dict[str, Any]:
    api_key = os.getenv("AI_API_KEY", "").strip() or os.getenv("OPENAI_API_KEY", "").strip()
    provider = os.getenv("AI_PROVIDER", "").strip()
    model = os.getenv("AI_MODEL", "").strip()
    base_url = os.getenv("AI_BASE_URL", "https://api.openai.com/v1").strip().rstrip("/")
    configured = bool(api_key and model and provider)
    return {
        "configured": configured,
        "provider": provider or None,
        "model": model or None,
        "base_url": base_url if configured else None,
        "api_key": api_key,
    }


def request_ai_question_generation(
    document: UploadedDocument,
    payload: DocumentGenerationInput,
    concepts: list[Concept],
) -> dict[str, Any]:
    configuration = ai_configuration()
    if not configuration["configured"]:
        raise HTTPException(
            status_code=503,
            detail=(
                "AI question generation is not configured. Set AI_PROVIDER, AI_MODEL, "
                "and AI_API_KEY (or OPENAI_API_KEY) in the backend environment."
            ),
        )
    if configuration["provider"] == "test-grounded" and configuration["model"] == "deterministic-test-double":
        passages = source_passages(document.extracted_text)
        questions = []
        for index in range(payload.number_of_questions):
            concept = concepts[index % len(concepts)]
            requested_type = payload.question_type
            question_type = (
                ["Multiple choice", "True or false", "Short answer", "Problem solving"][index % 4]
                if requested_type == "Mixed"
                else requested_type
            )
            data = generated_material_question(
                passages, index, question_type, concept, payload.topic, payload.include_calculations
            )
            questions.append({
                "concept_id": concept.id,
                "prompt": data["prompt"],
                "question_type": question_type,
                "correct_answer": data["correct_answer"],
                "choices": data["choices"],
                "explanation": f"The answer is supported by {data['source_locator']}.",
                "hint": f"Review {data['source_locator']}.",
                "solution_steps": data["solution_steps"],
                "source_locator": data["source_locator"],
                "learning_competency": payload.learning_competency,
                "estimated_cognitive_demand": 0.5,
                "distractor_rationales": data["distractor_rationales"],
                "is_calculation": data["is_calculation"],
                "validation": {key: True for key in (
                    "answerable", "answer_correct", "solution_consistent", "formula_units_valid",
                    "source_supported", "distractors_valid", "not_duplicate",
                    "difficulty_followed", "no_unrelated_topic",
                )},
                "validation_notes": ["Teacher review required: test provider output."],
            })
        return {"analysis": document.analysis or {}, "questions": questions}
    if APP_ENV == "production" and not configuration["base_url"].startswith("https://"):
        raise HTTPException(
            status_code=500,
            detail="The AI provider must use HTTPS in production.",
        )
    settings = payload.model_dump()
    concept_data = [
        {"id": item.id, "code": item.code, "name": item.name, "description": item.description}
        for item in concepts
    ]
    source_text = document.extracted_text[: int(os.getenv("AI_SOURCE_CHAR_LIMIT", "60000"))]
    system_prompt = (
        "You are an educational assessment author and verifier. Analyze only the supplied "
        "learning material, then create original questions grounded in it. Return strict JSON. "
        "Never introduce facts absent from the source. Verify answers, units, formulas, "
        "distractors, difficulty, uniqueness, and solution consistency before returning."
    )
    user_prompt = {
        "task": "Analyze the module and generate source-grounded questions with a quality-review pass.",
        "settings": settings,
        "allowed_concepts": concept_data,
        "existing_structural_analysis": document.analysis or {},
        "required_output": {
            "analysis": {
                "detected_subject": "string",
                "module_title": "string",
                "main_topics": ["string"],
                "subtopics": ["string"],
                "learning_objectives": ["string"],
                "important_concepts": ["string"],
                "definitions": ["string"],
                "formulas": ["string"],
                "formula_variables": ["string"],
                "worked_examples": ["string"],
                "prerequisites": ["string"],
                "estimated_learner_level": "string",
                "estimated_difficulty": "string",
                "misconceptions": ["string"],
                "pages_used_as_evidence": [1],
            },
            "questions": [{
                "concept_id": "integer from allowed_concepts",
                "prompt": "string",
                "question_type": "Multiple choice, True or false, Identification, Short answer, or Problem solving",
                "correct_answer": "string",
                "choices": [{"text": "string", "is_correct": True}],
                "explanation": "string",
                "hint": "string",
                "solution_steps": "scaffolded solution string",
                "source_locator": "Page N or Page N - section",
                "learning_competency": "string",
                "estimated_cognitive_demand": 0.5,
                "distractor_rationales": {"choice text": "rationale"},
                "is_calculation": False,
                "validation": {
                    "answerable": True,
                    "answer_correct": True,
                    "solution_consistent": True,
                    "formula_units_valid": True,
                    "source_supported": True,
                    "distractors_valid": True,
                    "not_duplicate": True,
                    "difficulty_followed": True,
                    "no_unrelated_topic": True,
                },
                "validation_notes": ["string"],
            }],
        },
        "learning_material": source_text,
    }
    try:
        import httpx

        response = httpx.post(
            f"{configuration['base_url']}/chat/completions",
            headers={"Authorization": f"Bearer {configuration['api_key']}"},
            json={
                "model": configuration["model"],
                "temperature": 0.2,
                "response_format": {"type": "json_object"},
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": json.dumps(user_prompt, ensure_ascii=False)},
                ],
            },
            timeout=float(os.getenv("AI_REQUEST_TIMEOUT_SECONDS", "90")),
        )
        response.raise_for_status()
        content = response.json()["choices"][0]["message"]["content"]
        result = json.loads(content)
    except HTTPException:
        raise
    except Exception as error:
        logger.exception("AI question generation failed")
        raise HTTPException(
            status_code=502,
            detail=f"The configured AI service could not generate validated questions: {type(error).__name__}",
        )
    if not isinstance(result, dict) or not isinstance(result.get("questions"), list):
        raise HTTPException(status_code=502, detail="The AI service returned an invalid question payload")
    return result


def replace_question_choices(
    db: Session, question: Question, choices: list[Any]
) -> None:
    for choice in list(
        db.scalars(
            select(AnswerChoice).where(AnswerChoice.question_id == question.id)
        )
    ):
        db.delete(choice)
    db.flush()
    for index, choice in enumerate(choices, start=1):
        text_value = (
            choice.text if hasattr(choice, "text") else str(choice["text"])
        )
        is_correct = (
            choice.is_correct
            if hasattr(choice, "is_correct")
            else bool(choice["is_correct"])
        )
        misconception_id = (
            getattr(choice, "misconception_id", None)
            if hasattr(choice, "text")
            else choice.get("misconception_id")
        )
        misconception_confidence = (
            getattr(choice, "misconception_confidence", None)
            if hasattr(choice, "text")
            else choice.get("misconception_confidence")
        )
        mapping_status = (
            getattr(choice, "mapping_status", "Unreviewed")
            if hasattr(choice, "text")
            else choice.get("mapping_status", "Unreviewed")
        )
        if misconception_id:
            rule = db.get(Misconception, misconception_id)
            if not rule or rule.concept_id != question.concept_id:
                raise HTTPException(
                    status_code=422,
                    detail="A distractor misconception must belong to the question competency",
                )
            if is_correct:
                raise HTTPException(
                    status_code=422,
                    detail="The correct choice cannot be mapped to a misconception",
                )
            if mapping_status in {"Teacher reviewed", "Validated"} and misconception_confidence is None:
                raise HTTPException(
                    status_code=422,
                    detail="Reviewed misconception mappings require a confidence value",
                )
        db.add(
            AnswerChoice(
                question_id=question.id,
                text=clean_content(text_value, 3000),
                is_correct=is_correct,
                position=index,
                misconception_id=misconception_id,
                misconception_confidence=misconception_confidence,
                mapping_status=("Validated" if is_correct else mapping_status),
            )
        )


def apply_question_bank_input(
    question: Question, payload: QuestionBankInput
) -> None:
    question.concept_id = payload.concept_id
    question.prompt = clean_content(payload.prompt)
    question.question_type = payload.question_type
    question.correct_answer = clean_content(payload.correct_answer, 3000)
    question.explanation = clean_content(payload.explanation)
    question.feedback = question.explanation
    question.hint = clean_content(payload.hint, 3000)
    question.difficulty_label = clean_content(payload.difficulty, 20)
    question.cognitive_level = clean_content(payload.cognitive_level, 40)
    question.subject = clean_content(payload.subject, 100)
    question.topic = clean_content(payload.topic, 160)
    question.learning_competency = clean_content(
        payload.learning_competency, 1000
    )
    question.points = payload.points
    question.source_locator = clean_content(payload.source_locator, 300)
    question.solution_steps = clean_content(payload.solution_steps)
    question.solution_structure = payload.solution_structure
    question.estimated_cognitive_demand = payload.estimated_cognitive_demand
    question.prerequisite_concept_id = payload.prerequisite_concept_id
    question.validation_status = clean_content(payload.validation_status, 30)
    question.validation_flags = [
        clean_content(flag, 250) for flag in payload.validation_flags
    ]
    question.distractor_rationales = {
        clean_content(key, 500): clean_content(value, 1000)
        for key, value in payload.distractor_rationales.items()
    }
    question.is_calculation = payload.is_calculation
    question.status = payload.status


def utc_value(value: datetime | None) -> datetime | None:
    if value is None:
        return None
    return value if value.tzinfo else value.replace(tzinfo=timezone.utc)


def refresh_assessment_status(db: Session, assessment: Assessment) -> None:
    now = datetime.now(timezone.utc)
    available_from = utc_value(assessment.available_from)
    due_at = utc_value(assessment.due_at)
    activity = db.get(Activity, assessment.activity_id) if assessment.activity_id else None
    if assessment.status == "Scheduled" and (
        available_from is None or available_from <= now
    ):
        assessment.status = "Published"
        assessment.published_at = now
        if activity:
            activity.active = True
    if assessment.status == "Published" and due_at and due_at < now:
        assessment.status = "Closed"
        if activity:
            activity.active = False


def assessment_assigned_to_student(
    db: Session, assessment: Assessment, student: User
) -> bool:
    assignments = list(
        db.scalars(
            select(AssessmentAssignment).where(
                AssessmentAssignment.assessment_id == assessment.id
            )
        )
    )
    if not assignments:
        return True
    profile = db.scalar(
        select(StudentProfile).where(StudentProfile.user_id == student.id)
    )
    return any(
        assignment.student_id == student.id
        or (
            assignment.section
            and profile
            and profile.section
            and assignment.section.casefold() == profile.section.casefold()
        )
        for assignment in assignments
    )


def onboarding_payload(profile: StudentProfile | None) -> dict[str, Any]:
    completed_at = profile.onboarding_completed_at if profile else None
    return {
        "completed": completed_at is not None,
        "completed_at": completed_at,
        "version": profile.onboarding_version if profile else "1.0",
    }


def onboarding_diagnostic_payload(db: Session, student: User) -> dict[str, Any]:
    activity = db.scalar(
        select(Activity).where(
            Activity.is_onboarding_diagnostic.is_(True),
            Activity.active.is_(True),
        )
    )
    if not activity:
        return {
            "available": False,
            "item_count": 0,
            "completed": False,
            "analysis_complete": False,
        }
    item_count = db.scalar(
        select(func.count(Question.id)).where(
            Question.activity_id == activity.id,
            Question.active.is_(True),
            func.lower(Question.question_type) == "multiple choice",
        )
    ) or 0
    attempt = db.scalar(
        select(AssessmentAttempt)
        .where(
            AssessmentAttempt.student_id == student.id,
            AssessmentAttempt.activity_id == activity.id,
        )
        .order_by(AssessmentAttempt.submitted_at.desc())
    )
    effort = (
        db.scalar(
            select(MentalEffortRating).where(
                MentalEffortRating.attempt_id == attempt.id
            )
        )
        if attempt
        else None
    )
    gaps = list(
        db.execute(
            select(LearningGap, Concept)
            .join(Concept, Concept.id == LearningGap.concept_id)
            .where(
                LearningGap.student_id == student.id,
                LearningGap.resolved_at.is_(None),
            )
            .order_by(LearningGap.mastery_score, Concept.name)
            .limit(4)
        )
    )
    return {
        "available": item_count == ONBOARDING_DIAGNOSTIC_ITEMS,
        "activity_id": activity.id,
        "title": activity.title,
        "item_count": item_count,
        "completed": attempt is not None,
        "analysis_complete": effort is not None,
        "pending_mental_effort_attempt_id": attempt.id if attempt and not effort else None,
        "mental_effort_boundaries": {
            "low_max": int(get_setting(db, "mental_effort_low_max")),
            "moderate_max": int(get_setting(db, "mental_effort_moderate_max")),
        },
        "latest_result": (
            {
                "attempt_id": attempt.id,
                "score": attempt.score,
                "max_score": attempt.max_score,
                "accuracy": attempt.accuracy,
                "submitted_at": attempt.submitted_at,
                "mental_effort": effort.rating if effort else None,
                "cognitive_load_category": effort.category if effort else "Pending",
            }
            if attempt
            else None
        ),
        "priority_gaps": [
            {
                "concept_id": concept.id,
                "concept": concept.name,
                "mastery_score": gap.mastery_score,
            }
            for gap, concept in gaps
        ],
    }


def dashboard_explainability(
    db: Session,
    student: User,
    profile: StudentProfile | None,
    mastery: dict[int, MasteryRecord],
    gaps: list[LearningGap],
    pathway: PathwayRecommendation | None,
    pathway_data: dict[str, Any] | None,
) -> dict[str, Any]:
    threshold = float(get_setting(db, "mastery_threshold"))
    mastery_rows = sorted(
        mastery.values(), key=lambda record: (db.get(Concept, record.concept_id).name, record.id)
    )
    concept_evidence: list[dict[str, Any]] = []
    for record in mastery_rows:
        concept = db.get(Concept, record.concept_id)
        evidence_rows = db.execute(
            select(
                AssessmentAttempt.id,
                AssessmentAttempt.submitted_at,
                func.sum(ItemResponse.earned_points),
                func.sum(ItemResponse.max_points),
            )
            .join(ItemResponse, ItemResponse.attempt_id == AssessmentAttempt.id)
            .join(Question, Question.id == ItemResponse.question_id)
            .where(
                AssessmentAttempt.student_id == student.id,
                Question.concept_id == record.concept_id,
            )
            .group_by(AssessmentAttempt.id, AssessmentAttempt.submitted_at)
            .order_by(AssessmentAttempt.submitted_at)
        ).all()
        attempts = []
        for position, (attempt_id, submitted_at, earned, maximum) in enumerate(
            evidence_rows, start=1
        ):
            earned_value = float(earned or 0)
            maximum_value = float(maximum or 0)
            attempts.append(
                {
                    "attempt_id": attempt_id,
                    "earned": earned_value,
                    "maximum": maximum_value,
                    "mastery": earned_value / maximum_value if maximum_value else None,
                    "weight": position if record.calculation_mode == "weighted" else None,
                    "submitted_at": submitted_at,
                }
            )
        concept_evidence.append(
            {
                "concept_id": concept.id,
                "concept": concept.name,
                "score": record.mastery_score,
                "classification": record.classification,
                "calculation_mode": record.calculation_mode,
                "attempts": attempts,
                "latest_evidence_at": record.created_at,
                "below_threshold": record.mastery_score < threshold,
            }
        )
    mastery_average = (
        sum(record.mastery_score for record in mastery_rows) / len(mastery_rows)
        if mastery_rows
        else None
    )
    below_threshold = [row for row in concept_evidence if row["below_threshold"]]
    score_breakdown_count = sum(bool(row["attempts"]) for row in concept_evidence)
    latest_mastery_at = max(
        (record.created_at for record in mastery_rows), default=None
    )
    if mastery_average is None:
        mastery_interpretation = (
            "NeuroLearn-X needs more completed learning activities before it can "
            "calculate this reliably."
        )
    elif mastery_average >= threshold:
        mastery_interpretation = (
            f"Your Average Mastery is {mastery_average:.0%}. Your assessed concepts "
            "show strong current evidence overall; individual concepts are still checked separately."
        )
    elif mastery_average >= threshold * 0.7:
        mastery_interpretation = (
            f"Your Average Mastery is {mastery_average:.0%}. You are showing partial "
            "understanding, and some concepts still need practice."
        )
    else:
        mastery_interpretation = (
            f"Your Average Mastery is {mastery_average:.0%}. NeuroLearn-X will provide "
            "more guided prerequisite practice while you build stronger evidence."
        )

    latest_interaction = db.scalar(
        select(InteractionLog)
        .where(InteractionLog.student_id == student.id)
        .order_by(InteractionLog.submission_time.desc())
    )
    prediction = db.scalar(
        select(CognitiveLoadPrediction)
        .where(CognitiveLoadPrediction.student_id == student.id)
        .order_by(CognitiveLoadPrediction.evidence_date.desc(), CognitiveLoadPrediction.id.desc())
    )
    latest_effort = db.scalar(
        select(MentalEffortRating)
        .where(MentalEffortRating.student_id == student.id)
        .order_by(MentalEffortRating.created_at.desc())
    )
    prediction_is_current = bool(
        prediction
        and (
            latest_interaction is None
            or utc_value(prediction.evidence_date) >= utc_value(latest_interaction.submission_time)
        )
    )
    prediction_version = db.get(ModelVersion, prediction.model_version_id) if prediction else None
    probabilities = prediction.probabilities if prediction_is_current and prediction else {}
    model_count = len(
        (prediction_version.metadata_json or {}).get("ensemble_members", [])
    ) if prediction_version else 0
    if prediction_is_current and prediction:
        load_interpretation = (
            f"Your model-predicted cognitive load is {prediction.predicted_category}. "
            + {
                "Low": "NeuroLearn-X can offer efficient review or an optional challenge while monitoring mastery.",
                "Moderate": "NeuroLearn-X will provide guided support while keeping the activity challenging enough for learning.",
                "High": "NeuroLearn-X will favor shorter steps and stronger scaffolding before increasing difficulty.",
            }.get(prediction.predicted_category, "NeuroLearn-X will continue monitoring your learning evidence.")
        )
        prediction_quality = (
            "All required model outputs are available."
            if not prediction.missing_features
            else "The prediction is available, but some optional evidence was missing: "
            + ", ".join(str(item).replace("_", " ") for item in prediction.missing_features)
            + "."
        )
    else:
        load_interpretation = (
            "NeuroLearn-X needs more completed learning activities before it can "
            "calculate this reliably."
        )
        prediction_quality = (
            "A newer completed activity exists, so the older model prediction is not displayed."
            if prediction and latest_interaction
            else "No current learner-level ensemble prediction has been recorded."
        )

    target = db.get(Concept, profile.target_concept_id) if profile and profile.target_concept_id else None
    target_mastery = mastery.get(target.id) if target else None
    target_gap = next((gap for gap in gaps if target and gap.concept_id == target.id), None)
    prerequisite_rows = (
        list(
            db.execute(
                select(Concept)
                .join(
                    PrerequisiteEdge,
                    PrerequisiteEdge.prerequisite_concept_id == Concept.id,
                )
                .where(
                    PrerequisiteEdge.succeeding_concept_id == target.id,
                    PrerequisiteEdge.active.is_(True),
                )
                .order_by(Concept.difficulty, Concept.name)
            ).scalars()
        )
        if target
        else []
    )
    target_prerequisites = [
        {
            "concept_id": concept.id,
            "concept": concept.name,
            "mastery": mastery[concept.id].mastery_score if concept.id in mastery else None,
            "below_threshold": (
                concept.id not in mastery or mastery[concept.id].mastery_score < threshold
            ),
        }
        for concept in prerequisite_rows
    ]
    target_reason = None
    if target:
        if target_gap and target_mastery:
            target_reason = (
                f"{target.name} is the active target because its current mastery is "
                f"{target_mastery.mastery_score:.0%}, below the configured {threshold:.0%} threshold."
            )
        elif target_mastery:
            target_reason = (
                f"{target.name} is the active focus saved in your learner profile. Its current "
                f"mastery is {target_mastery.mastery_score:.0%}; the pathway still checks connected prerequisites."
            )
        else:
            target_reason = (
                f"{target.name} is the active focus saved in your learner profile, but no valid "
                "mastery evidence has been recorded for it yet."
            )

    required_steps = [step for step in (pathway_data or {}).get("steps", []) if step["required"]]
    completed_steps = [step for step in required_steps if step["completed_at"]]
    remaining_steps = [step for step in required_steps if not step["completed_at"]]
    progress_percent = (
        len(completed_steps) / len(required_steps) if required_steps else None
    )
    latest_progress_at = max(
        [step["completed_at"] for step in completed_steps]
        + ([pathway.created_at] if pathway else []),
        default=None,
    )
    next_step = remaining_steps[0] if remaining_steps else None
    next_concept = db.get(Concept, next_step["concept_id"]) if next_step else None
    next_gap = next(
        (gap for gap in gaps if next_step and gap.concept_id == next_step["concept_id"]),
        None,
    )
    next_prerequisites = []
    if next_concept:
        next_prerequisites = [
            concept.name
            for concept in db.execute(
                select(Concept)
                .join(
                    PrerequisiteEdge,
                    PrerequisiteEdge.prerequisite_concept_id == Concept.id,
                )
                .where(
                    PrerequisiteEdge.succeeding_concept_id == next_concept.id,
                    PrerequisiteEdge.active.is_(True),
                )
                .order_by(Concept.name)
            ).scalars()
        ]
    decision = (pathway.decision_explanation or {}) if pathway else {}
    weights = decision.get("weights") or {
        "alpha": float(get_setting(db, "alpha")),
        "beta": float(get_setting(db, "beta")),
        "gamma": float(get_setting(db, "gamma")),
    }
    aps_available = bool(
        pathway
        and pathway.source_type == "Automatic"
        and decision.get("adaptive_pathway_score") is not None
    )
    difficulty_labels = {1: "Foundational", 2: "Developing", 3: "Moderate", 4: "Advanced", 5: "Challenging"}
    next_activity = db.get(Activity, next_step["activity_id"]) if next_step else None

    return {
        "mastery_threshold": threshold,
        "average_mastery": {
            "available": mastery_average is not None,
            "value": mastery_average,
            "meaning": "Average Mastery shows how well you currently understand the concepts measured by your completed assessments and learning activities.",
            "why_matters": "It summarizes current concept evidence while keeping each concept's mastery result separate.",
            "interpretation": mastery_interpretation,
            "concepts": concept_evidence,
            "concepts_below_threshold": below_threshold,
            "sum_mastery": sum(record.mastery_score for record in mastery_rows),
            "concept_count": len(mastery_rows),
            "score_breakdown_count": score_breakdown_count,
            "threshold": threshold,
            "latest_evidence_at": latest_mastery_at,
            "pathway_effect": (
                f"{len(below_threshold)} concept(s) are below the {threshold:.0%} threshold and may be prioritized as learning gaps."
                if mastery_rows
                else "No mastery value is sent to the pathway until valid assessment evidence exists."
            ),
            "data_quality": (
                f"Calculated from {len(mastery_rows)} concept(s) with valid saved mastery evidence. "
                + (
                    "Item-score breakdowns are available for every concept."
                    if score_breakdown_count == len(mastery_rows)
                    else f"Item-score breakdowns are available for {score_breakdown_count}; the remaining {len(mastery_rows) - score_breakdown_count} use their saved mastery records because older item-level provenance is unavailable."
                )
                if mastery_rows
                else "NeuroLearn-X needs more completed learning activities before it can calculate this reliably."
            ),
            "related_path": "/student/mastery",
        },
        "model_predicted_cognitive_load": {
            "available": prediction_is_current,
            "category": prediction.predicted_category if prediction_is_current and prediction else None,
            "probabilities": probabilities,
            "index": prediction.expected_index if prediction_is_current and prediction else None,
            "confidence": prediction.confidence if prediction_is_current and prediction else None,
            "model_version": prediction_version.version if prediction_is_current and prediction_version else None,
            "model_count": model_count,
            "meaning": "Predicted Cognitive Load estimates how mentally demanding your recent learning activities may have been based on your performance and behavior.",
            "why_matters": "It helps NeuroLearn-X choose an appropriate amount of scaffolding without treating effort as ability.",
            "interpretation": load_interpretation,
            "evidence": prediction.evidence if prediction_is_current and prediction else None,
            "reported_mental_effort": (
                {
                    "rating": latest_effort.rating,
                    "category": latest_effort.category,
                    "reported_at": latest_effort.created_at,
                }
                if latest_effort
                else None
            ),
            "latest_evidence_at": prediction.evidence_date if prediction_is_current and prediction else None,
            "pathway_effect": (
                prediction.recommended_action
                if prediction_is_current and prediction
                else "No model-predicted load is applied as a learner-level dashboard result until a current prediction is available."
            ),
            "data_quality": prediction_quality,
            "disclaimer": "This is an educational model prediction for learning support, not a medical or psychological diagnosis.",
            "related_path": "/student/history",
        },
        "current_target": {
            "available": target is not None,
            "concept": concept_payload(target) if target else None,
            "mastery": target_mastery.mastery_score if target_mastery else None,
            "threshold": threshold,
            "detected_gap": target_gap is not None,
            "gap_amount": (
                max(0.0, threshold - target_mastery.mastery_score)
                if target_mastery and target_gap
                else None
            ),
            "prerequisites": target_prerequisites,
            "reason": target_reason,
            "meaning": "Current Target is the concept NeuroLearn-X recommends that you focus on now.",
            "why_matters": "A clear target keeps the pathway focused on the knowledge needed for your next goal.",
            "interpretation": target_reason or "Choose a General Physics target to begin personalized recommendations.",
            "latest_evidence_at": target_mastery.created_at if target_mastery else None,
            "pathway_effect": (
                "The active pathway follows the prerequisite graph toward this target and prioritizes unmastered connected concepts."
                if target
                else "A pathway cannot be generated until a target competency is selected."
            ),
            "data_quality": (
                "Target and mastery values come from the learner profile and latest valid mastery record."
                if target and target_mastery
                else "The target is available, but more completed learning activities are needed for a reliable mastery comparison."
                if target
                else "No target competency is currently selected."
            ),
            "related_path": "/student/targets",
        },
        "pathway_progress": {
            "available": bool(pathway and required_steps),
            "completed": len(completed_steps),
            "total": len(required_steps),
            "remaining": len(remaining_steps),
            "percentage": progress_percent,
            "current_activity": next_step["activity"] if next_step else None,
            "steps": [
                {
                    "activity": step["activity"],
                    "concept": step["concept"],
                    "status": "Completed" if step["completed_at"] else "Current" if next_step and step["id"] == next_step["id"] else "Remaining",
                    "completed_at": step["completed_at"],
                }
                for step in required_steps
            ],
            "meaning": "Pathway Progress shows how much of your current personalized learning pathway you have successfully completed.",
            "why_matters": "It shows completed evidence and what remains without counting an activity before its completion requirement is met.",
            "interpretation": (
                f"You completed {len(completed_steps)} of {len(required_steps)} required activities, so your Pathway Progress is {progress_percent:.1%}. {len(remaining_steps)} activit{'y' if len(remaining_steps) == 1 else 'ies'} remain."
                if progress_percent is not None
                else "NeuroLearn-X needs an active pathway with required activities before progress can be calculated."
            ),
            "latest_evidence_at": latest_progress_at,
            "pathway_effect": "The pathway may update when new results reveal a learning gap, a concept reaches mastery, or predicted cognitive load changes.",
            "data_quality": (
                "Only required steps with saved completion evidence are counted as completed."
                if required_steps
                else "NeuroLearn-X needs more completed learning activities before it can calculate this reliably."
            ),
            "related_path": "/student/pathway",
        },
        "next_recommended_step": {
            "available": next_step is not None,
            "activity": next_step["activity"] if next_step else None,
            "concept": next_step["concept"] if next_step else None,
            "learning_gap": (
                {
                    "concept": next_concept.name,
                    "mastery": next_gap.mastery_score,
                    "threshold": next_gap.threshold,
                }
                if next_gap and next_concept
                else None
            ),
            "prerequisites": next_prerequisites,
            "estimated_minutes": next_step["estimated_minutes"] if next_step else None,
            "difficulty": difficulty_labels.get(next_activity.difficulty, f"Level {next_activity.difficulty}") if next_activity else None,
            "predicted_load_index": next_step["predicted_load_index"] if next_step else None,
            "selection_reason": next_step["selection_reason"] if next_step else None,
            "meaning": "Next Recommended Step is the activity NeuroLearn-X considers most suitable for your current learning needs.",
            "why_matters": "It turns current mastery, prerequisite, load, and time evidence into one clear next action.",
            "interpretation": (
                next_step["selection_reason"]
                if next_step
                else "No required activity is currently waiting in the active pathway."
            ),
            "aps": (
                {
                    "available": True,
                    "gap_coverage": pathway.gap_coverage,
                    "predicted_cognitive_load": pathway.predicted_cognitive_load,
                    "normalized_learning_time": pathway.normalized_learning_time,
                    "score": pathway.adaptive_pathway_score,
                    "weights": weights,
                    "alternatives": decision.get("alternatives_not_selected", []),
                    "selection_reason": decision.get("selection_reason"),
                }
                if aps_available and pathway
                else {
                    "available": False,
                    "reason": (
                        "This is a teacher-assigned pathway, so an Adaptive Pathway Score was not used to override the teacher's selection."
                        if pathway and pathway.source_type == "Teacher"
                        else "No valid automatic pathway candidate comparison is available."
                    ),
                }
            ),
            "latest_evidence_at": pathway.created_at if pathway else None,
            "pathway_effect": "Completing this activity with the existing evidence requirement can update mastery, resolve or add gaps, and recalculate the remaining route.",
            "data_quality": (
                "Recommendation fields come from the current required pathway step and its saved selection evidence."
                if next_step
                else "NeuroLearn-X needs an active pathway before it can recommend a next step."
            ),
            "related_path": "/student/pathway",
        },
    }


def student_snapshot(db: Session, student: User) -> dict[str, Any]:
    mastery = latest_mastery_map(db, student.id)
    concepts = {row.id: row for row in db.scalars(select(Concept))}
    gaps = list(
        db.scalars(
            select(LearningGap).where(
                LearningGap.student_id == student.id, LearningGap.resolved_at.is_(None)
            )
        )
    )
    pathway = db.scalar(
        select(PathwayRecommendation)
        .where(
            PathwayRecommendation.student_id == student.id,
            PathwayRecommendation.active.is_(True),
            PathwayRecommendation.selected.is_(True),
        )
        .order_by(PathwayRecommendation.created_at.desc())
    )
    if pathway and pathway.source_type == "Automatic":
        pathway_steps = list(
            db.scalars(
                select(PathwayStep).where(PathwayStep.pathway_id == pathway.id)
            )
        )
        if pathway_steps and any(not step.content for step in pathway_steps):
            generate_pathways(db, student)
            pathway = db.scalar(
                select(PathwayRecommendation)
                .where(
                    PathwayRecommendation.student_id == student.id,
                    PathwayRecommendation.active.is_(True),
                    PathwayRecommendation.selected.is_(True),
                )
                .order_by(PathwayRecommendation.created_at.desc())
            )
    profile = db.scalar(select(StudentProfile).where(StudentProfile.user_id == student.id))
    attempts = list(
        db.scalars(
            select(AssessmentAttempt)
            .where(AssessmentAttempt.student_id == student.id)
            .order_by(AssessmentAttempt.submitted_at.desc())
            .limit(12)
        )
    )
    ratings = {
        row.attempt_id: row
        for row in db.scalars(
            select(MentalEffortRating).where(MentalEffortRating.student_id == student.id)
        )
    }
    mastery_rows = [
        {
            "concept_id": concept_id,
            "concept": concepts[concept_id].name,
            "subject": concepts[concept_id].subject,
            "score": record.mastery_score,
            "classification": record.classification,
            "updated_at": record.created_at,
        }
        for concept_id, record in mastery.items()
        if concept_id in concepts
    ]
    pathway_data = serialize_pathway(db, pathway) if pathway else None
    pathway_history = [
        serialize_pathway(db, item)
        for item in db.scalars(
            select(PathwayRecommendation)
            .where(PathwayRecommendation.student_id == student.id)
            .order_by(PathwayRecommendation.created_at.desc())
            .limit(12)
        )
    ]
    required_steps = (
        [step for step in pathway_data["steps"] if step["required"]]
        if pathway_data
        else []
    )
    completed = sum(1 for step in required_steps if step["completed_at"])
    total = len(required_steps)
    explainability = dashboard_explainability(
        db,
        student,
        profile,
        mastery,
        gaps,
        pathway,
        pathway_data,
    )
    return {
        "onboarding": onboarding_payload(profile),
        "diagnostic": onboarding_diagnostic_payload(db, student),
        "explainability": explainability,
        "student": {
            **user_payload(student),
            "grade_level": profile.grade_level if profile else None,
            "section": profile.section if profile else None,
        },
        "target": concept_payload(concepts[profile.target_concept_id])
        if profile and profile.target_concept_id
        else None,
        "mastery": mastery_rows,
        "mastery_average": (
            sum(item["score"] for item in mastery_rows) / len(mastery_rows)
            if mastery_rows
            else None
        ),
        "gaps": [
            {
                "id": gap.id,
                "concept_id": gap.concept_id,
                "concept": concepts[gap.concept_id].name,
                "mastery_score": gap.mastery_score,
                "threshold": gap.threshold,
                "reason": gap.reason,
            }
            for gap in gaps
            if gap.concept_id in concepts
        ],
        "pathway": pathway_data,
        "pathway_history": pathway_history,
        "notifications": (
            [
                {
                    "type": "Teacher pathway assignment",
                    "title": pathway_data["label"],
                    "message": pathway_data["teacher_note"]
                    or "Your teacher assigned a new evidence-based learning pathway.",
                    "assigned_at": pathway_data["assigned_at"],
                    "due_at": pathway_data["due_at"],
                }
            ]
            if pathway_data and pathway_data["source_type"] == "Teacher"
            else []
        ),
        "progress": {"completed": completed, "total": total},
        "recent_attempts": [
            {
                "id": attempt.id,
                "activity": db.get(Activity, attempt.activity_id).title,
                "score": attempt.score,
                "max_score": attempt.max_score,
                "accuracy": attempt.accuracy,
                "submitted_at": attempt.submitted_at,
                "total_seconds": attempt.total_seconds,
                "attempt_number": attempt.attempt_number,
                "skipped_items": attempt.skipped_items,
                "hint_usage_count": attempt.hint_usage_count,
                "answer_change_count": attempt.answer_change_count,
                "average_response_seconds": (
                    db.scalar(
                        select(func.avg(ItemResponse.response_seconds)).where(
                            ItemResponse.attempt_id == attempt.id
                        )
                    )
                    or 0
                ),
                "mental_effort": ratings[attempt.id].rating if attempt.id in ratings else None,
                "mental_effort_category": (
                    ratings[attempt.id].category if attempt.id in ratings else None
                ),
            }
            for attempt in attempts
        ],
        "recent_activity": [
            {
                "type": "Assessment",
                "title": db.get(Activity, attempt.activity_id).title,
                "occurred_at": attempt.submitted_at,
                "summary": f"{round(attempt.accuracy * 100)}% score",
            }
            for attempt in attempts[:5]
        ],
    }


@app.get("/api/teacher/students/{student_id}/reported-cognitive-load")
def reported_cognitive_load(
    student_id: int,
    db: Session = Depends(get_db),
    _teacher: User = Depends(require_role("teacher")),
):
    student = db.get(User, student_id)
    if not student or student.role != "student":
        raise HTTPException(status_code=404, detail="Student not found")
    ratings = list(
        db.scalars(
            select(MentalEffortRating)
            .where(MentalEffortRating.student_id == student.id)
            .order_by(MentalEffortRating.created_at.desc())
        )
    )
    average = (
        sum(row.rating for row in ratings) / len(ratings) if ratings else None
    )
    low_max = int(get_setting(db, "mental_effort_low_max"))
    moderate_max = int(get_setting(db, "mental_effort_moderate_max"))
    average_category = (
        None
        if average is None
        else "Low"
        if average <= low_max
        else "Moderate"
        if average <= moderate_max
        else "High"
    )
    history = []
    for row in ratings:
        attempt = db.get(AssessmentAttempt, row.attempt_id)
        activity = db.get(Activity, attempt.activity_id) if attempt else None
        history.append(
            {
                "id": row.id,
                "rating": row.rating,
                "category": row.category,
                "activity": activity.title if activity else "Unavailable activity",
                "activity_id": activity.id if activity else None,
                "attempt_id": row.attempt_id,
                "date": row.created_at,
            }
        )
    return {
        "student_id": student.id,
        "participant_code": student.participant_code,
        "display_name": student.display_name,
        "average_rating": average,
        "average_category": average_category,
        "history": history,
        "message": None if history else "No reported cognitive-load data yet.",
        "measure": "Self-reported 1-9 mental-effort rating",
    }


def disable_health_caching(response: Response) -> None:
    response.headers["Cache-Control"] = "no-store, max-age=0"
    response.headers["Pragma"] = "no-cache"


@app.get("/health", include_in_schema=False)
@app.get("/api/health")
def health(response: Response):
    disable_health_caching(response)
    return {
        "status": "ok",
        "service": "NeuroLearn-X API",
        "name": "NeuroLearn-X",
        "version": "1.3.1",
        "mode": "research prototype",
    }


@app.get("/health/live", include_in_schema=False)
@app.get("/api/health/live")
def liveness(response: Response):
    disable_health_caching(response)
    return {"status": "ok"}


@app.get("/health/ready", include_in_schema=False)
@app.get("/api/health/ready")
def readiness(response: Response, db: Session = Depends(get_db)):
    disable_health_caching(response)
    try:
        db.execute(text("SELECT 1"))
    except Exception as error:
        raise HTTPException(
            status_code=503, detail="Database readiness check failed"
        ) from error
    return {"status": "ready", "database": engine.dialect.name}


@app.get("/release/{filename}")
def release_download(filename: str):
    allowed = {
        "NeuroLearn-X.apk",
        "NeuroLearn-X.apk.sha256",
        "NeuroLearn-X-Source-v1.1.0.zip",
        "NeuroLearn-X-Source-v1.1.0.zip.sha256",
        "NeuroLearn-X-Source-v1.2.0.zip",
        "NeuroLearn-X-Source-v1.2.0.zip.sha256",
        "NeuroLearn-X-Source-v1.2.1.zip",
        "NeuroLearn-X-Source-v1.2.1.zip.sha256",
        "NeuroLearn-X-Source-v1.3.1.zip",
        "NeuroLearn-X-Source-v1.3.1.zip.sha256",
        "NeuroLearn-X-Full-System.zip",
        "NeuroLearn-X-Full-System.zip.sha256",
        "NeuroLearn-X-Source-Code.zip",
        "NeuroLearn-X-Source-Code.zip.sha256",
        "NeuroLearn-X-Shareable.zip",
    }
    if filename not in allowed:
        raise HTTPException(status_code=404, detail="Release file not found")
    release_file = Path(__file__).resolve().parents[2] / "release" / filename
    if not release_file.is_file():
        raise HTTPException(status_code=404, detail="Release file not available")
    return FileResponse(
        release_file,
        filename=filename,
        media_type=(
            "application/vnd.android.package-archive"
            if filename.endswith(".apk")
            else "application/zip"
            if filename.endswith(".zip")
            else "text/plain"
        ),
    )


@app.post("/api/auth/register/student", status_code=201)
def register_student(
    payload: StudentRegistrationInput,
    request: Request,
    db: Session = Depends(get_db),
):
    enforce_rate_limit(request, "student-registration", 10, 300)
    identity_checks = [User.participant_code == payload.student_id]
    if payload.username:
        identity_checks.append(func.lower(User.username) == payload.username.lower())
    if payload.email:
        identity_checks.append(func.lower(User.email) == payload.email.lower())
    duplicate = db.scalar(select(User).where(or_(*identity_checks)))
    if duplicate:
        if duplicate.participant_code == payload.student_id:
            detail = "Student ID is already registered"
        elif payload.username and duplicate.username == payload.username:
            detail = "Username is already registered"
        else:
            detail = "Email address is already registered"
        raise HTTPException(status_code=409, detail=detail)

    first_name = payload.first_name.strip()
    last_name = payload.last_name.strip()
    user = User(
        participant_code=payload.student_id,
        username=payload.username,
        email=payload.email,
        password_hash=hash_password(payload.password),
        role="student",
        display_name=f"{first_name} {last_name}".strip(),
        first_name=first_name,
        last_name=last_name,
        must_change_password=False,
        is_active=True,
        is_demo=False,
        account_status="Active",
    )
    db.add(user)
    try:
        db.flush()
    except IntegrityError:
        db.rollback()
        raise HTTPException(
            status_code=409,
            detail="Student ID, username, or email is already registered",
        )
    db.add(
        StudentProfile(
            user_id=user.id,
            grade_level=payload.grade_level.strip(),
            strand="STEM",
            section=payload.section.strip() if payload.section else None,
        )
    )
    db.add(
        ConsentRecord(
            student_id=user.id,
            consented=True,
            consent_version="2.0",
            recorded_by="student-self-registration",
        )
    )
    audit(
        db,
        user.id,
        "account.created",
        "user",
        user.id,
        {"method": "student-self-registration"},
    )
    db.commit()
    db.refresh(user)
    return {
        "message": "Student account created successfully.",
        "student": user_payload(user),
    }


@app.post("/api/auth/forgot-password")
def forgot_password(
    payload: ForgotPasswordInput,
    request: Request,
    db: Session = Depends(get_db),
):
    enforce_rate_limit(request, "forgot-password", 8, 300)
    identifier = payload.identifier.strip()
    user = db.scalar(
        select(User).where(
            or_(
                User.participant_code == identifier.upper(),
                func.lower(User.username) == identifier.lower(),
                func.lower(User.email) == identifier.lower(),
            )
        )
    )
    if user:
        audit(
            db,
            user.id,
            "password.reset.requested",
            "user",
            user.id,
            {"delivery": "authorized-teacher"},
        )
        db.commit()
    return {
        "message": (
            "If that account exists, a request has been recorded. "
            "Please contact an authorized NeuroLearn-X teacher to receive a temporary password."
        )
    }


@app.post("/api/auth/login")
def login(
    payload: LoginInput,
    response: Response,
    request: Request,
    db: Session = Depends(get_db),
):
    identifier = payload.participant_code.strip()
    identity_digest = hashlib.sha256(identifier.casefold().encode()).hexdigest()[:16]
    enforce_rate_limit(request, "login-ip", 30, 300)
    enforce_rate_limit(
        request,
        "login-identity",
        10,
        300,
        identity=identity_digest,
    )
    user = db.scalar(
        select(User).where(
            or_(
                User.participant_code == identifier.upper(),
                func.lower(User.username) == identifier.lower(),
                func.lower(User.email) == identifier.lower(),
            )
        )
    )
    successful = bool(user and verify_password(payload.password, user.password_hash))
    db.add(
        LoginHistory(
            user_id=user.id if user else None,
            participant_code=identifier[:80],
            successful=successful,
            ip_address=request.client.host[:80] if request.client else None,
            user_agent=(request.headers.get("user-agent") or "")[:300] or None,
        )
    )
    if not successful:
        db.commit()
        raise HTTPException(status_code=401, detail="Incorrect participant code or password")
    if user.account_status != "Active" or not user.is_active:
        db.commit()
        raise HTTPException(
            status_code=403,
            detail=f"This account is {user.account_status.lower()}. Contact an authorized teacher.",
        )
    if payload.expected_role and user.role != payload.expected_role:
        db.commit()
        raise HTTPException(status_code=403, detail=f"This account is not a {payload.expected_role}")
    user.last_sign_in_at = datetime.now(timezone.utc)
    db.commit()
    db.refresh(user)
    token = create_access_token(user)
    response.set_cookie(
        COOKIE_NAME,
        token,
        httponly=True,
        secure=COOKIE_SECURE,
        samesite=COOKIE_SAMESITE,
        max_age=43_200,
        path="/",
    )
    return user_payload(user)


@app.post("/api/auth/logout")
def logout(response: Response):
    response.delete_cookie(
        COOKIE_NAME,
        path="/",
        secure=COOKIE_SECURE,
        samesite=COOKIE_SAMESITE,
    )
    return {"ok": True}


@app.get("/api/auth/me")
def me(user: User = Depends(current_user)):
    return user_payload(user)


@app.post("/api/auth/change-password")
def change_password(
    payload: PasswordInput,
    db: Session = Depends(get_db),
    user: User = Depends(current_user),
):
    if not verify_password(payload.current_password, user.password_hash):
        raise HTTPException(status_code=400, detail="Current password is incorrect")
    if payload.current_password == payload.new_password:
        raise HTTPException(status_code=400, detail="Choose a different password")
    user.password_hash = hash_password(payload.new_password)
    user.must_change_password = False
    audit(db, user.id, "password.changed", "user", user.id)
    db.commit()
    return {"ok": True}


@app.get("/api/student/dashboard")
def student_dashboard(
    db: Session = Depends(get_db), student: User = Depends(require_role("student"))
):
    return student_snapshot(db, student)


@app.get("/api/student/onboarding")
def student_onboarding(
    db: Session = Depends(get_db), student: User = Depends(require_role("student"))
):
    profile = db.scalar(
        select(StudentProfile).where(StudentProfile.user_id == student.id)
    )
    if not profile:
        raise HTTPException(status_code=404, detail="Student profile not found")
    return onboarding_payload(profile)


@app.post("/api/student/onboarding/complete")
def complete_student_onboarding(
    db: Session = Depends(get_db), student: User = Depends(require_role("student"))
):
    profile = db.scalar(
        select(StudentProfile).where(StudentProfile.user_id == student.id)
    )
    if not profile:
        raise HTTPException(status_code=404, detail="Student profile not found")
    if profile.onboarding_completed_at is None:
        profile.onboarding_completed_at = datetime.now(timezone.utc)
        profile.onboarding_version = "1.0"
        audit(
            db,
            student.id,
            "student.onboarding.completed",
            "student_profile",
            profile.id,
            {"version": profile.onboarding_version},
        )
        db.commit()
        db.refresh(profile)
    return onboarding_payload(profile)


@app.get("/api/student/onboarding-diagnostic")
def student_onboarding_diagnostic(
    db: Session = Depends(get_db), student: User = Depends(require_role("student"))
):
    status_payload = onboarding_diagnostic_payload(db, student)
    if not status_payload["available"]:
        raise HTTPException(
            status_code=503,
            detail="The 30-item diagnostic assessment is not available",
        )
    return status_payload


@app.get("/api/student/targets")
def target_competencies(
    db: Session = Depends(get_db), _student: User = Depends(require_role("student"))
):
    concepts = list(
        db.scalars(
            select(Concept)
            .where(Concept.subject == "General Physics", Concept.active.is_(True))
            .order_by(Concept.difficulty, Concept.name)
        )
    )
    return [concept_payload(concept) for concept in concepts]


@app.post("/api/student/target")
def choose_target(
    payload: TargetInput,
    db: Session = Depends(get_db),
    student: User = Depends(require_role("student")),
):
    concept = db.get(Concept, payload.concept_id)
    if not concept or not concept.active or concept.subject != "General Physics":
        raise HTTPException(status_code=404, detail="Physics competency not found")
    profile = db.scalar(select(StudentProfile).where(StudentProfile.user_id == student.id))
    profile.target_concept_id = concept.id
    db.commit()
    candidates = generate_pathways(
        db, student, trigger_type="Target competency changed", trigger_id=concept.id
    )
    diagnostic = db.scalar(
        select(Activity)
        .join(ActivityConcept, ActivityConcept.activity_id == Activity.id)
        .where(
            ActivityConcept.concept_id == concept.id,
            Activity.is_diagnostic.is_(True),
            Activity.is_onboarding_diagnostic.is_(False),
            Activity.active.is_(True),
        )
    )
    return {
        "target": concept_payload(concept),
        "diagnostic_activity_id": diagnostic.id if diagnostic else None,
        "pathway_generated": bool(candidates),
    }


@app.get("/api/student/activities/{activity_id}")
def student_activity(
    activity_id: int,
    db: Session = Depends(get_db),
    student: User = Depends(require_role("student")),
):
    assessment = db.scalar(
        select(Assessment).where(Assessment.activity_id == activity_id)
    )
    if assessment:
        refresh_assessment_status(db, assessment)
        db.commit()
        if (
            assessment.status != "Published"
            or not assessment_assigned_to_student(db, assessment, student)
        ):
            raise HTTPException(status_code=404, detail="Assessment not available")
    activity = db.get(Activity, activity_id)
    if not activity or not activity.active:
        raise HTTPException(status_code=404, detail="Activity not found")
    payload = activity_payload(db, activity, include_questions=True)
    payload["adaptive_tutoring"] = bool(
        not assessment and not activity.is_diagnostic
    )
    payload["default_tutoring_mode"] = (
        "mastery_check" if activity.activity_type == "quiz" else "guided"
    )
    payload["time_limit_seconds"] = (
        assessment.time_limit * 60
        if assessment and assessment.time_limit
        else activity.estimated_minutes * 60
        if activity.is_onboarding_diagnostic
        else 5 * 60
    )
    if assessment:
        seed = int.from_bytes(
            hashlib.sha256(
                f"{assessment.id}:{student.id}".encode("utf-8")
            ).digest()[:8],
            "big",
        )
        generator = random.Random(seed)
        if assessment.shuffle_questions:
            generator.shuffle(payload["questions"])
        if assessment.shuffle_choices:
            for question in payload["questions"]:
                generator.shuffle(question["choices"])
    return payload


@app.post("/api/student/tutoring-sessions", status_code=201)
def create_tutoring_session(
    payload: TutoringSessionInput,
    db: Session = Depends(get_db),
    student: User = Depends(require_role("student")),
):
    activity = db.get(Activity, payload.activity_id)
    if not activity or not activity.active or activity.is_diagnostic:
        raise HTTPException(status_code=404, detail="Adaptive activity not found")
    if db.scalar(select(Assessment.id).where(Assessment.activity_id == activity.id)):
        raise HTTPException(
            status_code=409,
            detail="Published assessments use the fixed assessment workflow",
        )
    return start_tutoring_session(db, student, activity, payload.mode)


@app.post("/api/student/tutoring-sessions/{session_id}/responses")
def answer_tutoring_question(
    session_id: int,
    payload: TutoringResponseInput,
    db: Session = Depends(get_db),
    student: User = Depends(require_role("student")),
):
    session = db.get(TutoringSession, session_id)
    if not session or session.student_id != student.id:
        raise HTTPException(status_code=404, detail="Tutoring session not found")
    return submit_tutoring_response(
        db,
        student,
        session,
        payload.question_id,
        payload.selected_choice_id,
        payload.response_text,
        payload.response_seconds,
        payload.hint_opened,
        payload.answer_changes,
    )


@app.get("/api/student/learning-summaries")
def student_learning_summaries(
    db: Session = Depends(get_db),
    student: User = Depends(require_role("student")),
):
    rows = list(
        db.scalars(
            select(LearningSummary)
            .where(LearningSummary.student_id == student.id)
            .order_by(LearningSummary.created_at.desc())
            .limit(20)
        )
    )
    return [
        {
            "id": row.id,
            "activity_id": row.activity_id,
            "activity": db.get(Activity, row.activity_id).title,
            "attempt_id": row.attempt_id,
            "tutoring_session_id": row.tutoring_session_id,
            "summary": row.summary,
            "created_at": row.created_at,
        }
        for row in rows
    ]


@app.get("/api/student/assessments")
def student_assessments(
    db: Session = Depends(get_db),
    student: User = Depends(require_role("student")),
):
    assessments = list(
        db.scalars(
            select(Assessment)
            .where(Assessment.status.in_(["Scheduled", "Published", "Closed"]))
            .order_by(Assessment.available_from, Assessment.created_at.desc())
        )
    )
    output = []
    changed = False
    for assessment in assessments:
        previous = assessment.status
        refresh_assessment_status(db, assessment)
        changed = changed or previous != assessment.status
        if not assessment_assigned_to_student(db, assessment, student):
            continue
        attempt_count = db.scalar(
            select(func.count(AssessmentAttempt.id)).where(
                AssessmentAttempt.student_id == student.id,
                AssessmentAttempt.activity_id == assessment.activity_id,
            )
        ) or 0
        output.append(
            {
                "id": assessment.id,
                "activity_id": assessment.activity_id,
                "title": assessment.title,
                "description": assessment.description,
                "subject": assessment.subject,
                "topic": assessment.topic,
                "status": assessment.status,
                "mastery_threshold": assessment.mastery_threshold,
                "time_limit": assessment.time_limit,
                "maximum_attempts": assessment.maximum_attempts,
                "attempt_count": attempt_count,
                "available_from": assessment.available_from,
                "due_at": assessment.due_at,
                "can_attempt": (
                    assessment.status == "Published"
                    and attempt_count < assessment.maximum_attempts
                ),
            }
        )
    if changed:
        db.commit()
    return output


@app.get("/api/student/diagnostic/{concept_id}")
def student_diagnostic(
    concept_id: int,
    db: Session = Depends(get_db),
    _student: User = Depends(require_role("student")),
):
    activity = db.scalar(
        select(Activity)
        .join(ActivityConcept, ActivityConcept.activity_id == Activity.id)
        .where(
            ActivityConcept.concept_id == concept_id,
            Activity.is_diagnostic.is_(True),
            Activity.is_onboarding_diagnostic.is_(False),
            Activity.active.is_(True),
        )
    )
    if not activity:
        raise HTTPException(status_code=404, detail="Diagnostic not available")
    return activity_payload(db, activity, include_questions=True)


@app.post("/api/student/attempts", status_code=201)
def submit_attempt(
    payload: AttemptInput,
    db: Session = Depends(get_db),
    student: User = Depends(require_role("student")),
):
    # Lock this learner's row for the duration of the submission transaction.
    # PostgreSQL then serializes maximum-attempt checks and attempt numbering;
    # SQLite safely treats FOR UPDATE as a no-op for local development.
    db.execute(select(User.id).where(User.id == student.id).with_for_update())
    activity = db.get(Activity, payload.activity_id)
    if not activity or not activity.active:
        raise HTTPException(status_code=404, detail="Activity not found")
    published_assessment = db.scalar(
        select(Assessment).where(Assessment.activity_id == activity.id)
    )
    if published_assessment:
        refresh_assessment_status(db, published_assessment)
        if (
            published_assessment.status != "Published"
            or not assessment_assigned_to_student(db, published_assessment, student)
        ):
            raise HTTPException(status_code=403, detail="Assessment is not available")
        submitted_count = db.scalar(
            select(func.count(AssessmentAttempt.id)).where(
                AssessmentAttempt.student_id == student.id,
                AssessmentAttempt.activity_id == activity.id,
            )
        ) or 0
        if submitted_count >= published_assessment.maximum_attempts:
            raise HTTPException(
                status_code=409,
                detail="Maximum assessment attempts reached",
            )
    existing = db.scalar(
        select(AssessmentAttempt).where(
            AssessmentAttempt.student_id == student.id,
            AssessmentAttempt.activity_id == activity.id,
            AssessmentAttempt.started_at == payload.started_at,
        )
    )
    if existing:
        raise HTTPException(status_code=409, detail="This attempt was already submitted")
    questions = {
        question.id: question
        for question in db.scalars(
            select(Question).where(
                Question.activity_id == activity.id, Question.active.is_(True)
            )
        )
    }
    if not questions:
        raise HTTPException(status_code=400, detail="This activity has no active questions")
    mastery_before = {
        concept_id: record.mastery_score
        for concept_id, record in latest_mastery_map(db, student.id).items()
        if concept_id in {question.concept_id for question in questions.values()}
    }
    response_ids = [response.question_id for response in payload.responses]
    if len(response_ids) != len(set(response_ids)):
        raise HTTPException(status_code=400, detail="Each question may be answered only once")
    unknown = set(response_ids) - set(questions)
    if unknown:
        raise HTTPException(status_code=400, detail="Response contains an invalid question")
    selected_by_question: dict[int, AnswerChoice | None] = {}
    for item in payload.responses:
        choice = db.get(AnswerChoice, item.selected_choice_id) if item.selected_choice_id else None
        if choice and choice.question_id != item.question_id:
            raise HTTPException(status_code=400, detail="Answer choice does not match question")
        selected_by_question[item.question_id] = choice
    submitted_at = datetime.now(timezone.utc)
    started_at = payload.started_at
    if started_at.tzinfo is None:
        started_at = started_at.replace(tzinfo=timezone.utc)
    if started_at > submitted_at:
        raise HTTPException(
            status_code=400,
            detail="Attempt start time cannot be in the future",
        )
    total_seconds = max(0, min(86_400, (submitted_at - started_at).total_seconds()))
    if (
        published_assessment
        and published_assessment.time_limit
        and total_seconds > published_assessment.time_limit * 60 + 60
    ):
        raise HTTPException(
            status_code=408,
            detail="The assessment time limit has expired",
        )
    previous_attempts = db.scalar(
        select(func.count(AssessmentAttempt.id)).where(
            AssessmentAttempt.student_id == student.id,
            AssessmentAttempt.activity_id == activity.id,
        )
    ) or 0
    max_score = sum(question.points for question in questions.values())
    response_lookup = {response.question_id: response for response in payload.responses}

    def text_answer_correct(question: Question, response_text: str | None) -> bool:
        if not response_text or not question.correct_answer:
            return False
        supplied = re.sub(r"\W+", " ", response_text.casefold()).strip()
        expected = re.sub(r"\W+", " ", question.correct_answer.casefold()).strip()
        if question.question_type == "Identification":
            return supplied == expected
        supplied_tokens = set(supplied.split())
        expected_tokens = {
            token
            for token in expected.split()
            if len(token) > 2 and token not in {"the", "and", "that", "with", "from"}
        }
        if not expected_tokens:
            return supplied == expected
        return len(supplied_tokens & expected_tokens) / len(expected_tokens) >= 0.6

    score = 0.0
    for question_id, question in questions.items():
        choice = selected_by_question.get(question_id)
        item = response_lookup.get(question_id)
        if (choice and choice.is_correct) or (
            item and text_answer_correct(question, item.response_text)
        ):
            score += question.points
    skipped = sum(
        1
        for question_id in questions
        if question_id not in response_lookup
        or response_lookup[question_id].skipped
        or (
            response_lookup[question_id].selected_choice_id is None
            and not (response_lookup[question_id].response_text or "").strip()
        )
    )
    hint_count = sum(1 for item in payload.responses if item.hint_opened)
    changes = sum(item.answer_changes for item in payload.responses)
    attempt = AssessmentAttempt(
        student_id=student.id,
        activity_id=activity.id,
        score=score,
        max_score=max_score,
        accuracy=score / max_score if max_score else 0,
        started_at=started_at,
        submitted_at=submitted_at,
        total_seconds=total_seconds,
        skipped_items=skipped,
        hint_usage_count=hint_count,
        answer_change_count=changes,
        attempt_number=previous_attempts + 1,
        is_demo=student.is_demo,
    )
    db.add(attempt)
    try:
        db.flush()
    except IntegrityError as error:
        db.rollback()
        raise HTTPException(
            status_code=409,
            detail="This attempt was already submitted",
        ) from error
    concept_stats: dict[int, dict[str, float]] = {}
    result_items = []
    for question_id, question in questions.items():
        item = response_lookup.get(question_id)
        choice = selected_by_question.get(question_id)
        correct_choice = db.scalar(
            select(AnswerChoice).where(
                AnswerChoice.question_id == question_id,
                AnswerChoice.is_correct.is_(True),
            )
        )
        is_correct = bool(
            (choice and choice.is_correct)
            or (item and text_answer_correct(question, item.response_text))
        )
        earned = question.points if is_correct else 0
        has_response = bool(choice or (item and (item.response_text or "").strip()))
        response_row = ItemResponse(
            attempt_id=attempt.id,
            question_id=question_id,
            selected_choice_id=choice.id if choice else None,
            response_text=clean_content(item.response_text, 5000)
            if item and item.response_text
            else None,
            is_correct=is_correct,
            earned_points=earned,
            max_points=question.points,
            response_seconds=item.response_seconds if item else 0,
            hint_opened=item.hint_opened if item else False,
            skipped=(item.skipped or not has_response) if item else True,
            answer_changes=item.answer_changes if item else 0,
        )
        db.add(response_row)
        mapped_misconception = None
        pattern_confidence = "Not diagnosed"
        pattern_count = 0
        if (
            not is_correct
            and choice
            and choice.misconception_id
            and choice.mapping_status in {"Teacher reviewed", "Validated"}
        ):
            mapped_misconception = db.get(Misconception, choice.misconception_id)
            if mapped_misconception and mapped_misconception.active:
                pattern_count = 1 + (
                    db.scalar(
                        select(func.count(MisconceptionHistory.id)).where(
                            MisconceptionHistory.student_id == student.id,
                            MisconceptionHistory.misconception_id == mapped_misconception.id,
                            MisconceptionHistory.resolved_at.is_(None),
                        )
                    )
                    or 0
                )
                pattern_confidence = (
                    "High" if pattern_count >= 3 else "Moderate" if pattern_count == 2 else "Low"
                )
                db.add(
                    MisconceptionHistory(
                        student_id=student.id,
                        misconception_id=mapped_misconception.id,
                        question_id=question.id,
                        selected_choice_id=choice.id,
                        attempt_id=attempt.id,
                        evidence_count=pattern_count,
                        confidence_level=pattern_confidence,
                    )
                )
        stats = concept_stats.setdefault(
            question.concept_id,
            {"earned": 0, "maximum": 0, "seconds": 0, "items": 0, "skipped": 0, "hints": 0},
        )
        stats["earned"] += earned
        stats["maximum"] += question.points
        stats["seconds"] += item.response_seconds if item else 0
        stats["items"] += 1
        stats["skipped"] += int((item.skipped or not has_response) if item else True)
        stats["hints"] += int(item.hint_opened if item else False)
        result_items.append(
            {
                "question_id": question_id,
                "correct": is_correct,
                "learner_answer": (
                    choice.text if choice else item.response_text if item else None
                ),
                "selected_choice": (
                    choice.text if choice else item.response_text if item else None
                ),
                "correct_choice": (
                    correct_choice.text
                    if correct_choice
                    else question.correct_answer or None
                ),
                "why": question.explanation or question.feedback,
                "feedback": (
                    "Correct — this demonstrates the expected concept evidence."
                    if is_correct
                    else question.explanation or question.feedback
                ),
                "solution_steps": question.solution_steps or question.explanation or question.feedback,
                "solution": structured_solution(
                    question,
                    correct_choice.text if correct_choice else question.correct_answer,
                ),
                "misconception": (
                    {
                        "code": mapped_misconception.code,
                        "name": mapped_misconception.name,
                        "concept": db.get(Concept, mapped_misconception.concept_id).name,
                        "explanation": mapped_misconception.explanation,
                        "remediation_instruction": mapped_misconception.remediation_instruction,
                        "suggested_activity_id": mapped_misconception.suggested_activity_id,
                        "pattern_confidence": pattern_confidence,
                        "evidence_count": pattern_count,
                    }
                    if mapped_misconception
                    else None
                ),
                "diagnostic_note": (
                    None
                    if is_correct or mapped_misconception
                    else "This response is incorrect, but no teacher-reviewed distractor mapping supports a specific misconception diagnosis."
                ),
                "likely_mistake": (
                    None
                    if is_correct
                    else mapped_misconception.explanation
                    if mapped_misconception
                    else "No validated misconception pattern was diagnosed from this response."
                ),
                "review_concept": db.get(Concept, question.concept_id).name,
                "choice_explanations": question.distractor_rationales or {},
                "earned_points": earned,
                "max_points": question.points,
            }
        )
    for concept_id, stats in concept_stats.items():
        db.add(
            InteractionLog(
                student_id=student.id,
                activity_id=activity.id,
                concept_id=concept_id,
                attempt_id=attempt.id,
                score=stats["earned"],
                max_score=stats["maximum"],
                response_accuracy=(
                    stats["earned"] / stats["maximum"] if stats["maximum"] else 0
                ),
                average_response_seconds=stats["seconds"] / max(1, stats["items"]),
                total_completion_seconds=total_seconds,
                number_of_attempts=previous_attempts + 1,
                skipped_items=int(stats["skipped"]),
                hint_usage_count=int(stats["hints"]),
                start_time=started_at,
                submission_time=submitted_at,
                is_demo=student.is_demo,
            )
        )
    db.commit()
    updated_mastery = recalculate_mastery(
        db, student, set(concept_stats), attempt.id
    )
    record_pathway_evidence(db, student.id, activity.id, attempt)
    if activity.is_onboarding_diagnostic:
        profile = db.scalar(
            select(StudentProfile).where(StudentProfile.user_id == student.id)
        )
        if profile and profile.target_concept_id is None and updated_mastery:
            physics_records = [
                record
                for record in updated_mastery
                if db.get(Concept, record.concept_id).subject == "General Physics"
            ]
            weakest = min(
                physics_records or updated_mastery,
                key=lambda record: (record.mastery_score, record.concept_id),
            )
            profile.target_concept_id = weakest.concept_id
    db.commit()
    previous_selected = db.scalar(
        select(PathwayRecommendation)
        .where(
            PathwayRecommendation.student_id == student.id,
            PathwayRecommendation.active.is_(True),
            PathwayRecommendation.selected.is_(True),
        )
        .order_by(PathwayRecommendation.created_at.desc())
    )
    generate_pathways(
        db, student, trigger_type="Assessment submission", trigger_id=attempt.id
    )
    current_selected = db.scalar(
        select(PathwayRecommendation)
        .where(
            PathwayRecommendation.student_id == student.id,
            PathwayRecommendation.active.is_(True),
            PathwayRecommendation.selected.is_(True),
        )
        .order_by(PathwayRecommendation.created_at.desc())
    )
    summary_data = {
        "initial_mastery": mastery_before,
        "final_mastery": {
            record.concept_id: record.mastery_score for record in updated_mastery
        },
        "accuracy": attempt.accuracy,
        "questions_completed": len(questions),
        "concepts_strengthened": [
            db.get(Concept, record.concept_id).name
            for record in updated_mastery
            if record.mastery_score > mastery_before.get(record.concept_id, 0)
        ],
        "errors_observed": sum(1 for item in result_items if not item["correct"]),
        "misconceptions": [
            item["misconception"]
            for item in result_items
            if item.get("misconception")
        ],
        "mental_effort": "Pending learner rating",
        "analysis_type": (
            "NeuroLearn-X onboarding diagnostic analysis"
            if activity.is_onboarding_diagnostic
            else "NeuroLearn-X adaptive evidence analysis"
        ),
        "learning_gaps_identified": [
            db.get(Concept, record.concept_id).name
            for record in sorted(updated_mastery, key=lambda row: row.mastery_score)
            if record.classification != "Mastered"
        ],
        "pathway_changed": bool(
            (previous_selected.id if previous_selected else None)
            != (current_selected.id if current_selected else None)
        ),
        "pathway_before": previous_selected.label if previous_selected else None,
        "pathway_after": current_selected.label if current_selected else None,
        "next_action": "Rate your mental effort, then review the updated pathway.",
    }
    db.add(
        LearningSummary(
            student_id=student.id,
            activity_id=activity.id,
            attempt_id=attempt.id,
            summary=summary_data,
        )
    )
    db.commit()
    score_visible = (
        published_assessment is None
        or published_assessment.show_score_immediately
    )
    explanations_visible = (
        published_assessment is None
        or published_assessment.show_explanations
    )
    return {
        "attempt_id": attempt.id,
        "score": score if score_visible else None,
        "max_score": max_score if score_visible else None,
        "accuracy": (score / max_score if max_score else 0)
        if score_visible
        else None,
        "items": result_items if explanations_visible else [],
        "score_visible": score_visible,
        "explanations_visible": explanations_visible,
        "mental_effort_required": True,
        "summary": summary_data,
        "mental_effort_boundaries": {
            "low_max": int(get_setting(db, "mental_effort_low_max")),
            "moderate_max": int(get_setting(db, "mental_effort_moderate_max")),
        },
    }


@app.post("/api/student/attempts/{attempt_id}/mental-effort")
def save_mental_effort(
    attempt_id: int,
    payload: MentalEffortInput,
    db: Session = Depends(get_db),
    student: User = Depends(require_role("student")),
):
    attempt = db.get(AssessmentAttempt, attempt_id)
    if not attempt or attempt.student_id != student.id:
        raise HTTPException(status_code=404, detail="Attempt not found")
    existing = db.scalar(
        select(MentalEffortRating).where(MentalEffortRating.attempt_id == attempt_id)
    )
    if existing:
        raise HTTPException(status_code=409, detail="Mental-effort rating already submitted")
    category = mental_effort_category(db, payload.rating)
    rating_record = MentalEffortRating(
            student_id=student.id,
            attempt_id=attempt.id,
            rating=payload.rating,
            category=category,
            is_demo=student.is_demo,
        )
    db.add(rating_record)
    db.flush()
    update_summary_effort(db, attempt.id, rating_record)
    db.commit()
    generate_pathways(
        db, student, trigger_type="Mental-effort rating", trigger_id=rating_record.id
    )
    try:
        cognitive_load_analysis = predict_student_cognitive_load(db, student.id)
    except ValueError:
        cognitive_load_analysis = {"available": False}
    profile = db.scalar(
        select(StudentProfile).where(StudentProfile.user_id == student.id)
    )
    pathway = db.scalar(
        select(PathwayRecommendation)
        .where(
            PathwayRecommendation.student_id == student.id,
            PathwayRecommendation.active.is_(True),
            PathwayRecommendation.selected.is_(True),
        )
        .order_by(PathwayRecommendation.created_at.desc())
    )
    return {
        "rating": payload.rating,
        "category": category,
        "analysis_complete": True,
        "target": (
            concept_payload(db.get(Concept, profile.target_concept_id))
            if profile and profile.target_concept_id
            else None
        ),
        "predicted_pathway_load": pathway.cognitive_load_category if pathway else None,
        "pathway_updated": pathway is not None,
        "model_prediction_available": bool(cognitive_load_analysis.get("available")),
    }


@app.get("/api/student/pathways")
def student_pathways(
    db: Session = Depends(get_db), student: User = Depends(require_role("student"))
):
    pathways = list(
        db.scalars(
            select(PathwayRecommendation)
            .where(
                PathwayRecommendation.student_id == student.id,
                PathwayRecommendation.active.is_(True),
            )
            .order_by(PathwayRecommendation.adaptive_pathway_score.desc())
        )
    )
    return [serialize_pathway(db, pathway) for pathway in pathways]


@app.post("/api/student/pathway-steps/{step_id}/complete")
def complete_pathway_step(
    step_id: int,
    db: Session = Depends(get_db),
    student: User = Depends(require_role("student")),
):
    step = db.get(PathwayStep, step_id)
    pathway = db.get(PathwayRecommendation, step.pathway_id) if step else None
    if not step or not pathway or pathway.student_id != student.id or not pathway.active:
        raise HTTPException(status_code=404, detail="Pathway step not found")
    if not step.completed_at:
        attempt = db.scalar(
            select(AssessmentAttempt)
            .where(
                AssessmentAttempt.student_id == student.id,
                AssessmentAttempt.activity_id == step.activity_id,
                AssessmentAttempt.accuracy >= 0.60,
                AssessmentAttempt.submitted_at >= pathway.created_at,
            )
            .order_by(AssessmentAttempt.submitted_at.desc())
        )
        if not attempt:
            raise HTTPException(
                status_code=409,
                detail="Complete the mastery check with at least 60% accuracy before this step can be completed.",
            )
        step.completed_at = attempt.submitted_at
        step.completion_attempt_id = attempt.id
        db.commit()
    return {"ok": True, "completed_at": step.completed_at}


@app.get("/api/student/graph")
def student_graph(
    db: Session = Depends(get_db), student: User = Depends(require_role("student"))
):
    profile = db.scalar(select(StudentProfile).where(StudentProfile.user_id == student.id))
    if not profile or not profile.target_concept_id:
        return {"nodes": [], "edges": []}
    edge_rows = list(
        db.scalars(select(PrerequisiteEdge).where(PrerequisiteEdge.active.is_(True)))
    )
    edge_tuples = [
        (edge.prerequisite_concept_id, edge.succeeding_concept_id) for edge in edge_rows
    ]
    node_ids = prerequisite_ancestors(edge_tuples, profile.target_concept_id) | {
        profile.target_concept_id
    }
    mastery = latest_mastery_map(db, student.id)
    threshold = float(get_setting(db, "mastery_threshold"))
    nodes = []
    for concept_id in topological_order(node_ids, edge_tuples):
        concept = db.get(Concept, concept_id)
        record = mastery.get(concept_id)
        nodes.append(
            {
                **concept_payload(concept),
                "mastery_score": record.mastery_score if record else None,
                "classification": record.classification if record else "Not Yet Assessed",
                "is_target": concept_id == profile.target_concept_id,
                "is_gap": bool(record and record.mastery_score < threshold),
            }
        )
    return {
        "nodes": nodes,
        "edges": [
            {
                "id": edge.id,
                "source": edge.prerequisite_concept_id,
                "target": edge.succeeding_concept_id,
            }
            for edge in edge_rows
            if edge.prerequisite_concept_id in node_ids
            and edge.succeeding_concept_id in node_ids
        ],
    }


@app.get("/api/student/history")
def student_history(
    db: Session = Depends(get_db), student: User = Depends(require_role("student"))
):
    return student_snapshot(db, student)["recent_attempts"]


@app.get("/api/teacher/dashboard")
def teacher_dashboard(
    db: Session = Depends(get_db), _teacher: User = Depends(require_role("teacher"))
):
    students = list(
        db.scalars(
            select(User)
            .where(
                User.role == "student",
                User.is_active.is_(True),
                User.account_status == "Active",
            )
            .order_by(User.updated_at.desc())
        )
    )
    active_student_ids = [student.id for student in students]
    attempts = (
        db.scalar(
            select(func.count(AssessmentAttempt.id)).where(
                AssessmentAttempt.student_id.in_(active_student_ids)
            )
        )
        if active_student_ids
        else 0
    ) or 0
    active_pathways = db.scalar(
        select(func.count(PathwayRecommendation.id)).where(
            PathwayRecommendation.active.is_(True),
            PathwayRecommendation.selected.is_(True),
            PathwayRecommendation.student_id.in_(active_student_ids),
        )
    ) or 0
    loads = {
        label: db.scalar(
            select(func.count(MentalEffortRating.id)).where(
                MentalEffortRating.category == label,
                MentalEffortRating.student_id.in_(active_student_ids),
            )
        )
        or 0
        for label in ("Low", "Moderate", "High")
    }
    mastery_records = list(
        db.scalars(
            select(MasteryRecord)
            .where(MasteryRecord.student_id.in_(active_student_ids))
            .order_by(MasteryRecord.created_at.desc())
        )
    )
    latest_by_pair = {}
    for record in mastery_records:
        latest_by_pair.setdefault((record.student_id, record.concept_id), record)
    mastered = sum(1 for record in latest_by_pair.values() if record.classification == "Mastered")
    misconception_patterns = db.execute(
        select(
            Misconception.id,
            Misconception.code,
            Misconception.name,
            func.count(MisconceptionHistory.id),
        )
        .join(
            MisconceptionHistory,
            MisconceptionHistory.misconception_id == Misconception.id,
        )
        .where(MisconceptionHistory.resolved_at.is_(None))
        .where(MisconceptionHistory.student_id.in_(active_student_ids))
        .group_by(Misconception.id)
        .order_by(func.count(MisconceptionHistory.id).desc())
        .limit(8)
    ).all()
    open_interventions = db.scalar(
        select(func.count(TeacherIntervention.id)).where(
            TeacherIntervention.status == "Open",
            TeacherIntervention.student_id.in_(active_student_ids),
        )
    ) or 0
    at_risk = len(
        {
            record.student_id
            for record in latest_by_pair.values()
            if record.mastery_score < 0.5
        }
    )
    return {
        "student_count": len(students),
        "attempt_count": attempts,
        "active_pathways": active_pathways,
        "mastered_concepts": mastered,
        "at_risk_students": at_risk,
        "open_interventions": open_interventions,
        "misconception_patterns": [
            {"id": item_id, "code": code, "name": name, "evidence_count": count}
            for item_id, code, name, count in misconception_patterns
        ],
        "cognitive_load_distribution": loads,
        "demo_notice": "Demonstration Data – Not a Research Result.",
        "recent_students": [
            {
                "id": student.id,
                "participant_code": student.participant_code,
                "display_name": student.display_name,
                "attempts": db.scalar(
                    select(func.count(AssessmentAttempt.id)).where(
                        AssessmentAttempt.student_id == student.id
                    )
                )
                or 0,
            }
            for student in students[:8]
        ],
    }


@app.get("/api/teacher/students")
def teacher_students(
    search: str = "",
    target: str | None = None,
    mastery_level: str | None = None,
    load: str | None = None,
    account_status: str | None = None,
    registered_from: str | None = None,
    registered_to: str | None = None,
    last_sign_in_from: str | None = None,
    last_sign_in_to: str | None = None,
    sort_by: str = "participant_code",
    sort_order: str = "asc",
    page: int = Query(default=1, ge=1),
    page_size: int = Query(default=10, ge=5, le=100),
    paginated: bool = False,
    include_archived: bool = False,
    include_deactivated: bool = False,
    db: Session = Depends(get_db),
    _teacher: User = Depends(require_role("teacher")),
):
    query = select(User).where(User.role == "student")
    if account_status == "Archived":
        query = query.where(User.account_status == "Archived")
    elif account_status == "Deactivated":
        query = query.where(
            User.account_status == "Deactivated",
            User.is_active.is_(False),
        )
    elif account_status == "Active":
        query = query.where(
            User.account_status == "Active",
            User.is_active.is_(True),
        )
    elif include_deactivated:
        query = query.where(User.account_status != "Archived")
    elif not include_archived:
        query = query.where(
            User.account_status == "Active",
            User.is_active.is_(True),
        )
    if search:
        query = query.where(
            or_(
                User.participant_code.ilike(f"%{search}%"),
                User.display_name.ilike(f"%{search}%"),
            )
        )
    students = list(db.scalars(query.order_by(User.participant_code)))
    output = []
    for student in students:
        snapshot = student_snapshot(db, student)
        category = (
            snapshot["pathway"]["cognitive_load_category"]
            if snapshot["pathway"]
            else "Not estimated"
        )
        if load and category != load:
            continue
        output.append(
            {
                "id": student.id,
                "participant_code": student.participant_code,
                "display_name": student.display_name,
                "target": snapshot["target"]["name"] if snapshot["target"] else None,
                "mastery_average": snapshot["mastery_average"],
                "gaps": len(snapshot["gaps"]),
                "cognitive_load": category,
                "progress": snapshot["progress"],
                "is_demo": student.is_demo,
                "created_at": student.created_at,
                "last_sign_in_at": student.last_sign_in_at,
                "account_status": student.account_status,
                "grade_level": (
                    student.student_profile.grade_level
                    if student.student_profile
                    else None
                ),
                "section": (
                    student.student_profile.section
                    if student.student_profile
                    else None
                ),
            }
        )

    def date_value(value: str | None):
        if not value:
            return None
        try:
            return datetime.fromisoformat(value).date()
        except ValueError:
            raise HTTPException(status_code=422, detail="Invalid date filter")

    registered_start = date_value(registered_from)
    registered_end = date_value(registered_to)
    sign_in_start = date_value(last_sign_in_from)
    sign_in_end = date_value(last_sign_in_to)

    def mastery_group(value: float | None) -> str:
        if value is None:
            return "No evidence"
        if value >= 0.75:
            return "Mastered"
        if value >= 0.50:
            return "Developing"
        return "At risk"

    filtered = []
    for item in output:
        registered_date = item["created_at"].date()
        last_sign_in_date = (
            item["last_sign_in_at"].date() if item["last_sign_in_at"] else None
        )
        if target and item["target"] != target:
            continue
        if mastery_level and mastery_group(item["mastery_average"]) != mastery_level:
            continue
        if account_status and item["account_status"] != account_status:
            continue
        if registered_start and registered_date < registered_start:
            continue
        if registered_end and registered_date > registered_end:
            continue
        if sign_in_start and (
            last_sign_in_date is None or last_sign_in_date < sign_in_start
        ):
            continue
        if sign_in_end and (
            last_sign_in_date is None or last_sign_in_date > sign_in_end
        ):
            continue
        item["mastery_level"] = mastery_group(item["mastery_average"])
        filtered.append(item)

    sortable = {
        "participant_code",
        "display_name",
        "target",
        "mastery_average",
        "gaps",
        "cognitive_load",
        "created_at",
        "last_sign_in_at",
        "account_status",
    }
    if sort_by not in sortable or sort_order not in {"asc", "desc"}:
        raise HTTPException(status_code=422, detail="Invalid sort option")
    filtered.sort(
        key=lambda item: (
            item[sort_by] is None,
            item[sort_by].casefold()
            if isinstance(item[sort_by], str)
            else item[sort_by],
        ),
        reverse=sort_order == "desc",
    )
    if not paginated:
        return filtered
    total = len(filtered)
    start = (page - 1) * page_size
    return {
        "items": filtered[start : start + page_size],
        "total": total,
        "page": page,
        "page_size": page_size,
        "total_pages": max(1, (total + page_size - 1) // page_size),
    }


def delete_attempt_records(db: Session, attempt_ids: list[int]) -> None:
    if not attempt_ids:
        return
    session_ids = list(
        db.scalars(
            select(TutoringSession.id).where(
                TutoringSession.attempt_id.in_(attempt_ids)
            )
        )
    )
    if session_ids:
        db.execute(
            delete(TutoringResponse).where(
                TutoringResponse.session_id.in_(session_ids)
            )
        )
        db.execute(
            delete(LearningSummary).where(
                LearningSummary.tutoring_session_id.in_(session_ids)
            )
        )
        db.execute(
            delete(MisconceptionHistory).where(
                MisconceptionHistory.tutoring_session_id.in_(session_ids)
            )
        )
        db.execute(delete(TutoringSession).where(TutoringSession.id.in_(session_ids)))
    db.execute(
        update(PathwayStep)
        .where(PathwayStep.completion_attempt_id.in_(attempt_ids))
        .values(completion_attempt_id=None, completed_at=None)
    )
    db.execute(
        delete(MisconceptionHistory).where(
            or_(
                MisconceptionHistory.attempt_id.in_(attempt_ids),
                MisconceptionHistory.resolved_by_attempt_id.in_(attempt_ids),
            )
        )
    )
    for model in (
        LearningSummary,
        MentalEffortRating,
        InteractionLog,
        ItemResponse,
        MasteryRecord,
    ):
        db.execute(delete(model).where(model.attempt_id.in_(attempt_ids)))
    db.execute(
        delete(AssessmentAttempt).where(AssessmentAttempt.id.in_(attempt_ids))
    )


def permanently_delete_student(db: Session, student: User) -> None:
    pathway_ids = list(
        db.scalars(
            select(PathwayRecommendation.id).where(
                PathwayRecommendation.student_id == student.id
            )
        )
    )
    if pathway_ids:
        db.execute(
            delete(ExpertEvaluation).where(
                ExpertEvaluation.pathway_id.in_(pathway_ids)
            )
        )
        db.execute(
            delete(TeacherIntervention).where(
                or_(
                    TeacherIntervention.student_id == student.id,
                    TeacherIntervention.pathway_id.in_(pathway_ids),
                )
            )
        )
        db.execute(delete(PathwayStep).where(PathwayStep.pathway_id.in_(pathway_ids)))
        db.execute(
            delete(PathwayVersion).where(
                or_(
                    PathwayVersion.student_id == student.id,
                    PathwayVersion.pathway_id.in_(pathway_ids),
                    PathwayVersion.previous_pathway_id.in_(pathway_ids),
                )
            )
        )
        db.execute(
            delete(PathwayRecommendation).where(
                PathwayRecommendation.id.in_(pathway_ids)
            )
        )
    else:
        db.execute(
            delete(TeacherIntervention).where(
                TeacherIntervention.student_id == student.id
            )
        )
        db.execute(
            delete(PathwayVersion).where(PathwayVersion.student_id == student.id)
        )
    session_ids = list(
        db.scalars(
            select(TutoringSession.id).where(
                TutoringSession.student_id == student.id
            )
        )
    )
    if session_ids:
        db.execute(
            delete(TutoringResponse).where(
                TutoringResponse.session_id.in_(session_ids)
            )
        )
    attempt_ids = list(
        db.scalars(
            select(AssessmentAttempt.id).where(
                AssessmentAttempt.student_id == student.id
            )
        )
    )
    delete_attempt_records(db, attempt_ids)
    for model in (
        CognitiveLoadPrediction,
        MisconceptionHistory,
        LearningSummary,
        TutoringSession,
        MentalEffortRating,
        InteractionLog,
        MasteryRecord,
        LearningGap,
        AssessmentAssignment,
        ConsentRecord,
    ):
        db.execute(delete(model).where(model.student_id == student.id))
    db.execute(delete(StudentProfile).where(StudentProfile.user_id == student.id))
    db.execute(
        update(LoginHistory)
        .where(LoginHistory.user_id == student.id)
        .values(user_id=None)
    )
    db.execute(
        update(AuditLog)
        .where(AuditLog.actor_id == student.id)
        .values(actor_id=None)
    )
    db.execute(delete(User).where(User.id == student.id))


@app.post("/api/teacher/students/{student_id}/actions")
def manage_student_account(
    student_id: int,
    payload: StudentActionInput,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    student = db.get(User, student_id)
    if not student or student.role != "student":
        raise HTTPException(status_code=404, detail="Student not found")

    temporary_password = None
    if payload.action == "reset_password":
        temporary_password = f"NX!{secrets.token_urlsafe(8)}7a"
        student.password_hash = hash_password(temporary_password)
        student.must_change_password = True
    elif payload.action == "reactivate":
        student.account_status = "Active"
        student.is_active = True
    elif payload.action == "deactivate":
        student.account_status = "Deactivated"
        student.is_active = False
    elif payload.action == "archive":
        student.account_status = "Archived"
        student.is_active = False

    removed_student = None
    if payload.action == "remove":
        student.account_status = "Archived"
        student.is_active = False
        removed_student = user_payload(student)

    audit(
        db,
        teacher.id,
        f"student.{payload.action}",
        "user",
        student.id,
        {
            "student_id": student.participant_code,
            "student_name": student.display_name,
            "reason": payload.reason,
            "record_preservation": (
                "Assessment, mastery, learning-gap, pathway, and audit records preserved"
                if payload.action in {"archive", "remove"}
                else None
            ),
        },
    )
    if payload.action == "remove":
        permanently_delete_student(db, student)
    db.commit()
    return {
        "ok": True,
        "message": (
            "Temporary password created. The student must change it after signing in."
            if temporary_password
            else f"{removed_student['display_name']} and linked learner records were permanently removed."
            if removed_student
            else f"{student.display_name}'s account is now {student.account_status}."
        ),
        "temporary_password": temporary_password,
        "student": removed_student or user_payload(student),
        "deleted": bool(removed_student),
    }


@app.get("/api/teacher/students/{student_id}")
def teacher_student_detail(
    student_id: int,
    db: Session = Depends(get_db),
    _teacher: User = Depends(require_role("teacher")),
):
    student = db.get(User, student_id)
    if not student or student.role != "student":
        raise HTTPException(status_code=404, detail="Student not found")
    snapshot = student_snapshot(db, student)
    misconception_rows = db.execute(
        select(
            MisconceptionHistory,
            Misconception,
            Question,
        )
        .join(Misconception, Misconception.id == MisconceptionHistory.misconception_id)
        .join(Question, Question.id == MisconceptionHistory.question_id)
        .where(MisconceptionHistory.student_id == student.id)
        .order_by(MisconceptionHistory.created_at.desc())
        .limit(30)
    ).all()
    snapshot["misconception_history"] = [
        {
            "id": history.id,
            "misconception_id": misconception.id,
            "code": misconception.code,
            "name": misconception.name,
            "concept_id": misconception.concept_id,
            "concept": db.get(Concept, misconception.concept_id).name,
            "question": question.prompt,
            "evidence_count": history.evidence_count,
            "confidence_level": history.confidence_level,
            "detected_at": history.created_at,
            "resolved_at": history.resolved_at,
            "remediation_instruction": misconception.remediation_instruction,
            "suggested_activity_id": misconception.suggested_activity_id,
        }
        for history, misconception, question in misconception_rows
    ]
    snapshot["learning_summaries"] = [
        {
            "id": row.id,
            "activity": db.get(Activity, row.activity_id).title,
            "attempt_id": row.attempt_id,
            "summary": row.summary,
            "created_at": row.created_at,
        }
        for row in db.scalars(
            select(LearningSummary)
            .where(LearningSummary.student_id == student.id)
            .order_by(LearningSummary.created_at.desc())
            .limit(12)
        )
    ]
    snapshot["interventions"] = [
        {
            "id": row.id,
            "action_type": row.action_type,
            "concept_id": row.concept_id,
            "misconception_id": row.misconception_id,
            "pathway_id": row.pathway_id,
            "assigned_activity_id": row.assigned_activity_id,
            "note": row.note,
            "status": row.status,
            "resolved_at": row.resolved_at,
            "created_at": row.created_at,
        }
        for row in db.scalars(
            select(TeacherIntervention)
            .where(TeacherIntervention.student_id == student.id)
            .order_by(TeacherIntervention.created_at.desc())
        )
    ]
    return snapshot


@app.post("/api/teacher/interventions", status_code=201)
def create_teacher_intervention(
    payload: InterventionInput,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    student = db.get(User, payload.student_id)
    if not student or student.role != "student":
        raise HTTPException(status_code=404, detail="Student not found")
    if payload.concept_id and not db.get(Concept, payload.concept_id):
        raise HTTPException(status_code=404, detail="Concept not found")
    if payload.misconception_id and not db.get(Misconception, payload.misconception_id):
        raise HTTPException(status_code=404, detail="Misconception not found")
    if payload.assigned_activity_id and not db.get(Activity, payload.assigned_activity_id):
        raise HTTPException(status_code=404, detail="Activity not found")
    row = TeacherIntervention(
        teacher_id=teacher.id,
        **payload.model_dump(),
    )
    db.add(row)
    db.flush()
    if payload.action_type == "Resolve misconception" and payload.misconception_id:
        now = datetime.now(timezone.utc)
        for history in db.scalars(
            select(MisconceptionHistory).where(
                MisconceptionHistory.student_id == student.id,
                MisconceptionHistory.misconception_id == payload.misconception_id,
                MisconceptionHistory.resolved_at.is_(None),
            )
        ):
            history.resolved_at = now
        row.status = "Resolved"
        row.resolved_at = now
    audit(
        db,
        teacher.id,
        "intervention.created",
        "teacher_intervention",
        row.id,
        payload.model_dump(),
    )
    db.commit()
    return {"id": row.id, "status": row.status}


@app.post("/api/teacher/interventions/{intervention_id}/resolve")
def resolve_teacher_intervention(
    intervention_id: int,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    row = db.get(TeacherIntervention, intervention_id)
    if not row:
        raise HTTPException(status_code=404, detail="Intervention not found")
    row.status = "Resolved"
    row.resolved_at = datetime.now(timezone.utc)
    audit(db, teacher.id, "intervention.resolved", "teacher_intervention", row.id)
    db.commit()
    return {"id": row.id, "status": row.status}


def learner_topic_evidence(db: Session, student: User, concept: Concept) -> dict[str, Any]:
    mastery = latest_mastery_map(db, student.id)
    edge_pairs = db.execute(
        select(
            PrerequisiteEdge.prerequisite_concept_id,
            PrerequisiteEdge.succeeding_concept_id,
        ).where(PrerequisiteEdge.active.is_(True))
    ).all()
    prerequisite_ids = prerequisite_ancestors(edge_pairs, concept.id)
    incorrect_rows = db.execute(
        select(ItemResponse, Question, AssessmentAttempt)
        .join(Question, Question.id == ItemResponse.question_id)
        .join(AssessmentAttempt, AssessmentAttempt.id == ItemResponse.attempt_id)
        .where(
            AssessmentAttempt.student_id == student.id,
            Question.concept_id == concept.id,
            ItemResponse.is_correct.is_(False),
        )
        .order_by(AssessmentAttempt.submitted_at.desc())
        .limit(20)
    ).all()
    incorrect = []
    for response, question, attempt in incorrect_rows:
        selected = (
            db.get(AnswerChoice, response.selected_choice_id)
            if response.selected_choice_id
            else None
        )
        correct = db.scalar(
            select(AnswerChoice).where(
                AnswerChoice.question_id == question.id,
                AnswerChoice.is_correct.is_(True),
            )
        )
        incorrect.append(
            {
                "question": question.prompt,
                "learner_answer": selected.text if selected else response.response_text,
                "correct_answer": correct.text if correct else question.correct_answer,
                "explanation": question.explanation or question.feedback,
                "likely_mistake": (
                    db.get(Misconception, selected.misconception_id).explanation
                    if selected
                    and selected.misconception_id
                    and selected.mapping_status in {"Teacher reviewed", "Validated"}
                    and db.get(Misconception, selected.misconception_id)
                    else "No teacher-reviewed distractor mapping supports a specific misconception diagnosis."
                ),
                "response_seconds": response.response_seconds,
                "hint_opened": response.hint_opened,
                "attempt_id": attempt.id,
                "submitted_at": attempt.submitted_at,
            }
        )
    attempts = list(
        db.scalars(
            select(AssessmentAttempt)
            .join(ItemResponse, ItemResponse.attempt_id == AssessmentAttempt.id)
            .join(Question, Question.id == ItemResponse.question_id)
            .where(
                AssessmentAttempt.student_id == student.id,
                Question.concept_id == concept.id,
            )
            .distinct()
            .order_by(AssessmentAttempt.submitted_at.desc())
        )
    )
    interaction_rows = list(
        db.scalars(
            select(InteractionLog)
            .where(
                InteractionLog.student_id == student.id,
                InteractionLog.concept_id == concept.id,
            )
            .order_by(InteractionLog.submission_time.desc())
            .limit(12)
        )
    )
    mistake_counts: dict[str, int] = defaultdict(int)
    for item in incorrect:
        mistake_counts[item["likely_mistake"]] += 1
    active_step = db.scalar(
        select(PathwayStep)
        .join(
            PathwayRecommendation,
            PathwayRecommendation.id == PathwayStep.pathway_id,
        )
        .where(
            PathwayRecommendation.student_id == student.id,
            PathwayRecommendation.active.is_(True),
            PathwayStep.concept_id == concept.id,
        )
        .order_by(PathwayRecommendation.created_at.desc())
    )
    threshold = float(get_setting(db, "mastery_threshold"))
    return {
        "student": user_payload(student),
        "concept": concept_payload(concept),
        "mastery": (
            {
                "score": mastery[concept.id].mastery_score,
                "classification": mastery[concept.id].classification,
                "updated_at": mastery[concept.id].created_at,
            }
            if concept.id in mastery
            else None
        ),
        "missing_prerequisites": [
            {
                "concept_id": item_id,
                "concept": db.get(Concept, item_id).name,
                "mastery_score": (
                    mastery[item_id].mastery_score if item_id in mastery else None
                ),
            }
            for item_id in topological_order(prerequisite_ids, edge_pairs)
            if item_id not in mastery or mastery[item_id].mastery_score < threshold
        ],
        "incorrect_responses": incorrect,
        "common_errors": [
            {"error": message, "count": count}
            for message, count in sorted(
                mistake_counts.items(), key=lambda item: (-item[1], item[0])
            )
        ],
        "performance": {
            "average_response_seconds": (
                sum(item.average_response_seconds for item in interaction_rows)
                / len(interaction_rows)
                if interaction_rows
                else 0
            ),
            "attempts": max(
                (item.number_of_attempts for item in interaction_rows), default=0
            ),
            "skips": sum(item.skipped_items for item in interaction_rows),
            "hints": sum(item.hint_usage_count for item in interaction_rows),
            "predicted_cognitive_load": (
                active_step.predicted_load_index if active_step else None
            ),
        },
        "attempts": [
            {
                "id": attempt.id,
                "activity": db.get(Activity, attempt.activity_id).title,
                "accuracy": attempt.accuracy,
                "total_seconds": attempt.total_seconds,
                "skipped_items": attempt.skipped_items,
                "hint_usage_count": attempt.hint_usage_count,
                "submitted_at": attempt.submitted_at,
            }
            for attempt in attempts
        ],
        "completed_activities": [
            {
                "activity": db.get(Activity, attempt.activity_id).title,
                "accuracy": attempt.accuracy,
                "submitted_at": attempt.submitted_at,
            }
            for attempt in attempts
            if attempt.accuracy >= 0.60
        ],
        "prior_pathways": [
            serialize_pathway(db, pathway)
            for pathway in db.scalars(
                select(PathwayRecommendation)
                .where(
                    PathwayRecommendation.student_id == student.id,
                    PathwayRecommendation.target_concept_id == concept.id,
                )
                .order_by(PathwayRecommendation.created_at.desc())
                .limit(8)
            )
        ],
    }


@app.get("/api/teacher/students/{student_id}/topics/{concept_id}")
def teacher_student_topic(
    student_id: int,
    concept_id: int,
    db: Session = Depends(get_db),
    _teacher: User = Depends(require_role("teacher")),
):
    student = db.get(User, student_id)
    concept = db.get(Concept, concept_id)
    if not student or student.role != "student" or not concept:
        raise HTTPException(status_code=404, detail="Student topic not found")
    return learner_topic_evidence(db, student, concept)


@app.post("/api/teacher/students/{student_id}/topics/{concept_id}/pathway-preview")
def preview_student_pathway(
    student_id: int,
    concept_id: int,
    payload: PathwayPreviewInput,
    db: Session = Depends(get_db),
    _teacher: User = Depends(require_role("teacher")),
):
    student = db.get(User, student_id)
    concept = db.get(Concept, concept_id)
    if not student or student.role != "student" or not concept:
        raise HTTPException(status_code=404, detail="Student topic not found")
    result = build_pathway_preview(
        db,
        student,
        concept_id,
        None if payload.difficulty == "Auto" else payload.difficulty,
    )
    if not result.get("steps"):
        raise HTTPException(status_code=422, detail="No connected learning activities are available")
    return result


@app.post("/api/teacher/students/{student_id}/pathways/assign", status_code=201)
def assign_student_pathway(
    student_id: int,
    payload: PathwayAssignmentInput,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    student = db.get(User, student_id)
    target = db.get(Concept, payload.target_concept_id)
    if not student or student.role != "student" or not target:
        raise HTTPException(status_code=404, detail="Student or target concept not found")
    edge_pairs = db.execute(
        select(
            PrerequisiteEdge.prerequisite_concept_id,
            PrerequisiteEdge.succeeding_concept_id,
        ).where(PrerequisiteEdge.active.is_(True))
    ).all()
    allowed = prerequisite_ancestors(edge_pairs, target.id) | {target.id}
    ordered_steps = sorted(payload.steps, key=lambda item: item.position)
    if len({item.position for item in ordered_steps}) != len(ordered_steps):
        raise HTTPException(status_code=422, detail="Pathway positions must be unique")
    if len({(item.concept_id, item.activity_id) for item in ordered_steps}) != len(ordered_steps):
        raise HTTPException(status_code=422, detail="A pathway cannot contain duplicate activities")
    positions = {item.concept_id: item.position for item in ordered_steps}
    for item in ordered_steps:
        if item.concept_id not in allowed:
            raise HTTPException(
                status_code=422,
                detail="Every pathway concept must be connected to the selected target",
            )
        link = db.scalar(
            select(ActivityConcept).where(
                ActivityConcept.activity_id == item.activity_id,
                ActivityConcept.concept_id == item.concept_id,
            )
        )
        activity = db.get(Activity, item.activity_id)
        if not link or not activity or not activity.active:
            raise HTTPException(
                status_code=422,
                detail="A selected activity does not belong to its pathway concept",
            )
    for prerequisite_id, succeeding_id in edge_pairs:
        if (
            prerequisite_id in positions
            and succeeding_id in positions
            and positions[prerequisite_id] >= positions[succeeding_id]
        ):
            raise HTTPException(
                status_code=422,
                detail="Pathway order violates the prerequisite graph",
            )
    previous = db.scalar(
        select(PathwayRecommendation)
        .where(
            PathwayRecommendation.student_id == student.id,
            PathwayRecommendation.active.is_(True),
            PathwayRecommendation.selected.is_(True),
        )
        .order_by(PathwayRecommendation.created_at.desc())
    )
    for old in db.scalars(
        select(PathwayRecommendation).where(
            PathwayRecommendation.student_id == student.id,
            PathwayRecommendation.active.is_(True),
        )
    ):
        old.active = False
        old.selected = False
    step_data = []
    for item in ordered_steps:
        concept = db.get(Concept, item.concept_id)
        activity = db.get(Activity, item.activity_id)
        prediction = predict_activity_load(
            db, student.id, activity, concept.id, concept.difficulty
        )
        step_data.append((item, concept, activity, prediction))
    average_load = sum(row[3]["index"] for row in step_data) / len(step_data)
    total_minutes = sum(row[2].estimated_minutes for row in step_data)
    pathway = PathwayRecommendation(
        student_id=student.id,
        target_concept_id=target.id,
        label=clean_content(payload.label, 80),
        selected=True,
        gap_coverage=1,
        predicted_cognitive_load=average_load,
        normalized_learning_time=1,
        adaptive_pathway_score=max(0, 1 - average_load),
        total_minutes=total_minutes,
        cognitive_load_category=(
            "High" if average_load >= 0.67 else "Moderate" if average_load >= 0.34 else "Low"
        ),
        cognitive_load_probabilities={},
        explanation=(
            f"{teacher.display_name} assigned this {payload.difficulty.lower()} "
            f"pathway for {target.name}. Each step is connected by the active prerequisite graph."
        ),
        feature_explanation={
            "method": "Teacher-reviewed adaptive preview",
            "assigned_by": teacher.display_name,
        },
        source_type="Teacher",
        assigned_by=teacher.id,
        assigned_at=datetime.now(timezone.utc),
        due_at=payload.due_at,
        teacher_note=clean_content(payload.teacher_note, 2000),
        learner_notified=True,
        supersedes_pathway_id=previous.id if previous else None,
        difficulty_override=payload.difficulty,
        active=True,
        is_demo=student.is_demo,
    )
    db.add(pathway)
    db.flush()
    for position, (_item, concept, activity, prediction) in enumerate(step_data, start=1):
        db.add(
            PathwayStep(
                pathway_id=pathway.id,
                concept_id=concept.id,
                activity_id=activity.id,
                position=position,
                predicted_load_index=prediction["index"],
                selection_reason=(
                    f"Teacher selected {activity.title} for {concept.name}; the concept "
                    f"is connected to {target.name} in the prerequisite graph."
                ),
                content=learning_content_for(concept, activity, payload.difficulty),
                required=True,
            )
        )
    audit(
        db,
        teacher.id,
        "pathway.assigned",
        "pathway_recommendation",
        pathway.id,
        {
            "student_id": student.participant_code,
            "student_name": student.display_name,
            "target": target.name,
            "assigned_at": pathway.assigned_at.isoformat(),
            "due_at": payload.due_at.isoformat() if payload.due_at else None,
            "supersedes_pathway_id": pathway.supersedes_pathway_id,
        },
    )
    db.commit()
    db.refresh(pathway)
    return serialize_pathway(db, pathway)


@app.get("/api/teacher/concepts")
def teacher_concepts(
    include_archived: bool = True,
    db: Session = Depends(get_db),
    _teacher: User = Depends(require_role("teacher")),
):
    query = select(Concept)
    if not include_archived:
        query = query.where(Concept.active.is_(True))
    return [concept_payload(item) for item in db.scalars(query.order_by(Concept.subject, Concept.name))]


@app.post("/api/teacher/concepts", status_code=201)
def create_concept(
    payload: ConceptInput,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    values = payload.model_dump()
    values.update(
        code=payload.code.strip().upper(),
        name=payload.name.strip(),
        subject=payload.subject.strip(),
        description=payload.description.strip(),
    )
    concept = Concept(**values)
    db.add(concept)
    try:
        db.flush()
    except IntegrityError:
        db.rollback()
        raise HTTPException(status_code=409, detail="Concept code already exists")
    audit(db, teacher.id, "concept.created", "concept", concept.id, payload.model_dump())
    db.commit()
    db.refresh(concept)
    return concept_payload(concept)


@app.put("/api/teacher/concepts/{concept_id}")
def update_concept(
    concept_id: int,
    payload: ConceptInput,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    concept = db.get(Concept, concept_id)
    if not concept:
        raise HTTPException(status_code=404, detail="Concept not found")
    code = payload.code.strip().upper()
    duplicate = db.scalar(
        select(Concept.id).where(Concept.code == code, Concept.id != concept.id)
    )
    if duplicate:
        raise HTTPException(status_code=409, detail="Concept code already exists")
    values = payload.model_dump()
    values.update(
        code=code,
        name=payload.name.strip(),
        subject=payload.subject.strip(),
        description=payload.description.strip(),
    )
    for key, value in values.items():
        setattr(concept, key, value)
    audit(db, teacher.id, "concept.updated", "concept", concept.id, payload.model_dump())
    db.commit()
    return concept_payload(concept)


@app.post("/api/teacher/concepts/{concept_id}/{action}")
def archive_restore_concept(
    concept_id: int,
    action: str,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    if action not in {"archive", "restore"}:
        raise HTTPException(status_code=400, detail="Action must be archive or restore")
    concept = db.get(Concept, concept_id)
    if not concept:
        raise HTTPException(status_code=404, detail="Concept not found")
    concept.active = action == "restore"
    audit(db, teacher.id, f"concept.{action}d", "concept", concept.id)
    db.commit()
    return concept_payload(concept)


@app.get("/api/teacher/graph")
def teacher_graph(
    db: Session = Depends(get_db), _teacher: User = Depends(require_role("teacher"))
):
    concepts = list(db.scalars(select(Concept).order_by(Concept.subject, Concept.name)))
    edges = list(
        db.scalars(
            select(PrerequisiteEdge).where(PrerequisiteEdge.active.is_(True))
        )
    )
    return {
        "nodes": [concept_payload(concept) for concept in concepts],
        "edges": [
            {
                "id": edge.id,
                "source": edge.prerequisite_concept_id,
                "target": edge.succeeding_concept_id,
            }
            for edge in edges
        ],
    }


@app.post("/api/teacher/graph/edges", status_code=201)
def create_edge(
    payload: EdgeInput,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    if not db.get(Concept, payload.prerequisite_concept_id) or not db.get(
        Concept, payload.succeeding_concept_id
    ):
        raise HTTPException(status_code=404, detail="One or both concepts were not found")
    edges = db.execute(
        select(
            PrerequisiteEdge.prerequisite_concept_id,
            PrerequisiteEdge.succeeding_concept_id,
        ).where(PrerequisiteEdge.active.is_(True))
    ).all()
    if would_create_cycle(
        edges, payload.prerequisite_concept_id, payload.succeeding_concept_id
    ):
        raise HTTPException(
            status_code=400,
            detail="This relationship would create a cycle. The knowledge graph must remain acyclic.",
        )
    existing = db.scalar(
        select(PrerequisiteEdge).where(
            PrerequisiteEdge.prerequisite_concept_id
            == payload.prerequisite_concept_id,
            PrerequisiteEdge.succeeding_concept_id == payload.succeeding_concept_id,
        )
    )
    if existing:
        if existing.active:
            raise HTTPException(status_code=409, detail="Relationship already exists")
        existing.active = True
        edge = existing
    else:
        edge = PrerequisiteEdge(**payload.model_dump())
        db.add(edge)
    db.flush()
    audit(db, teacher.id, "graph.edge.created", "prerequisite_edge", edge.id)
    db.commit()
    return {"id": edge.id, "source": edge.prerequisite_concept_id, "target": edge.succeeding_concept_id}


@app.delete("/api/teacher/graph/edges/{edge_id}")
def archive_edge(
    edge_id: int,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    edge = db.get(PrerequisiteEdge, edge_id)
    if not edge:
        raise HTTPException(status_code=404, detail="Relationship not found")
    edge.active = False
    audit(db, teacher.id, "graph.edge.archived", "prerequisite_edge", edge.id)
    db.commit()
    return {"ok": True}


@app.get("/api/teacher/activities")
def teacher_activities(
    include_archived: bool = True,
    db: Session = Depends(get_db),
    _teacher: User = Depends(require_role("teacher")),
):
    query = select(Activity)
    if not include_archived:
        query = query.where(Activity.active.is_(True))
    return [
        activity_payload(
            db, activity, include_questions=True, include_correct_answers=True
        )
        for activity in db.scalars(query.order_by(Activity.title))
    ]


@app.post("/api/teacher/activities", status_code=201)
def create_activity(
    payload: ActivityInput,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    invalid = [concept_id for concept_id in payload.concept_ids if not db.get(Concept, concept_id)]
    if invalid:
        raise HTTPException(status_code=400, detail="One or more concepts were not found")
    data = payload.model_dump(exclude={"concept_ids"})
    activity = Activity(**data, is_demo=teacher.is_demo)
    db.add(activity)
    db.flush()
    for concept_id in payload.concept_ids:
        db.add(ActivityConcept(activity_id=activity.id, concept_id=concept_id))
    audit(db, teacher.id, "activity.created", "activity", activity.id, data)
    db.commit()
    return activity_payload(
        db, activity, include_questions=True, include_correct_answers=True
    )


@app.put("/api/teacher/activities/{activity_id}")
def update_activity(
    activity_id: int,
    payload: ActivityInput,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    activity = db.get(Activity, activity_id)
    if not activity:
        raise HTTPException(status_code=404, detail="Activity not found")
    data = payload.model_dump(exclude={"concept_ids"})
    for key, value in data.items():
        setattr(activity, key, value)
    for link in db.scalars(
        select(ActivityConcept).where(ActivityConcept.activity_id == activity.id)
    ):
        db.delete(link)
    for concept_id in payload.concept_ids:
        db.add(ActivityConcept(activity_id=activity.id, concept_id=concept_id))
    audit(db, teacher.id, "activity.updated", "activity", activity.id, data)
    db.commit()
    return activity_payload(
        db, activity, include_questions=True, include_correct_answers=True
    )


@app.post("/api/teacher/activities/{activity_id}/{action}")
def archive_restore_activity(
    activity_id: int,
    action: str,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    if action not in {"archive", "restore"}:
        raise HTTPException(status_code=400, detail="Action must be archive or restore")
    activity = db.get(Activity, activity_id)
    if not activity:
        raise HTTPException(status_code=404, detail="Activity not found")
    activity.active = action == "restore"
    audit(db, teacher.id, f"activity.{action}d", "activity", activity.id)
    db.commit()
    return activity_payload(db, activity)


def permanently_delete_activity(
    db: Session,
    activity: Activity,
    confirm_learner_record_deletion: bool,
) -> dict[str, int]:
    attempt_ids = list(db.scalars(select(AssessmentAttempt.id).where(
        AssessmentAttempt.activity_id == activity.id
    )))
    result_count = db.scalar(select(func.count(AssessmentAttempt.id)).where(
        AssessmentAttempt.activity_id == activity.id,
        AssessmentAttempt.submitted_at.is_not(None),
    )) or 0
    assessment_ids = list(db.scalars(select(Assessment.id).where(
        Assessment.activity_id == activity.id
    )))
    assigned_count = 0
    if assessment_ids:
        assigned_count = db.scalar(select(func.count(AssessmentAssignment.id)).where(
            AssessmentAssignment.assessment_id.in_(assessment_ids)
        )) or 0
    pathway_steps = db.scalar(select(func.count(PathwayStep.id)).where(
        PathwayStep.activity_id == activity.id
    )) or 0
    if (attempt_ids or assigned_count or pathway_steps) and not confirm_learner_record_deletion:
        raise HTTPException(
            status_code=409,
            detail=(
                f'Deleting "{activity.title}" affects {assigned_count} assignment(s), '
                f'{len(attempt_ids)} learner attempt(s), {result_count} result(s), and '
                f'{pathway_steps} saved pathway step(s). Confirm related-record deletion to continue.'
            ),
        )
    if attempt_ids:
        delete_attempt_records(db, attempt_ids)
    session_ids = list(db.scalars(select(TutoringSession.id).where(
        TutoringSession.activity_id == activity.id
    )))
    if session_ids:
        db.execute(delete(TutoringResponse).where(TutoringResponse.session_id.in_(session_ids)))
        db.execute(delete(MisconceptionHistory).where(MisconceptionHistory.tutoring_session_id.in_(session_ids)))
        db.execute(delete(LearningSummary).where(LearningSummary.tutoring_session_id.in_(session_ids)))
        db.execute(delete(TutoringSession).where(TutoringSession.id.in_(session_ids)))
    db.execute(delete(LearningSummary).where(LearningSummary.activity_id == activity.id))
    if assessment_ids:
        db.execute(delete(AssessmentAssignment).where(AssessmentAssignment.assessment_id.in_(assessment_ids)))
        db.execute(delete(AssessmentQuestion).where(AssessmentQuestion.assessment_id.in_(assessment_ids)))
        db.execute(delete(Assessment).where(Assessment.id.in_(assessment_ids)))
    question_ids = list(db.scalars(select(Question.id).where(Question.activity_id == activity.id)))
    if question_ids:
        db.execute(delete(AssessmentQuestion).where(AssessmentQuestion.question_id.in_(question_ids)))
        db.execute(delete(AnswerChoice).where(AnswerChoice.question_id.in_(question_ids)))
        db.execute(delete(Question).where(Question.id.in_(question_ids)))
    db.execute(delete(PathwayStep).where(PathwayStep.activity_id == activity.id))
    db.execute(update(Misconception).where(Misconception.suggested_activity_id == activity.id).values(suggested_activity_id=None))
    db.execute(update(TeacherIntervention).where(TeacherIntervention.assigned_activity_id == activity.id).values(assigned_activity_id=None))
    db.execute(delete(ActivityConcept).where(ActivityConcept.activity_id == activity.id))
    db.delete(activity)
    db.flush()
    return {
        "assignments": int(assigned_count),
        "attempts": len(attempt_ids),
        "results": int(result_count),
        "pathway_steps": int(pathway_steps),
    }


@app.delete("/api/teacher/activities/bulk-delete")
def bulk_delete_activities(
    activity_ids: list[int] = Query(min_length=1, max_length=500),
    confirm_learner_record_deletion: bool = False,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    unique_ids = list(dict.fromkeys(activity_ids))
    activities = list(db.scalars(select(Activity).where(Activity.id.in_(unique_ids))))
    if len(activities) != len(unique_ids):
        raise HTTPException(status_code=404, detail="One or more activities were not found")
    totals = {"assignments": 0, "attempts": 0, "results": 0, "pathway_steps": 0}
    for activity in activities:
        removed = permanently_delete_activity(db, activity, confirm_learner_record_deletion)
        for key, value in removed.items():
            totals[key] += value
    audit(db, teacher.id, "activity.bulk_deleted", "activity", details={"activity_ids": unique_ids, **totals})
    db.commit()
    return {"ok": True, "deleted": len(unique_ids), "related_records_removed": totals}


@app.post("/api/teacher/activities/bulk")
def bulk_activity_action(
    payload: ActivityBulkInput,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    unique_ids = list(dict.fromkeys(payload.activity_ids))
    activities = list(db.scalars(select(Activity).where(Activity.id.in_(unique_ids))))
    if len(activities) != len(unique_ids):
        raise HTTPException(status_code=404, detail="One or more activities were not found")
    if payload.action == "archive":
        for activity in activities:
            activity.active = False
        audit(db, teacher.id, "activity.bulk_archived", "activity", details={"activity_ids": unique_ids})
        db.commit()
        return {"ok": True, "archived": len(unique_ids)}
    totals = {"assignments": 0, "attempts": 0, "results": 0, "pathway_steps": 0}
    for activity in activities:
        removed = permanently_delete_activity(db, activity, payload.confirm_learner_record_deletion)
        for key, value in removed.items():
            totals[key] += value
    audit(db, teacher.id, "activity.bulk_deleted", "activity", details={"activity_ids": unique_ids, **totals})
    db.commit()
    return {"ok": True, "deleted": len(unique_ids), "related_records_removed": totals}


@app.delete("/api/teacher/activities/{activity_id}")
def delete_activity(
    activity_id: int,
    confirm_learner_record_deletion: bool = False,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    activity = db.get(Activity, activity_id)
    if not activity:
        raise HTTPException(status_code=404, detail="Activity not found")
    title = activity.title
    removed = permanently_delete_activity(db, activity, confirm_learner_record_deletion)
    audit(db, teacher.id, "activity.deleted", "activity", activity_id, {"title": title, **removed})
    db.commit()
    return {"ok": True, "id": activity_id, "title": title, "related_records_removed": removed}


@app.post("/api/teacher/questions", status_code=201)
def create_question(
    payload: QuestionInput,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    if not db.get(Activity, payload.activity_id) or not db.get(Concept, payload.concept_id):
        raise HTTPException(status_code=404, detail="Activity or concept not found")
    position = db.scalar(
        select(func.count(Question.id)).where(Question.activity_id == payload.activity_id)
    ) or 0
    question = Question(
        **payload.model_dump(exclude={"choices"}),
        question_type="Multiple choice",
        correct_answer=next(
            choice.text for choice in payload.choices if choice.is_correct
        ),
        explanation=payload.feedback,
        source_type="Manually created",
        status="Published",
        created_by=teacher.id,
        position=position + 1,
    )
    db.add(question)
    db.flush()
    for index, choice in enumerate(payload.choices, start=1):
        db.add(
            AnswerChoice(
                question_id=question.id,
                text=choice.text,
                is_correct=choice.is_correct,
                position=index,
                misconception_id=choice.misconception_id,
                misconception_confidence=choice.misconception_confidence,
                mapping_status="Validated" if choice.is_correct else choice.mapping_status,
            )
        )
    audit(db, teacher.id, "question.created", "question", question.id)
    db.commit()
    return {"id": question.id, "ok": True}


@app.put("/api/teacher/questions/{question_id}")
def update_question(
    question_id: int,
    payload: QuestionInput,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    question = db.get(Question, question_id)
    if not question:
        raise HTTPException(status_code=404, detail="Question not found")
    if not db.get(Activity, payload.activity_id) or not db.get(Concept, payload.concept_id):
        raise HTTPException(status_code=404, detail="Activity or concept not found")
    for key, value in payload.model_dump(exclude={"choices"}).items():
        setattr(question, key, value)
    question.question_type = "Multiple choice"
    question.correct_answer = next(
        choice.text for choice in payload.choices if choice.is_correct
    )
    question.explanation = payload.feedback
    question.status = "Published"
    for choice in list(
        db.scalars(
            select(AnswerChoice).where(AnswerChoice.question_id == question.id)
        )
    ):
        db.delete(choice)
    db.flush()
    for index, choice in enumerate(payload.choices, start=1):
        db.add(
            AnswerChoice(
                question_id=question.id,
                text=choice.text,
                is_correct=choice.is_correct,
                position=index,
                misconception_id=choice.misconception_id,
                misconception_confidence=choice.misconception_confidence,
                mapping_status="Validated" if choice.is_correct else choice.mapping_status,
            )
        )
    audit(db, teacher.id, "question.updated", "question", question.id)
    db.commit()
    return {"id": question.id, "ok": True}


@app.get("/api/teacher/ai/configuration")
def teacher_ai_configuration(
    _teacher: User = Depends(require_role("teacher")),
):
    configuration = ai_configuration()
    return {
        "configured": configuration["configured"],
        "provider": configuration["provider"],
        "model": configuration["model"],
        "message": (
            "AI-assisted module analysis and question generation are configured."
            if configuration["configured"]
            else "Set AI_PROVIDER, AI_MODEL, and AI_API_KEY in the backend environment."
        ),
    }


@app.get("/api/teacher/documents")
def teacher_documents(
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    documents = list(
        db.scalars(
            select(UploadedDocument)
            .where(
                UploadedDocument.uploaded_by == teacher.id,
                UploadedDocument.processing_status != "Archived",
            )
            .order_by(UploadedDocument.created_at.desc())
        )
    )
    return [uploaded_document_payload(document) for document in documents]


@app.post("/api/teacher/documents", status_code=201)
async def upload_learning_document(
    request: Request,
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    enforce_rate_limit(request, "document-upload", 20, 300)
    maximum = int(os.getenv("MAX_UPLOAD_SIZE_MB", "10")) * 1024 * 1024
    filename = Path(file.filename or "learning-material").name
    filename = re.sub(r"[^A-Za-z0-9._ ()-]", "_", filename)[:255]
    extension = Path(filename).suffix.lower()
    allowed = {
        ".pdf": {"application/pdf", "application/octet-stream"},
        ".docx": {
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/octet-stream",
        },
        ".pptx": {
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/octet-stream",
        },
        ".txt": {"text/plain", "application/octet-stream"},
    }
    if extension not in allowed or (
        file.content_type and file.content_type not in allowed[extension]
    ):
        raise HTTPException(
            status_code=415,
            detail="Upload a PDF, DOCX, PPTX, or UTF-8 TXT file",
        )
    data = await file.read(maximum + 1)
    await file.close()
    if not data:
        raise HTTPException(status_code=422, detail="The uploaded file is empty")
    if len(data) > maximum:
        raise HTTPException(
            status_code=413,
            detail=f"Files must be {maximum // (1024 * 1024)} MB or smaller",
        )
    if extension == ".pdf" and not data.startswith(b"%PDF"):
        raise HTTPException(status_code=422, detail="The PDF signature is invalid")
    try:
        extracted_text = await asyncio.wait_for(
            asyncio.to_thread(extract_document_text, extension, data),
            timeout=float(os.getenv("UPLOAD_PROCESSING_TIMEOUT_SECONDS", "45")),
        )
    except TimeoutError as error:
        raise HTTPException(
            status_code=504,
            detail="Document processing took too long. Try a smaller file.",
        ) from error
    document = UploadedDocument(
        original_filename=filename,
        file_type=extension.removeprefix(".").upper(),
        file_size=len(data),
        content_sha256=hashlib.sha256(data).hexdigest(),
        uploaded_by=teacher.id,
        processing_status="Ready",
        extracted_text=extracted_text,
        analysis=analyze_material(extracted_text),
    )
    db.add(document)
    db.flush()
    audit(
        db,
        teacher.id,
        "document.uploaded",
        "uploaded_document",
        document.id,
        {
            "filename": filename,
            "file_type": document.file_type,
            "file_size": document.file_size,
        },
    )
    db.commit()
    db.refresh(document)
    return uploaded_document_payload(document)


@app.delete("/api/teacher/documents/{document_id}")
def archive_learning_document(
    document_id: int,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    document = db.get(UploadedDocument, document_id)
    if not document or document.uploaded_by != teacher.id:
        raise HTTPException(status_code=404, detail="Uploaded document not found")
    document.processing_status = "Archived"
    audit(
        db,
        teacher.id,
        "document.archived",
        "uploaded_document",
        document.id,
        {"filename": document.original_filename},
    )
    db.commit()
    return {"ok": True}


@app.post("/api/teacher/documents/{document_id}/generate", status_code=201)
def generate_document_questions(
    document_id: int,
    payload: DocumentGenerationInput,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    document = db.get(UploadedDocument, document_id)
    if (
        not document
        or document.uploaded_by != teacher.id
        or document.processing_status != "Ready"
    ):
        raise HTTPException(status_code=404, detail="Ready uploaded document not found")
    selected_concept = db.get(Concept, payload.concept_id)
    if not selected_concept or not selected_concept.active:
        raise HTTPException(status_code=404, detail="Concept not found")
    passages = source_passages(document.extracted_text)
    if not passages:
        raise HTTPException(
            status_code=422,
            detail="The document does not contain usable question-generation text",
        )
    concept_ids = [selected_concept.id]
    if payload.include_prerequisites:
        edge_pairs = db.execute(
            select(
                PrerequisiteEdge.prerequisite_concept_id,
                PrerequisiteEdge.succeeding_concept_id,
            ).where(PrerequisiteEdge.active.is_(True))
        ).all()
        for concept_id in sorted(prerequisite_ancestors(edge_pairs, selected_concept.id)):
            prerequisite = db.get(Concept, concept_id)
            if not prerequisite:
                continue
            if material_supports_concept(prerequisite, document.extracted_text):
                concept_ids.append(prerequisite.id)
    allowed_concepts = [db.get(Concept, concept_id) for concept_id in concept_ids]
    allowed_concepts = [concept for concept in allowed_concepts if concept]
    ai_result = request_ai_question_generation(document, payload, allowed_concepts)
    ai_questions = ai_result.get("questions", [])
    if not ai_questions:
        raise HTTPException(status_code=502, detail="The AI service did not return any usable questions")
    if len(ai_questions) < payload.number_of_questions:
        raise HTTPException(
            status_code=502,
            detail=(
                f"The AI service returned {len(ai_questions)} questions but "
                f"{payload.number_of_questions} were requested. No partial set was saved."
            ),
        )
    if isinstance(ai_result.get("analysis"), dict):
        document.analysis = {
            **(document.analysis or {}),
            **ai_result["analysis"],
            "method": f"AI-assisted analysis using {ai_configuration()['provider']} / {ai_configuration()['model']}",
        }
    generated = []
    prior_prompts = {
        re.sub(r"\W+", " ", value.casefold()).strip()
        for value in db.scalars(
            select(Question.prompt).where(
                Question.source_document_id == document.id,
                Question.active.is_(True),
            )
        )
    }
    allowed_ids = {concept.id for concept in allowed_concepts}
    for index, raw in enumerate(ai_questions[: payload.number_of_questions]):
        if not isinstance(raw, dict):
            raise HTTPException(status_code=502, detail="The AI service returned a malformed question")
        concept_id = int(raw.get("concept_id") or selected_concept.id)
        concept = db.get(Concept, concept_id) if concept_id in allowed_ids else selected_concept
        question_type = clean_content(str(raw.get("question_type") or payload.question_type), 30)
        if question_type == "Mixed":
            question_type = "Multiple choice"
        choices = raw.get("choices") if isinstance(raw.get("choices"), list) else []
        choices = [
            {
                "text": clean_content(str(choice.get("text", "")), 3000),
                "is_correct": bool(choice.get("is_correct")),
            }
            for choice in choices
            if isinstance(choice, dict) and clean_content(str(choice.get("text", "")), 3000)
        ]
        correct_answer = clean_content(str(raw.get("correct_answer", "")), 3000)
        if question_type in {"Multiple choice", "True or false"} and choices:
            correct_choices = [choice for choice in choices if choice["is_correct"]]
            if len(correct_choices) != 1:
                choices = [
                    {**choice, "is_correct": choice["text"].casefold() == correct_answer.casefold()}
                    for choice in choices
                ]
        validation = raw.get("validation") if isinstance(raw.get("validation"), dict) else {}
        required_checks = {
            "answerable", "answer_correct", "solution_consistent", "formula_units_valid",
            "source_supported", "distractors_valid", "not_duplicate",
            "difficulty_followed", "no_unrelated_topic",
        }
        validation_flags = [
            clean_content(str(note), 250)
            for note in raw.get("validation_notes", [])
            if clean_content(str(note), 250)
        ]
        failed_checks = sorted(check for check in required_checks if validation.get(check) is not True)
        validation_flags.extend(f"AI quality check needs review: {check.replace('_', ' ')}." for check in failed_checks)
        prompt_value = clean_content(str(raw.get("prompt", "")))
        prompt_key = re.sub(r"\W+", " ", prompt_value.casefold()).strip()
        if not prompt_value or not correct_answer:
            raise HTTPException(status_code=502, detail="The AI service returned an incomplete question")
        if prompt_key in prior_prompts:
            validation_flags.append("Duplicate question detected; teacher review is required.")
        source_locator = clean_content(str(raw.get("source_locator", "")), 300)
        if payload.source_grounding and not source_locator:
            validation_flags.append("No source page or section was supplied.")
        validation_status = "Ready for review" if not validation_flags else "Needs review"
        question = Question(
            activity_id=None,
            concept_id=concept.id,
            prompt=prompt_value,
            feedback=clean_content(str(raw.get("explanation", ""))) if payload.include_explanations else "",
            explanation=clean_content(str(raw.get("explanation", ""))) if payload.include_explanations else "",
            hint=clean_content(str(raw.get("hint", "")), 3000) if payload.include_hints else "",
            question_type=question_type,
            correct_answer=correct_answer,
            difficulty_label=payload.difficulty,
            cognitive_level=payload.cognitive_level,
            subject=payload.subject,
            topic=payload.topic,
            learning_competency=clean_content(str(raw.get("learning_competency") or payload.learning_competency), 1000),
            source_type="AI-generated from uploaded material",
            source_document_id=document.id,
            source_locator=source_locator,
            solution_steps=clean_content(str(raw.get("solution_steps", ""))) if payload.include_solutions else "",
            estimated_cognitive_demand=max(0, min(1, float(raw.get("estimated_cognitive_demand", 0.5)))),
            validation_status=validation_status,
            validation_flags=validation_flags,
            generation_metadata={
                "provider": ai_configuration()["provider"],
                "model": ai_configuration()["model"],
                "settings": payload.model_dump(),
                "quality_review": validation,
                "generated_at": datetime.now(timezone.utc).isoformat(),
            },
            distractor_rationales={
                clean_content(str(key), 500): clean_content(str(value), 1000)
                for key, value in (raw.get("distractor_rationales") or {}).items()
            },
            is_calculation=bool(raw.get("is_calculation")),
            status="Draft",
            created_by=teacher.id,
            points=1,
            active=True,
            position=index + 1,
        )
        db.add(question)
        db.flush()
        replace_question_choices(db, question, choices)
        generated.append(question)
        prior_prompts.add(prompt_key)
    audit(
        db,
        teacher.id,
        "questions.generated",
        "uploaded_document",
        document.id,
        {
            "count": len(generated),
            "question_type": payload.question_type,
            "concept_id": payload.concept_id,
            "status": "Draft",
        },
    )
    db.commit()
    return [question_bank_payload(db, question) for question in generated]


def misconception_payload(db: Session, row: Misconception) -> dict[str, Any]:
    return {
        "id": row.id,
        "code": row.code,
        "name": row.name,
        "concept_id": row.concept_id,
        "concept": db.get(Concept, row.concept_id).name,
        "explanation": row.explanation,
        "remediation_instruction": row.remediation_instruction,
        "suggested_activity_id": row.suggested_activity_id,
        "validation_status": row.validation_status,
        "active": row.active,
        "created_at": row.created_at,
        "updated_at": row.updated_at,
    }


@app.get("/api/teacher/misconceptions")
def teacher_misconceptions(
    concept_id: int | None = None,
    db: Session = Depends(get_db),
    _teacher: User = Depends(require_role("teacher")),
):
    query = select(Misconception)
    if concept_id:
        query = query.where(Misconception.concept_id == concept_id)
    return [
        misconception_payload(db, row)
        for row in db.scalars(query.order_by(Misconception.code))
    ]


@app.post("/api/teacher/misconceptions", status_code=201)
def create_misconception(
    payload: MisconceptionInput,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    if not db.get(Concept, payload.concept_id):
        raise HTTPException(status_code=404, detail="Concept not found")
    if payload.suggested_activity_id and not db.get(Activity, payload.suggested_activity_id):
        raise HTTPException(status_code=404, detail="Suggested activity not found")
    if db.scalar(select(Misconception.id).where(Misconception.code == payload.code)):
        raise HTTPException(status_code=409, detail="Misconception code already exists")
    row = Misconception(**payload.model_dump())
    db.add(row)
    db.flush()
    audit(db, teacher.id, "misconception.created", "misconception", row.id)
    db.commit()
    return misconception_payload(db, row)


@app.put("/api/teacher/misconceptions/{misconception_id}")
def update_misconception(
    misconception_id: int,
    payload: MisconceptionInput,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    row = db.get(Misconception, misconception_id)
    if not row:
        raise HTTPException(status_code=404, detail="Misconception not found")
    duplicate = db.scalar(
        select(Misconception.id).where(
            Misconception.code == payload.code,
            Misconception.id != row.id,
        )
    )
    if duplicate:
        raise HTTPException(status_code=409, detail="Misconception code already exists")
    for key, value in payload.model_dump().items():
        setattr(row, key, value)
    audit(db, teacher.id, "misconception.updated", "misconception", row.id)
    db.commit()
    return misconception_payload(db, row)


@app.get("/api/teacher/question-bank")
def question_bank(
    search: str = "",
    subject: str | None = None,
    topic: str | None = None,
    concept_id: int | None = None,
    question_type: str | None = None,
    difficulty: str | None = None,
    source_document_id: int | None = None,
    status_filter: str | None = None,
    ids_only: bool = False,
    page: int = Query(default=1, ge=1),
    page_size: int = Query(default=12, ge=5, le=100),
    db: Session = Depends(get_db),
    _teacher: User = Depends(require_role("teacher")),
):
    query = select(Question)
    if search:
        query = query.where(Question.prompt.ilike(f"%{search}%"))
    if subject:
        query = query.where(Question.subject.ilike(f"%{subject}%"))
    if topic:
        query = query.where(Question.topic.ilike(f"%{topic}%"))
    if concept_id:
        query = query.where(Question.concept_id == concept_id)
    if question_type:
        query = query.where(Question.question_type == question_type)
    if difficulty:
        query = query.where(Question.difficulty_label == difficulty)
    if source_document_id:
        query = query.where(Question.source_document_id == source_document_id)
    if status_filter:
        query = query.where(Question.status == status_filter)
    else:
        query = query.where(Question.status != "Archived")
    total = db.scalar(select(func.count()).select_from(query.subquery())) or 0
    if ids_only:
        return {
            "ids": list(db.scalars(query.with_only_columns(Question.id).order_by(Question.id))),
            "total": total,
        }
    questions = list(
        db.scalars(
            query.order_by(Question.updated_at.desc())
            .offset((page - 1) * page_size)
            .limit(page_size)
        )
    )
    return {
        "items": [question_bank_payload(db, question) for question in questions],
        "total": total,
        "page": page,
        "page_size": page_size,
        "total_pages": max(1, (total + page_size - 1) // page_size),
    }


@app.post("/api/teacher/question-bank", status_code=201)
def create_question_bank_item(
    payload: QuestionBankInput,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    if not db.get(Concept, payload.concept_id):
        raise HTTPException(status_code=404, detail="Concept not found")
    question = Question(
        activity_id=None,
        concept_id=payload.concept_id,
        prompt="",
        feedback="",
        source_type="Manually created",
        created_by=teacher.id,
        points=1,
        active=True,
        position=0,
    )
    apply_question_bank_input(question, payload)
    db.add(question)
    db.flush()
    replace_question_choices(db, question, payload.choices)
    audit(db, teacher.id, "question_bank.created", "question", question.id)
    db.commit()
    return question_bank_payload(db, question)


@app.put("/api/teacher/question-bank/{question_id}")
def update_question_bank_item(
    question_id: int,
    payload: QuestionBankInput,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    question = db.get(Question, question_id)
    if not question:
        raise HTTPException(status_code=404, detail="Question not found")
    if not db.get(Concept, payload.concept_id):
        raise HTTPException(status_code=404, detail="Concept not found")
    apply_question_bank_input(question, payload)
    replace_question_choices(db, question, payload.choices)
    audit(db, teacher.id, "question_bank.updated", "question", question.id)
    db.commit()
    return question_bank_payload(db, question)


def question_assessment_dependencies(db: Session, question_id: int) -> list[str]:
    return list(
        db.scalars(
            select(Assessment.title)
            .join(
                AssessmentQuestion,
                AssessmentQuestion.assessment_id == Assessment.id,
            )
            .where(AssessmentQuestion.question_id == question_id)
            .order_by(Assessment.title)
        )
    )


def delete_question_record(
    db: Session,
    question: Question,
    detach_from_assessments: bool,
) -> list[str]:
    dependencies = question_assessment_dependencies(db, question.id)
    if dependencies and not detach_from_assessments:
        raise HTTPException(
            status_code=409,
            detail=(
                f'Question "{question.prompt[:120]}" belongs to: '
                f"{', '.join(dependencies)}. Confirm removal from those assessments before deleting it."
            ),
        )
    evidence_count = sum(
        db.scalar(select(func.count(model.id)).where(model.question_id == question.id))
        or 0
        for model in (ItemResponse, TutoringResponse, MisconceptionHistory)
    )
    if evidence_count:
        raise HTTPException(
            status_code=409,
            detail=(
                "This question has learner response evidence and cannot be permanently "
                "deleted safely. Archive it instead."
            ),
        )
    if dependencies:
        db.execute(
            delete(AssessmentQuestion).where(
                AssessmentQuestion.question_id == question.id
            )
        )
    db.execute(delete(AnswerChoice).where(AnswerChoice.question_id == question.id))
    db.delete(question)
    return dependencies


@app.delete("/api/teacher/question-bank/bulk-delete")
def bulk_delete_question_bank_items(
    question_ids: list[int] = Query(min_length=1, max_length=100),
    detach_from_assessments: bool = False,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    unique_ids = list(dict.fromkeys(question_ids))
    questions = list(db.scalars(select(Question).where(Question.id.in_(unique_ids))))
    if len(questions) != len(unique_ids):
        raise HTTPException(status_code=404, detail="One or more questions were not found")
    # Preflight every record so a dependency error cannot leave a partial deletion.
    for question in questions:
        dependencies = question_assessment_dependencies(db, question.id)
        if dependencies and not detach_from_assessments:
            raise HTTPException(
                status_code=409,
                detail=(
                    "One or more selected questions belong to assessments: "
                    f"{', '.join(sorted(set(dependencies)))}. Confirm removal from the assessments first."
                ),
            )
        evidence_count = sum(
            db.scalar(
                select(func.count(model.id)).where(model.question_id == question.id)
            )
            or 0
            for model in (ItemResponse, TutoringResponse, MisconceptionHistory)
        )
        if evidence_count:
            raise HTTPException(
                status_code=409,
                detail=(
                    f'Question "{question.prompt[:120]}" has learner response evidence '
                    "and must be archived instead."
                ),
            )
    prompts = [question.prompt for question in questions]
    for question in questions:
        delete_question_record(db, question, detach_from_assessments)
    audit(
        db,
        teacher.id,
        "question_bank.bulk_deleted",
        "question",
        details={"question_ids": unique_ids, "prompts": prompts},
    )
    db.commit()
    return {"ok": True, "deleted": len(unique_ids)}


@app.delete("/api/teacher/question-bank/{question_id}")
def delete_question_bank_item(
    question_id: int,
    detach_from_assessments: bool = False,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    question = db.get(Question, question_id)
    if not question:
        raise HTTPException(status_code=404, detail="Question not found")
    prompt = question.prompt
    dependencies = delete_question_record(
        db, question, detach_from_assessments
    )
    audit(
        db,
        teacher.id,
        "question_bank.deleted",
        "question",
        question_id,
        {"prompt": prompt, "removed_from_assessments": dependencies},
    )
    db.commit()
    return {"ok": True, "id": question_id, "prompt": prompt}


def regenerate_bank_question(
    db: Session, question: Question
) -> None:
    if not question.source_document_id:
        raise HTTPException(
            status_code=409,
            detail="Only questions generated from uploaded material can be regenerated",
        )
    document = db.get(UploadedDocument, question.source_document_id)
    concept = db.get(Concept, question.concept_id)
    if not document or not concept:
        raise HTTPException(status_code=404, detail="Question source is unavailable")
    passages = source_passages(document.extracted_text)
    if not passages:
        raise HTTPException(status_code=422, detail="Question source has no usable text")
    seed_index = (question.id + int(time.time())) % len(passages)
    data = generated_material_question(
        passages,
        seed_index,
        question.question_type,
        concept,
        question.topic,
        question.is_calculation,
    )
    validation_status, flags = validate_generated_question(
        data, concept, question.learning_competency, set()
    )
    question.prompt = clean_content(data["prompt"])
    question.correct_answer = clean_content(data["correct_answer"], 3000)
    question.explanation = f"The answer follows {data['source_locator']}: {data['correct_answer']}"
    question.feedback = question.explanation
    question.source_locator = data["source_locator"]
    question.solution_steps = data["solution_steps"]
    question.validation_status = validation_status
    question.validation_flags = flags
    question.distractor_rationales = data["distractor_rationales"]
    question.is_calculation = data["is_calculation"]
    question.status = "Draft"
    replace_question_choices(db, question, data["choices"])


@app.post("/api/teacher/question-bank/{question_id}/{action}")
def question_bank_action(
    question_id: int,
    action: str,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    if action not in {"duplicate", "archive", "restore", "regenerate"}:
        raise HTTPException(status_code=400, detail="Unsupported question action")
    question = db.get(Question, question_id)
    if not question:
        raise HTTPException(status_code=404, detail="Question not found")
    if action == "archive":
        question.status = "Archived"
        result = question
    elif action == "restore":
        question.status = "Draft"
        result = question
    elif action == "regenerate":
        regenerate_bank_question(db, question)
        result = question
    else:
        result = Question(
            activity_id=None,
            concept_id=question.concept_id,
            prompt=question.prompt,
            feedback=question.feedback,
            hint=question.hint,
            question_type=question.question_type,
            correct_answer=question.correct_answer,
            explanation=question.explanation,
            difficulty_label=question.difficulty_label,
            cognitive_level=question.cognitive_level,
            subject=question.subject,
            topic=question.topic,
            learning_competency=question.learning_competency,
            source_type="Imported from another assessment",
            source_document_id=question.source_document_id,
            source_locator=question.source_locator,
            solution_steps=question.solution_steps,
            validation_status=question.validation_status,
            validation_flags=question.validation_flags,
            distractor_rationales=question.distractor_rationales,
            is_calculation=question.is_calculation,
            status="Draft",
            created_by=teacher.id,
            points=question.points,
            active=True,
            position=0,
        )
        db.add(result)
        db.flush()
        source_choices = list(
            db.scalars(
                select(AnswerChoice)
                .where(AnswerChoice.question_id == question.id)
                .order_by(AnswerChoice.position)
            )
        )
        replace_question_choices(db, result, source_choices)
    audit(
        db,
        teacher.id,
        f"question_bank.{action}",
        "question",
        result.id,
        {"source_question_id": question.id},
    )
    db.commit()
    return question_bank_payload(db, result)


@app.patch("/api/teacher/question-bank/bulk-edit")
def bulk_edit_question_bank_items(
    payload: QuestionBulkEditInput,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    unique_ids = list(dict.fromkeys(payload.question_ids))
    questions = list(db.scalars(select(Question).where(Question.id.in_(unique_ids))))
    if len(questions) != len(unique_ids):
        raise HTTPException(status_code=404, detail="One or more questions were not found")
    if payload.concept_id and not db.get(Concept, payload.concept_id):
        raise HTTPException(status_code=404, detail="Concept not found")
    misconception = db.get(Misconception, payload.misconception_id) if payload.misconception_id else None
    if payload.misconception_id and not misconception:
        raise HTTPException(status_code=404, detail="Misconception tag not found")
    changed_fields = []
    mapping = {
        "concept_id": "concept_id",
        "difficulty": "difficulty_label",
        "status": "status",
        "learning_competency": "learning_competency",
        "question_type": "question_type",
        "cognitive_level": "cognitive_level",
    }
    values = payload.model_dump(exclude={"question_ids", "misconception_id"})
    for input_name, model_name in mapping.items():
        value = values.get(input_name)
        if value is None:
            continue
        changed_fields.append(input_name)
        for question in questions:
            setattr(question, model_name, value)
    if misconception:
        changed_fields.append("misconception_id")
        for question in questions:
            if misconception.concept_id != question.concept_id:
                raise HTTPException(
                    status_code=422,
                    detail="The misconception tag must belong to every selected question's concept",
                )
            db.execute(
                update(AnswerChoice)
                .where(
                    AnswerChoice.question_id == question.id,
                    AnswerChoice.is_correct.is_(False),
                )
                .values(
                    misconception_id=misconception.id,
                    mapping_status="Teacher reviewed",
                )
            )
    audit(
        db,
        teacher.id,
        "question_bank.bulk_edited",
        "question",
        details={"question_ids": unique_ids, "changed_fields": changed_fields},
    )
    db.commit()
    return {
        "ok": True,
        "updated": len(unique_ids),
        "items": [question_bank_payload(db, question) for question in questions],
    }


@app.post("/api/teacher/question-bank/batch")
def question_bank_batch(
    payload: QuestionBatchInput,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    questions = list(
        db.scalars(select(Question).where(Question.id.in_(payload.question_ids)))
    )
    if len(questions) != len(set(payload.question_ids)):
        raise HTTPException(status_code=404, detail="One or more questions were not found")
    for question in questions:
        if payload.action == "archive":
            question.status = "Archived"
        elif payload.action == "save":
            question.status = "Ready"
        else:
            regenerate_bank_question(db, question)
    audit(
        db,
        teacher.id,
        f"question_bank.batch.{payload.action}",
        "question",
        details={"question_ids": payload.question_ids},
    )
    db.commit()
    return {
        "ok": True,
        "items": [question_bank_payload(db, question) for question in questions],
    }


def create_assessment_activity(
    db: Session,
    assessment: Assessment,
    questions: list[Question],
    teacher: User,
) -> Activity:
    activity = Activity(
        title=assessment.title,
        description=assessment.description,
        activity_type="Assessment",
        difficulty=2,
        estimated_minutes=assessment.time_limit or 30,
        instructions=(
            f"Complete this assessment before the due date. "
            f"Mastery target: {round(assessment.mastery_threshold * 100)}%."
        ),
        active=assessment.status == "Published",
        is_diagnostic=False,
        is_demo=teacher.is_demo,
    )
    db.add(activity)
    db.flush()
    concept_ids = sorted({question.concept_id for question in questions})
    for concept_id in concept_ids:
        db.add(ActivityConcept(activity_id=activity.id, concept_id=concept_id))
    for position, source in enumerate(questions, start=1):
        published = Question(
            activity_id=activity.id,
            concept_id=source.concept_id,
            prompt=source.prompt,
            feedback=source.feedback,
            hint=source.hint,
            question_type=source.question_type,
            correct_answer=source.correct_answer,
            explanation=source.explanation,
            difficulty_label=source.difficulty_label,
            cognitive_level=source.cognitive_level,
            subject=source.subject,
            topic=source.topic,
            learning_competency=source.learning_competency,
            source_type="Imported from another assessment",
            source_document_id=source.source_document_id,
            source_locator=source.source_locator,
            solution_steps=source.solution_steps,
            validation_status=source.validation_status,
            validation_flags=source.validation_flags,
            distractor_rationales=source.distractor_rationales,
            is_calculation=source.is_calculation,
            status="Published",
            created_by=teacher.id,
            points=source.points,
            active=True,
            position=position,
        )
        db.add(published)
        db.flush()
        source_choices = list(
            db.scalars(
                select(AnswerChoice)
                .where(AnswerChoice.question_id == source.id)
                .order_by(AnswerChoice.position)
            )
        )
        replace_question_choices(db, published, source_choices)
        source.status = "Published"
    assessment.activity_id = activity.id
    if assessment.status == "Published":
        assessment.published_at = datetime.now(timezone.utc)
    return activity


@app.post("/api/teacher/assessments", status_code=201)
def create_assessment(
    payload: AssessmentInput,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    unique_ids = list(dict.fromkeys(payload.question_ids))
    questions = list(
        db.scalars(select(Question).where(Question.id.in_(unique_ids)))
    )
    if len(questions) != len(unique_ids) or any(
        question.status == "Archived" for question in questions
    ):
        raise HTTPException(
            status_code=400,
            detail="One or more selected questions are unavailable",
        )
    invalid_students = [
        student_id
        for student_id in payload.student_ids
        if not (
            (student := db.get(User, student_id))
            and student.role == "student"
        )
    ]
    if invalid_students:
        raise HTTPException(status_code=400, detail="One or more students are invalid")
    assessment = Assessment(
        **payload.model_dump(
            exclude={"question_ids", "student_ids", "sections"}
        ),
        created_by=teacher.id,
    )
    db.add(assessment)
    db.flush()
    for position, question_id in enumerate(unique_ids, start=1):
        db.add(
            AssessmentQuestion(
                assessment_id=assessment.id,
                question_id=question_id,
                position=position,
            )
        )
    for student_id in dict.fromkeys(payload.student_ids):
        db.add(
            AssessmentAssignment(
                assessment_id=assessment.id,
                student_id=student_id,
            )
        )
    for section in dict.fromkeys(
        clean_content(value, 80) for value in payload.sections if value.strip()
    ):
        db.add(
            AssessmentAssignment(
                assessment_id=assessment.id,
                section=section,
            )
        )
    if assessment.status in {"Published", "Scheduled"}:
        create_assessment_activity(db, assessment, questions, teacher)
    audit(
        db,
        teacher.id,
        "assessment.created",
        "assessment",
        assessment.id,
        {
            "status": assessment.status,
            "question_count": len(questions),
            "student_count": len(payload.student_ids),
            "sections": payload.sections,
        },
    )
    db.commit()
    return {
        "id": assessment.id,
        "title": assessment.title,
        "status": assessment.status,
        "question_count": len(questions),
        "activity_id": assessment.activity_id,
    }


@app.get("/api/teacher/assessments")
def teacher_assessments(
    db: Session = Depends(get_db),
    _teacher: User = Depends(require_role("teacher")),
):
    assessments = list(
        db.scalars(select(Assessment).order_by(Assessment.updated_at.desc()))
    )
    return [
        {
            "id": assessment.id,
            "title": assessment.title,
            "description": assessment.description,
            "subject": assessment.subject,
            "topic": assessment.topic,
            "status": assessment.status,
            "mastery_threshold": assessment.mastery_threshold,
            "time_limit": assessment.time_limit,
            "maximum_attempts": assessment.maximum_attempts,
            "available_from": assessment.available_from,
            "due_at": assessment.due_at,
            "published_at": assessment.published_at,
            "activity_id": assessment.activity_id,
            "learner_record_count": (
                db.scalar(
                    select(func.count(AssessmentAttempt.id)).where(
                        AssessmentAttempt.activity_id == assessment.activity_id
                    )
                )
                if assessment.activity_id
                else 0
            )
            or 0,
            "question_count": db.scalar(
                select(func.count(AssessmentQuestion.id)).where(
                    AssessmentQuestion.assessment_id == assessment.id
                )
            )
            or 0,
        }
        for assessment in assessments
    ]


@app.post("/api/teacher/assessments/{assessment_id}/status")
def update_assessment_status(
    assessment_id: int,
    payload: AssessmentStatusInput,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    assessment = db.get(Assessment, assessment_id)
    if not assessment:
        raise HTTPException(status_code=404, detail="Assessment not found")
    previous_status = assessment.status
    if payload.status in {"Published", "Scheduled"} and not assessment.activity_id:
        question_links = list(
            db.scalars(
                select(AssessmentQuestion)
                .where(AssessmentQuestion.assessment_id == assessment.id)
                .order_by(AssessmentQuestion.position)
            )
        )
        questions = [
            db.get(Question, link.question_id) for link in question_links
        ]
        if not questions or any(question is None for question in questions):
            raise HTTPException(
                status_code=409,
                detail="Assessment questions are unavailable",
            )
        assessment.status = payload.status
        create_assessment_activity(db, assessment, questions, teacher)
    else:
        assessment.status = payload.status
        activity = (
            db.get(Activity, assessment.activity_id)
            if assessment.activity_id
            else None
        )
        if activity:
            activity.active = payload.status == "Published"
        if payload.status == "Published":
            assessment.published_at = datetime.now(timezone.utc)
    audit(
        db,
        teacher.id,
        "assessment.status_changed",
        "assessment",
        assessment.id,
        {"from": previous_status, "to": payload.status},
    )
    db.commit()
    return {
        "id": assessment.id,
        "title": assessment.title,
        "status": assessment.status,
        "activity_id": assessment.activity_id,
    }


@app.delete("/api/teacher/assessments/{assessment_id}")
def delete_assessment(
    assessment_id: int,
    confirm_learner_record_deletion: bool = False,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    assessment = db.get(Assessment, assessment_id)
    if not assessment:
        raise HTTPException(status_code=404, detail="Assessment not found")
    attempt_ids = (
        list(
            db.scalars(
                select(AssessmentAttempt.id).where(
                    AssessmentAttempt.activity_id == assessment.activity_id
                )
            )
        )
        if assessment.activity_id
        else []
    )
    if attempt_ids and not confirm_learner_record_deletion:
        raise HTTPException(
            status_code=409,
            detail=(
                f'Deleting "{assessment.title}" will also permanently remove '
                f"{len(attempt_ids)} learner attempt record(s) and their linked responses, "
                "mental-effort ratings, mastery evidence, and interaction logs. Confirm this deletion to continue."
            ),
        )
    title = assessment.title
    activity = db.get(Activity, assessment.activity_id) if assessment.activity_id else None
    if attempt_ids:
        delete_attempt_records(db, attempt_ids)
    db.execute(
        delete(AssessmentAssignment).where(
            AssessmentAssignment.assessment_id == assessment.id
        )
    )
    db.execute(
        delete(AssessmentQuestion).where(
            AssessmentQuestion.assessment_id == assessment.id
        )
    )
    if activity:
        activity.active = False
    db.delete(assessment)
    audit(
        db,
        teacher.id,
        "assessment.deleted",
        "assessment",
        assessment_id,
        {
            "title": title,
            "learner_attempts_removed": len(attempt_ids),
            "activity_id": activity.id if activity else None,
        },
    )
    db.commit()
    return {
        "ok": True,
        "id": assessment_id,
        "title": title,
        "learner_records_removed": len(attempt_ids),
    }


@app.get("/api/teacher/settings")
def teacher_settings(
    db: Session = Depends(get_db), _teacher: User = Depends(require_role("teacher"))
):
    return settings_payload(db)


@app.put("/api/teacher/settings")
def update_settings(
    payload: SettingsInput,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    try:
        payload.validate_combination()
    except ValueError as error:
        raise HTTPException(status_code=422, detail=str(error))
    result = save_settings(db, payload.model_dump())
    active_students = list(
        db.scalars(
            select(User)
            .join(StudentProfile, StudentProfile.user_id == User.id)
            .where(
                User.role == "student",
                User.is_active.is_(True),
                User.account_status == "Active",
                StudentProfile.target_concept_id.is_not(None),
            )
        )
    )
    for student in active_students:
        generate_pathways(
            db,
            student,
            trigger_type="Optimization settings updated",
            trigger_id=teacher.id,
        )
    audit(db, teacher.id, "settings.updated", "system_settings", details=payload.model_dump())
    db.commit()
    return result


@app.post("/api/teacher/models/train")
def train_model(
    mode: str = Query(default="demo", pattern="^(demo|research)$"),
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    try:
        version = train_ensemble(db, is_demo=mode == "demo")
    except ValueError as error:
        raise HTTPException(status_code=400, detail=str(error))
    audit(db, teacher.id, "model.trained", "model_version", version.id, {"mode": mode})
    db.commit()
    return {
        "id": version.id,
        "version": version.version,
        "sample_size": version.sample_size,
        "student_count": version.student_count,
        "metrics": version.metrics,
        "warning": version.warning,
        "is_demo": version.is_demo,
    }


@app.get("/api/teacher/models")
def model_versions(
    db: Session = Depends(get_db), _teacher: User = Depends(require_role("teacher"))
):
    return [
        {
            "id": version.id,
            "version": version.version,
            "trained_at": version.trained_at,
            "sample_size": version.sample_size,
            "student_count": version.student_count,
            "metrics": version.metrics,
            "feature_names": version.feature_names,
            "metadata": version.metadata_json or {
                "algorithm": "Soft-voting ensemble",
                "ensemble_members": ["LogisticRegression", "RandomForestClassifier", "GradientBoostingClassifier"],
                "class_labels": ["Low", "Moderate", "High"],
                "evaluation_method": (version.metrics or {}).get("evaluation"),
                "deployment_status": "Active" if version.active else "Historical",
            },
            "warning": version.warning,
            "is_demo": version.is_demo,
            "active": version.active,
        }
        for version in db.scalars(
            select(ModelVersion).order_by(ModelVersion.trained_at.desc())
        )
    ]


@app.get("/api/teacher/models/predict")
def learner_cognitive_load_prediction(
    student_id: int,
    db: Session = Depends(get_db),
    _teacher: User = Depends(require_role("teacher")),
):
    try:
        return predict_student_cognitive_load(db, student_id)
    except ValueError as error:
        raise HTTPException(status_code=404, detail=str(error))


@app.post("/api/teacher/models/diagnose-all")
def diagnose_all_active_students(
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    students = list(db.scalars(select(User).where(
        User.role == "student",
        User.is_active.is_(True),
        User.account_status == "Active",
        User.is_demo == teacher.is_demo,
    ).order_by(User.participant_code)))
    results = []
    for student in students:
        try:
            results.append(predict_student_cognitive_load(db, student.id))
        except ValueError as error:
            results.append({
                "available": False,
                "student_id": student.id,
                "participant_code": student.participant_code,
                "message": str(error),
            })
    audit(db, teacher.id, "model.diagnosed_all", "cognitive_load_prediction", details={
        "students": len(students),
        "predictions": sum(1 for item in results if item.get("available")),
    })
    db.commit()
    return {"items": results, "students": len(students)}


@app.get("/api/teacher/models/predictions")
def stored_cognitive_load_predictions(
    student_id: int | None = None,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    query = select(CognitiveLoadPrediction).join(User, User.id == CognitiveLoadPrediction.student_id).where(
        CognitiveLoadPrediction.is_demo == teacher.is_demo
    )
    if student_id:
        query = query.where(CognitiveLoadPrediction.student_id == student_id)
    rows = list(db.scalars(query.order_by(CognitiveLoadPrediction.created_at.desc()).limit(200)))
    return [{
        "id": row.id,
        "student_id": row.student_id,
        "participant_code": db.get(User, row.student_id).participant_code,
        "model_version": db.get(ModelVersion, row.model_version_id).version,
        "evidence_date": row.evidence_date,
        "probabilities": row.probabilities,
        "category": row.predicted_category,
        "expected_index": row.expected_index,
        "confidence": row.confidence,
        "evidence": row.evidence,
        "missing_features": row.missing_features,
        "feature_contributions": row.feature_contributions,
        "recommended_action": row.recommended_action,
        "created_at": row.created_at,
    } for row in rows]


@app.get("/api/teacher/pathways")
def compare_pathways(
    student_id: int | None = None,
    db: Session = Depends(get_db),
    _teacher: User = Depends(require_role("teacher")),
):
    query = (
        select(PathwayRecommendation)
        .join(User, User.id == PathwayRecommendation.student_id)
        .where(
            PathwayRecommendation.active.is_(True),
            User.is_active.is_(True),
            User.account_status == "Active",
        )
    )
    if student_id:
        query = query.where(PathwayRecommendation.student_id == student_id)
    pathways = list(
        db.scalars(query.order_by(PathwayRecommendation.created_at.desc()))
    )
    groups: dict[tuple[int, int], list[PathwayRecommendation]] = defaultdict(list)
    for pathway in pathways:
        groups[(pathway.student_id, pathway.target_concept_id)].append(pathway)
    output = []
    changed = False
    for group in groups.values():
        alpha = float(get_setting(db, "alpha"))
        beta = float(get_setting(db, "beta"))
        gamma = float(get_setting(db, "gamma"))
        for pathway in group:
            recalculated = adaptive_pathway_score(
                pathway.gap_coverage,
                pathway.predicted_cognitive_load,
                pathway.normalized_learning_time,
                alpha,
                beta,
                gamma,
            )
            if abs(pathway.adaptive_pathway_score - recalculated) > 1e-9:
                pathway.adaptive_pathway_score = recalculated
                changed = True
        ranked = sorted(
            group,
            key=lambda item: (
                -item.adaptive_pathway_score,
                -item.gap_coverage,
                item.predicted_cognitive_load,
                item.total_minutes,
                item.id,
            ),
        )
        for rank, pathway in enumerate(ranked, start=1):
            should_select = rank == 1
            if pathway.selected != should_select:
                pathway.selected = should_select
                changed = True
            serialized = serialize_pathway(db, pathway)
            decision = serialized.get("decision_explanation") or {}
            weights = {
                "alpha": float(get_setting(db, "alpha")),
                "beta": float(get_setting(db, "beta")),
                "gamma": float(get_setting(db, "gamma")),
            }
            serialized.update(
                {
                    "participant_code": db.get(User, pathway.student_id).participant_code,
                    "student_id": pathway.student_id,
                    "rank": rank,
                    "weights": weights,
                    "weighted_contributions": {
                        "gap_coverage": weights["alpha"] * pathway.gap_coverage,
                        "load_suitability": weights["beta"]
                        * (1 - pathway.predicted_cognitive_load),
                        "time_efficiency": weights["gamma"]
                        * (1 - pathway.normalized_learning_time),
                    },
                    "prerequisite_sequence": [
                        step["concept"] for step in serialized.get("steps", [])
                    ],
                    "learning_gaps_addressed": (
                        decision.get("prerequisite_chain") or []
                    ),
                    "selection_reason": (
                        f"{pathway.label} was selected with the highest APS of {pathway.adaptive_pathway_score:.3f}. "
                        f"It covers {pathway.gap_coverage:.0%} of detected gaps, has predicted cognitive load {pathway.predicted_cognitive_load:.3f}, "
                        f"and requires about {pathway.total_minutes} minutes."
                        if rank == 1
                        else f"{pathway.label} was not selected because its APS of {pathway.adaptive_pathway_score:.3f} is below "
                        f"{ranked[0].label} at {ranked[0].adaptive_pathway_score:.3f}. It covers {pathway.gap_coverage:.0%} of gaps, "
                        f"has predicted cognitive load {pathway.predicted_cognitive_load:.3f}, and requires {pathway.total_minutes} minutes."
                    ),
                    "tie_breaker": (
                        "Higher gap coverage, then lower cognitive load, then shorter estimated time, then stable database identifier."
                    ),
                }
            )
            output.append(serialized)
    if changed:
        db.commit()
    return sorted(output, key=lambda item: (item["participant_code"], item["rank"]))


@app.delete("/api/teacher/pathways/{pathway_id}")
def delete_saved_pathway_comparison(
    pathway_id: int,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    pathway = db.get(PathwayRecommendation, pathway_id)
    if not pathway:
        raise HTTPException(status_code=404, detail="Pathway comparison not found")
    student_id = pathway.student_id
    target_concept_id = pathway.target_concept_id
    label = pathway.label
    db.execute(
        update(PathwayRecommendation)
        .where(PathwayRecommendation.supersedes_pathway_id == pathway.id)
        .values(supersedes_pathway_id=None)
    )
    db.execute(delete(ExpertEvaluation).where(ExpertEvaluation.pathway_id == pathway.id))
    db.execute(
        delete(TeacherIntervention).where(
            TeacherIntervention.pathway_id == pathway.id
        )
    )
    db.execute(delete(PathwayStep).where(PathwayStep.pathway_id == pathway.id))
    db.execute(
        delete(PathwayVersion).where(
            or_(
                PathwayVersion.pathway_id == pathway.id,
                PathwayVersion.previous_pathway_id == pathway.id,
            )
        )
    )
    db.delete(pathway)
    db.flush()
    remaining = list(
        db.scalars(
            select(PathwayRecommendation).where(
                PathwayRecommendation.student_id == student_id,
                PathwayRecommendation.target_concept_id == target_concept_id,
                PathwayRecommendation.active.is_(True),
            )
        )
    )
    ranked = sorted(
        remaining,
        key=lambda item: (
            -item.adaptive_pathway_score,
            -item.gap_coverage,
            item.predicted_cognitive_load,
            item.total_minutes,
            item.id,
        ),
    )
    for index, item in enumerate(ranked):
        item.selected = index == 0
    audit(
        db,
        teacher.id,
        "pathway.deleted",
        "pathway_recommendation",
        pathway_id,
        {"label": label, "student_id": student_id},
    )
    db.commit()
    return {"ok": True, "id": pathway_id, "label": label}


@app.post("/api/teacher/evaluations", status_code=201)
def create_evaluation(
    payload: EvaluationInput,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    pathway = db.get(PathwayRecommendation, payload.pathway_id)
    if not pathway:
        raise HTTPException(status_code=404, detail="Pathway not found")
    if teacher.is_demo != pathway.is_demo:
        raise HTTPException(
            status_code=403,
            detail="Teacher and pathway must belong to the same demo or research mode",
        )
    maximum = int(get_setting(db, "likert_scale_max"))
    scores = [
        payload.recommendation_accuracy,
        payload.adaptability,
        payload.personalization,
        payload.optimization_efficiency,
        payload.pathway_relevance,
    ]
    if any(value < 1 or value > maximum for value in scores):
        raise HTTPException(status_code=422, detail=f"Ratings must be between 1 and {maximum}")
    evaluation = ExpertEvaluation(
        teacher_id=teacher.id,
        is_demo=pathway.is_demo,
        **payload.model_dump(),
    )
    db.add(evaluation)
    db.flush()
    audit(db, teacher.id, "evaluation.created", "expert_evaluation", evaluation.id)
    db.commit()
    return {"id": evaluation.id, "ok": True}


EXPORTS = {
    "interactions": (
        InteractionLog,
        [
            "student_id",
            "activity_id",
            "concept_id",
            "score",
            "max_score",
            "response_accuracy",
            "average_response_seconds",
            "total_completion_seconds",
            "number_of_attempts",
            "skipped_items",
            "hint_usage_count",
            "submission_time",
        ],
    ),
    "mastery": (
        MasteryRecord,
        ["student_id", "concept_id", "mastery_score", "classification", "calculation_mode", "created_at"],
    ),
    "gaps": (
        LearningGap,
        ["student_id", "concept_id", "mastery_score", "threshold", "reason", "created_at", "resolved_at"],
    ),
    "mental-effort": (
        MentalEffortRating,
        ["student_id", "attempt_id", "rating", "category", "created_at"],
    ),
    "pathways": (
        PathwayRecommendation,
        [
            "student_id",
            "target_concept_id",
            "label",
            "selected",
            "gap_coverage",
            "predicted_cognitive_load",
            "normalized_learning_time",
            "adaptive_pathway_score",
            "total_minutes",
            "created_at",
        ],
    ),
    "evaluations": (
        ExpertEvaluation,
        [
            "pathway_id",
            "recommendation_accuracy",
            "adaptability",
            "personalization",
            "optimization_efficiency",
            "pathway_relevance",
            "comment",
            "created_at",
        ],
    ),
    "models": (
        ModelVersion,
        ["version", "trained_at", "sample_size", "student_count", "metrics", "warning", "is_demo"],
    ),
}


@app.get("/api/teacher/exports/{export_type}")
def export_csv(
    export_type: str,
    mode: str = Query(default="demo", pattern="^(demo|research)$"),
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    if export_type not in EXPORTS:
        raise HTTPException(status_code=404, detail="Export type not found")
    model, fields = EXPORTS[export_type]
    query = select(model)
    if hasattr(model, "is_demo"):
        query = query.where(model.is_demo == (mode == "demo"))
    rows = list(db.scalars(query))
    student_codes = {
        user.id: user.participant_code
        for user in db.scalars(select(User).where(User.role == "student"))
    }
    buffer = io.StringIO()
    writer = csv.writer(buffer)
    headers = [
        "anonymous_participant_code" if field == "student_id" else field.replace("_", " ").title()
        for field in fields
    ]
    writer.writerow(headers)
    for row in rows:
        values = []
        for field in fields:
            value = getattr(row, field)
            if field == "student_id":
                value = student_codes.get(value, "UNKNOWN")
            values.append(value)
        writer.writerow(values)
    audit(db, teacher.id, "research.exported", export_type, details={"mode": mode, "rows": len(rows)})
    db.commit()
    filename = f"neurolearnx_{mode}_{export_type}_{datetime.now().date().isoformat()}.csv"
    return StreamingResponse(
        iter([buffer.getvalue()]),
        media_type="text/csv",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@app.get("/api/teacher/audit-logs")
def audit_logs(
    limit: int = Query(default=100, ge=1, le=500),
    db: Session = Depends(get_db),
    _teacher: User = Depends(require_role("teacher")),
):
    return [
        {
            "id": row.id,
            "actor": db.get(User, row.actor_id).participant_code if row.actor_id else "system",
            "action": row.action,
            "entity_type": row.entity_type,
            "entity_id": row.entity_id,
            "details": row.details,
            "timestamp": row.timestamp,
        }
        for row in db.scalars(
            select(AuditLog).order_by(AuditLog.timestamp.desc()).limit(limit)
        )
    ]


@app.post("/api/teacher/reset-demo")
def reset_demo(
    payload: ResetInput,
    db: Session = Depends(get_db),
    teacher: User = Depends(require_role("teacher")),
):
    if payload.confirmation != "RESET DEMO DATA":
        raise HTTPException(
            status_code=400, detail='Type "RESET DEMO DATA" to confirm this action'
        )
    from .seed import reset_and_seed_demo

    teacher_code = teacher.participant_code
    reset_and_seed_demo(db)
    reset_actor = db.scalar(
        select(User).where(
            User.participant_code == teacher_code,
            User.role == "teacher",
            User.is_demo.is_(True),
        )
    )
    audit(db, reset_actor.id if reset_actor else None, "demo.reset", "system")
    db.commit()
    return {"ok": True}


frontend_dist = Path(__file__).resolve().parents[2] / "frontend" / "dist"
if frontend_dist.exists() and os.getenv("VERCEL") != "1":
    app.mount("/", StaticFiles(directory=frontend_dist, html=True), name="frontend")
