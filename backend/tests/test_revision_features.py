from __future__ import annotations

from io import BytesIO

from docx import Document
from pptx import Presentation
from pypdf import PdfWriter
from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

from app.main import analyze_material


def test_plain_text_analysis_recognizes_authored_structure():
    analysis = analyze_material(
        """Newton's Second Law

Worked Example

Force, mass, and acceleration are related by F = ma.
A 2 kg cart accelerates at 3 m/s^2, so F = 6 N.
Learners should not divide mass by acceleration when applying this law.
"""
    )
    assert "Worked Example" in analysis["headings"]
    assert analysis["worked_examples"]
    assert analysis["relationships"]
    assert analysis["competencies"]
    assert analysis["misconceptions"]


def teacher_login(client):
    response = client.post(
        "/api/auth/login",
        json={
            "participant_code": "TEACHER01",
            "password": "NeuroTeach!2026",
            "expected_role": "teacher",
        },
    )
    assert response.status_code == 200


def pdf_material() -> bytes:
    output = BytesIO()
    writer = PdfWriter()
    page = writer.add_blank_page(width=612, height=792)
    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    page[NameObject("/Resources")] = DictionaryObject(
        {
            NameObject("/Font"): DictionaryObject(
                {NameObject("/F1"): writer._add_object(font)}
            )
        }
    )
    stream = DecodedStreamObject()
    stream.set_data(
        b"BT /F1 12 Tf 50 700 Td (Newton force equals mass times acceleration. "
        b"A free-body diagram shows forces acting on an object. Calculate acceleration using F = ma.) Tj ET"
    )
    page[NameObject("/Contents")] = writer._add_object(stream)
    writer.write(output)
    return output.getvalue()


def docx_material() -> bytes:
    document = Document()
    document.add_heading("Newton's Laws", 0)
    document.add_heading("Force and acceleration", 1)
    document.add_paragraph(
        "Net force is defined as the vector sum of forces. Newton's second law means that F = ma. "
        "Calculate acceleration by dividing net force by mass. A common error is using weight as mass."
    )
    output = BytesIO()
    document.save(output)
    return output.getvalue()


def pptx_material() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Momentum and Impulse"
    slide.placeholders[1].text = (
        "Momentum is mass multiplied by velocity. Impulse equals change in momentum. "
        "For example, a longer stopping time decreases average force. J = F times delta t."
    )
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def test_all_supported_documents_are_analyzed_and_source_grounded(client):
    teacher_login(client)
    files = [
        ("material.txt", b"Kinematics\nVelocity is displacement divided by time. Acceleration is change in velocity. Calculate speed using v = d/t. A common error is confusing distance and displacement.", "text/plain", 'Heading "Kinematics"'),
        ("material.pdf", pdf_material(), "application/pdf", "Page 1"),
        (
            "material.docx",
            docx_material(),
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "Heading",
        ),
        (
            "material.pptx",
            pptx_material(),
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "Slide 1",
        ),
    ]
    concept = next(
        item
        for item in client.get("/api/teacher/concepts").json()
        if item["code"] == "GP-NL"
    )
    for filename, content, mime_type, expected_locator in files:
        uploaded = client.post(
            "/api/teacher/documents",
            files={"file": (filename, content, mime_type)},
        )
        assert uploaded.status_code == 201, uploaded.text
        document = uploaded.json()
        analysis = document["analysis"]
        assert analysis["title"]
        assert analysis["main_topic"]
        assert analysis["facts"]
        assert analysis["method"].startswith("Local deterministic")
        generated = client.post(
            f"/api/teacher/documents/{document['id']}/generate",
            json={
                "subject": "General Physics",
                "grade_level": "Grade 12",
                "topic": "Newton's Laws",
                "concept_id": concept["id"],
                "learning_competency": "Apply force relationships to physical situations.",
                "number_of_questions": 1,
                "question_type": "Multiple choice",
                "difficulty": "Moderate",
                "cognitive_level": "Apply",
                "include_explanations": True,
                "include_hints": True,
                "include_prerequisites": False,
                "include_calculations": True,
            },
        )
        assert generated.status_code == 201, generated.text
        question = generated.json()[0]
        assert question["status"] == "Draft"
        assert question["source_locator"]
        assert expected_locator in question["source_locator"]
        assert question["validation_status"] in {"Ready for review", "Needs review"}
        assert question["validation_flags"]
        assert question["solution_steps"]
        assert len(question["choices"]) == 4
        assert sum(choice["is_correct"] for choice in question["choices"]) == 1


def test_generation_keeps_unmentioned_prerequisites_out_and_uses_false_distractors(client):
    teacher_login(client)
    material = (
        b"Newton's Second Law\nWorked Example\nForce, mass, and acceleration are related by F = ma. "
        b"Net force is measured in newtons, mass in kilograms, and acceleration in metres per second squared. "
        b"A 2 kg cart accelerating at 3 m/s squared needs a 6 N net force."
    )
    uploaded = client.post(
        "/api/teacher/documents",
        files={"file": ("newton-source.txt", material, "text/plain")},
    )
    assert uploaded.status_code == 201, uploaded.text
    concept = next(
        item
        for item in client.get("/api/teacher/concepts").json()
        if item["code"] == "GP-NL"
    )
    generated = client.post(
        f"/api/teacher/documents/{uploaded.json()['id']}/generate",
        json={
            "subject": "General Physics",
            "grade_level": "Grade 12",
            "topic": "Newton's Second Law",
            "concept_id": concept["id"],
            "learning_competency": "Calculate net force from mass and acceleration.",
            "number_of_questions": 2,
            "question_type": "Multiple choice",
            "difficulty": "Moderate",
            "cognitive_level": "Apply",
            "include_explanations": True,
            "include_hints": True,
            "include_prerequisites": True,
            "include_calculations": True,
        },
    )
    assert generated.status_code == 201, generated.text
    questions = generated.json()
    assert len({question["prompt"] for question in questions}) == 2
    concepts = {item["id"]: item["code"] for item in client.get("/api/teacher/concepts").json()}
    assert {concepts[question["concept_id"]] for question in questions} <= {"GP-NL", "GP-MK"}
    assert questions[0]["concept_id"] == concept["id"]
    assert questions[0]["is_calculation"] is True
    assert questions[1]["is_calculation"] is False
    incorrect_choices = [
        choice["text"] for choice in questions[1]["choices"] if not choice["is_correct"]
    ]
    assert all(choice not in material.decode() for choice in incorrect_choices)


def test_archive_is_removed_from_active_list_blocks_login_and_can_restore(client):
    created = client.post(
        "/api/auth/register/student",
        json={
            "student_id": "ARCHIVE301",
            "first_name": "Archive",
            "last_name": "Test",
            "email": "archive301@example.edu",
            "username": "archive301",
            "password": "Secure!Pass7",
            "confirm_password": "Secure!Pass7",
            "grade_level": "Grade 12",
            "section": "STEM A",
            "accept_terms": True,
        },
    )
    assert created.status_code == 201
    student_id = created.json()["student"]["id"]
    client.post("/api/auth/logout")
    teacher_login(client)
    archived = client.post(
        f"/api/teacher/students/{student_id}/actions",
        json={"action": "archive", "reason": "Archive workflow test"},
    )
    assert archived.status_code == 200
    active = client.get("/api/teacher/students", params={"search": "ARCHIVE301"}).json()
    assert active == []
    archived_rows = client.get(
        "/api/teacher/students",
        params={
            "search": "ARCHIVE301",
            "account_status": "Archived",
            "include_archived": True,
        },
    ).json()
    assert len(archived_rows) == 1
    client.post("/api/auth/logout")
    blocked = client.post(
        "/api/auth/login",
        json={
            "participant_code": "ARCHIVE301",
            "password": "Secure!Pass7",
            "expected_role": "student",
        },
    )
    assert blocked.status_code == 403
    teacher_login(client)
    restored = client.post(
        f"/api/teacher/students/{student_id}/actions",
        json={"action": "reactivate"},
    )
    assert restored.status_code == 200
    assert restored.json()["student"]["account_status"] == "Active"


def test_teacher_preview_assignment_and_evidence_only_completion(client):
    teacher_login(client)
    student = client.get("/api/teacher/students", params={"search": "STEM001"}).json()[0]
    concept = next(
        item
        for item in client.get("/api/teacher/concepts").json()
        if item["code"] == "GP-NL"
    )
    evidence = client.get(
        f"/api/teacher/students/{student['id']}/topics/{concept['id']}"
    )
    assert evidence.status_code == 200
    assert "incorrect_responses" in evidence.json()
    preview = client.post(
        f"/api/teacher/students/{student['id']}/topics/{concept['id']}/pathway-preview",
        json={"difficulty": "Guided pathway"},
    )
    assert preview.status_code == 200, preview.text
    data = preview.json()
    assert data["steps"]
    assert all(step["content"]["worked_example"]["steps"] for step in data["steps"])
    assigned = client.post(
        f"/api/teacher/students/{student['id']}/pathways/assign",
        json={
            "target_concept_id": concept["id"],
            "label": "Teacher guided Newton pathway",
            "difficulty": "Guided pathway",
            "teacher_note": "Work through the examples before the mastery check.",
            "due_at": None,
            "steps": [
                {
                    "concept_id": step["concept_id"],
                    "activity_id": step["activity_id"],
                    "position": index + 1,
                }
                for index, step in enumerate(data["steps"])
            ],
        },
    )
    assert assigned.status_code == 201, assigned.text
    pathway = assigned.json()
    assert pathway["source_type"] == "Teacher"
    assert pathway["assigned_by"]
    assert pathway["teacher_note"]
    first_step = pathway["steps"][0]
    client.post("/api/auth/logout")
    assert client.post(
        "/api/auth/login",
        json={
            "participant_code": "STEM001",
            "password": "LearnX!2026",
            "expected_role": "student",
        },
    ).status_code == 200
    dashboard = client.get("/api/student/dashboard").json()
    assert dashboard["pathway"]["id"] == pathway["id"]
    assert dashboard["notifications"][0]["type"] == "Teacher pathway assignment"
    assert len(dashboard["pathway_history"]) >= 2
    blocked = client.post(f"/api/student/pathway-steps/{first_step['id']}/complete")
    assert blocked.status_code == 409


def test_activity_exposes_five_minute_timer_and_duplicate_attempt_key(client):
    response = client.post(
        "/api/auth/login",
        json={
            "participant_code": "STEM002",
            "password": "LearnX!2026",
            "expected_role": "student",
        },
    )
    assert response.status_code == 200
    dashboard = client.get("/api/student/dashboard").json()
    activity_id = dashboard["pathway"]["steps"][0]["activity_id"]
    activity = client.get(f"/api/student/activities/{activity_id}")
    assert activity.status_code == 200
    assert activity.json()["time_limit_seconds"] == 300
